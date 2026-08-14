# -*- coding: utf-8 -*-
"""
Cortex MVA 操作手冊 v3 — BG/BU 權限隔離版
對齊 Cortex_互動Demo_v0.11.html · cleansheet-mva-sd v0.4
新增章節:
  - BG/BU 組織架構 + 權限模型
  - 設備類別(取代實體設備)
  - DL 角色費率(取代廠級單一 wage)
  - 設定 master UI(廠 baseline + 類別 + DL/IDL + admin copy + 權限 grant)
~ 50 slides · 16:9
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

INK    = RGBColor(0x0B, 0x1F, 0x3A)
TEXT   = RGBColor(0x37, 0x47, 0x55)
MUTED  = RGBColor(0x6B, 0x72, 0x80)
OCEAN  = RGBColor(0x02, 0xC3, 0x9A)
NAVY   = RGBColor(0x1C, 0x72, 0x93)
LINE   = RGBColor(0xE5, 0xE7, 0xEB)
BG_SOFT= RGBColor(0xFA, 0xFB, 0xFC)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GOLD   = RGBColor(0xCA, 0x8A, 0x04)
RED    = RGBColor(0xDC, 0x26, 0x26)
GREEN  = RGBColor(0x16, 0xA3, 0x4A)
PURPLE = RGBColor(0x7C, 0x3A, 0xED)
CYAN   = RGBColor(0x08, 0x91, 0xB2)
SKY    = RGBColor(0x0E, 0xA5, 0xE9)
ORANGE = RGBColor(0xF5, 0xA5, 0x24)


def add_slide(): return prs.slides.add_slide(prs.slide_layouts[6])


def add_text(s, x, y, w, h, text, size=12, bold=False, color=None,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font='Microsoft JhengHei'):
    tx = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tx.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    if color:
        r.font.color.rgb = color
    return tx


def add_rect(s, x, y, w, h, fill, line=None, line_width=0.5):
    rect = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    rect.fill.solid()
    rect.fill.fore_color.rgb = fill
    if line is None:
        rect.line.fill.background()
    else:
        rect.line.color.rgb = line
        rect.line.width = Pt(line_width)
    rect.shadow.inherit = False
    return rect


def add_rounded(s, x, y, w, h, fill, line=None, line_width=0.5, radius=0.05):
    rect = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    rect.adjustments[0] = radius
    rect.fill.solid()
    rect.fill.fore_color.rgb = fill
    if line is None:
        rect.line.fill.background()
    else:
        rect.line.color.rgb = line
        rect.line.width = Pt(line_width)
    rect.shadow.inherit = False
    return rect


def header(s, title, subtitle=None, color=NAVY, badge=None):
    add_rect(s, 0, 0, 13.333, 0.16, color)
    if badge:
        add_rounded(s, 0.4, 0.3, 1.5, 0.45, color, radius=0.3)
        add_text(s, 0.4, 0.32, 1.5, 0.4, badge, size=12, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, 2.05, 0.25, 11, 0.5, title, size=22, bold=True, color=INK)
        if subtitle:
            add_text(s, 2.05, 0.78, 11, 0.32, subtitle, size=11, color=MUTED)
    else:
        add_text(s, 0.4, 0.25, 11, 0.5, title, size=24, bold=True, color=INK)
        if subtitle:
            add_text(s, 0.4, 0.8, 11, 0.32, subtitle, size=12, color=MUTED)


def footer(s, n, total):
    add_rect(s, 0.4, 7.08, 12.55, 0.02, LINE)
    add_text(s, 0.4, 7.12, 9, 0.22, 'Cortex MVA 操作手冊 v3 · BG/BU 權限隔離版 · 對齊 v0.11 demo',
             size=8.5, color=MUTED)
    add_text(s, 12.0, 7.12, 1, 0.22, f'{n} / {total}', size=8.5, color=MUTED, align=PP_ALIGN.RIGHT)


# ════════════════════════════════════════════════════════════════════════
# SLIDES
# ════════════════════════════════════════════════════════════════════════

def slide_cover():
    s = add_slide()
    add_rect(s, 0, 0, 13.333, 7.5, INK)
    add_rect(s, 0, 2.7, 13.333, 0.06, OCEAN)
    add_text(s, 0.7, 1.2, 12, 0.4, 'FOXLINK 正崴 · Cortex BOM × MVA', size=12, color=OCEAN)
    add_text(s, 0.7, 1.6, 12, 0.95, 'MVA 操作手冊 v3', size=48, bold=True, color=WHITE)
    add_text(s, 0.7, 3.0, 12, 0.5, 'BG/BU 權限隔離 + 設備類別 + DL 角色 + 設定維護 UI', size=20, color=WHITE)

    add_rounded(s, 0.7, 3.85, 12, 2.0, RGBColor(0x14, 0x2F, 0x4E), radius=0.04)
    add_text(s, 1.0, 4.0, 11.5, 0.4, 'v3 5 大規格修正(對齊 v0.11 demo)', size=14, bold=True, color=OCEAN)
    items = [
        ('🔒', 'BG/BU 隔離', '同廠不同 BG 看不到對方數字 · baseline 按 BG 切版', SKY),
        ('🏭', '設備改類別', '報價階段 SMT 整線 × qty · 試產才細到 DEK / FUJI', CYAN),
        ('👷', 'DL 角色費率', '取代廠級單一 $4.95/hr · 7 個工種各自費率(對齊 IDL)', RED),
        ('🔧', '設定 master UI', '廠 baseline + 類別 + DL/IDL + admin copy + 權限 grant', PURPLE),
        ('🔐', '權限 grant 表', 'user × bg × bu × factory × scope · 對齊既有資料權限管理', ORANGE),
    ]
    for i, (ico, t, sub, c) in enumerate(items):
        y = 4.5 + i * 0.28
        add_text(s, 1.0, y, 0.3, 0.28, ico, size=14, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, 1.4, y, 2.3, 0.28, t, size=11, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, 3.7, y, 8.5, 0.28, sub, size=10, color=RGBColor(0xCB, 0xD5, 0xE1), anchor=MSO_ANCHOR.MIDDLE)

    add_text(s, 0.7, 6.45, 12, 0.4, 'v3 · 2026-06 · 對應 cleansheet-mva-sd v0.4 + Cortex_互動Demo_v0.11.html', size=11, color=MUTED)


def slide_overview(total):
    s = add_slide()
    header(s, 'v3 vs v2 vs v1 差異', '本手冊涵蓋的章節 + 對齊 demo 版本')

    # 版本演進
    add_text(s, 0.4, 1.3, 12, 0.3, '📈 v0.6 → v0.11 規格演進(僅展示與本 v3 手冊相關)',
             size=12, bold=True, color=INK)

    add_rounded(s, 0.4, 1.75, 12.55, 1.8, BG_SOFT, line=NAVY, line_width=1, radius=0.03)
    versions = [
        ('v0.6', '初版 Cleansheet 9×9 matrix · CN 廠 baseline', RGBColor(0x6B, 0x72, 0x80)),
        ('v0.7', '+ qty scenario · MVA 隨 qty 動態變動', RGBColor(0x6B, 0x72, 0x80)),
        ('v0.8', '+ WHOOP 案 · MOUSE_STD vs WHOOP_WEARABLE 兩案', RGBColor(0x6B, 0x72, 0x80)),
        ('v0.9', '+ §MVA 操作流程(Phase A-G)', RGBColor(0x6B, 0x72, 0x80)),
        ('v0.10', '+ Step 1-6 案級操作每步展開 + 5 區 compute trace + 可編輯 + scope toggle', SKY),
        ('v0.11', '⭐ BG/BU 隔離 + 設備類別取代實體 + DL 角色費率 + 設定 master UI + admin copy', RED),
    ]
    for i, (ver, desc, color) in enumerate(versions):
        y = 1.9 + i * 0.27
        add_rounded(s, 0.55, y, 0.7, 0.22, color, radius=0.3)
        add_text(s, 0.55, y, 0.7, 0.22, ver, size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, 1.35, y, 11.5, 0.22, desc, size=10, color=TEXT, anchor=MSO_ANCHOR.MIDDLE)

    # 章節
    add_text(s, 0.4, 3.85, 12, 0.3, '📚 本手冊章節', size=12, bold=True, color=INK)
    sections = [
        ('🔒 BG/BU 隔離', 5, '組織 + 使用者主檔 + 權限 grant + 案綁定', SKY),
        ('🏭 設備類別', 5, '為何改 + 26 類 catalog + admin copy + 案級 binding', CYAN),
        ('👷 DL 角色費率', 3, '7 個 default + 取代廠級單一 wage', RED),
        ('🔧 設定 master UI', 5, '廠 baseline + 7 個 sub-tab + 編輯流程', PURPLE),
        ('📊 公式變更', 2, 'Step 5 區 A DL role-based + 區 C 類別 rep_acq', GREEN),
        ('🎯 Checklist', 2, 'EPM / DPM / admin 各自 checklist', ORANGE),
    ]
    for i, (name, n, desc, color) in enumerate(sections):
        col = i % 3
        row = i // 3
        x = 0.4 + col * 4.25
        y = 4.3 + row * 1.4
        add_rounded(s, x, y, 4.05, 1.25, BG_SOFT, line=color, line_width=1.5, radius=0.03)
        add_rounded(s, x, y, 4.05, 0.32, color, radius=0.03)
        add_text(s, x + 0.15, y + 0.04, 3, 0.27, name, size=11, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + 3, y + 0.04, 1, 0.27, f'{n} slides', size=9, color=WHITE, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT)
        add_text(s, x + 0.15, y + 0.4, 3.8, 0.85, desc, size=10, color=TEXT)

    footer(s, 2, total)


# ────────────────────────────────────────────────────────────────────────
# Section 1: BG/BU 隔離(5 slides)
# ────────────────────────────────────────────────────────────────────────

def slide_bg_intro(total, n):
    s = add_slide()
    header(s, 'BG / BU 組織架構', '對齊既有「資料權限管理」+ 使用者主檔事業群', SKY, '🔒 隔離')

    # 左:組織樹
    add_rounded(s, 0.4, 1.3, 6.2, 5.4, BG_SOFT, line=SKY, line_width=1.5, radius=0.03)
    add_rounded(s, 0.4, 1.3, 6.2, 0.4, SKY, radius=0.03)
    add_text(s, 0.55, 1.35, 5.9, 0.3, '🏢 集團組織架構(demo 4 個 BG)', size=12, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)

    org_tree = """FOXLINK 正崴
├─ 光電事業群 (OPTO) ⭐ SteelSeries 案
│   ├─ 滑鼠 BU (MOUSE)
│   ├─ 鍵盤 BU (KEYBOARD)
│   └─ 線材 BU (CABLE)
│
├─ 消費電子事業群 (CONSUMER) ⭐ WHOOP 案
│   ├─ 可穿戴 BU (WEARABLE)
│   └─ 音訊 BU (AUDIO)
│
├─ 連接器事業群 (CONN)
│   └─ 標準連接器 BU (STD)
│
└─ 中央及貿易 (HQ)
    └─ 資訊工程處 (IT) ← admin 來源

廠區跨 BG 共用:
  CN / VN / TW 三廠 都可服務多個 BG
  但每 BG 有獨立 baseline · 不共用單價"""
    add_text(s, 0.6, 1.8, 5.9, 4.8, org_tree, size=10.5, color=TEXT, font='Consolas')

    # 右:使用者主檔對齊
    add_rounded(s, 6.8, 1.3, 6.15, 5.4, BG_SOFT, line=NAVY, line_width=1.5, radius=0.03)
    add_rounded(s, 6.8, 1.3, 6.15, 0.4, NAVY, radius=0.03)
    add_text(s, 6.95, 1.35, 5.85, 0.3, '👤 使用者主檔(對齊截圖)', size=12, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    user_md = """使用者主檔已有以下 BG/BU 欄位(從 LDAP / ERP 同步):

  ✓ 部門代碼          (7759)
  ✓ 部門名稱          (總管理處-資訊-系統支援部)
  ✓ 利潤中心代碼      (X4)
  ✓ 利潤中心名稱      (資訊工程處)
  ✓ 事業處代碼        (U)
  ✓ 事業處名稱        (中央單位)
  ✓ 事業群名稱        (中央及貿易) ← bg_code
  ✓ 廠區碼            (HQ)
  ✓ 離職日

判斷邏輯:
  • current_user.bg_code → 隱含 filter
  • 同 BG 可看 baseline / 類別單價
  • 跨 BG → 預設 404 / 看不到
  • admin role (HQ + IT)→ 跨 BG 可看可編
  • 細部 grant → bom_settings_admin_grant"""
    add_text(s, 6.95, 1.8, 5.85, 4.8, user_md, size=10, color=TEXT)

    footer(s, n, total)


def slide_bg_isolation(total, n):
    s = add_slide()
    header(s, 'BG 隔離核心:baseline 按 BG 切版', '同廠 (CN) 可有 BG3-baseline + BG5-baseline 各一份', SKY, '🔒 隔離')

    add_text(s, 0.4, 1.3, 12, 0.3, '🎯 設計關鍵 — baseline 隔離 vs 子表加 bg_code 比較',
             size=12, bold=True, color=INK)

    # 兩個方案比較
    add_rounded(s, 0.4, 1.7, 6.2, 5.0, BG_SOFT, line=GREEN, line_width=2, radius=0.03)
    add_rounded(s, 0.4, 1.7, 6.2, 0.45, GREEN, radius=0.03)
    add_text(s, 0.55, 1.75, 5.9, 0.35, '✓ 方案 A(採用):baseline 隔離', size=12, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    plan_a = """ALTER TABLE bom_factory_baseline
  ADD bg_code VARCHAR2(40) NOT NULL,
  ADD bu_code VARCHAR2(40);

UNIQUE (factory_code, bg_code, version_label)

範例:
  CN-OPTO-2026Q2     ← 光電 EPM 看這份
  CN-CONSUMER-2026Q2 ← 消費 EPM 看這份

子表 (price/dl_role/idl_role/consumable)
  保留 ON DELETE CASCADE baseline_id
  → 透過 baseline.bg_code 自動 BG 隔離
  → 子表 schema 完全不變

優點:
  ✓ 子表不需每個都加 bg_code
  ✓ 月度更新各 BG 各更各的
  ✓ 維護成本低
  ✓ 一查 baseline 就明確 BG"""
    add_text(s, 0.55, 2.2, 5.9, 4.4, plan_a, size=10, color=TEXT, font='Consolas')

    # 方案 B 對比
    add_rounded(s, 6.8, 1.7, 6.15, 5.0, BG_SOFT, line=MUTED, line_width=1.5, radius=0.03)
    add_rounded(s, 6.8, 1.7, 6.15, 0.45, MUTED, radius=0.03)
    add_text(s, 6.95, 1.75, 5.85, 0.35, '✗ 方案 B:每個子表加 bg_code', size=12, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    plan_b = """ALTER TABLE bom_factory_equip_category_price
  ADD bg_code VARCHAR2(40);
ALTER TABLE bom_factory_dl_role
  ADD bg_code VARCHAR2(40);
ALTER TABLE bom_factory_idl_role
  ADD bg_code VARCHAR2(40);
ALTER TABLE bom_factory_consumable
  ADD bg_code VARCHAR2(40);

每個子表 query 都要寫 bg_code = ?

缺點:
  ✗ 每個子表都要 ALTER
  ✗ 子表 bg_code 跟 baseline bg_code
    可能不一致(資料不一致風險)
  ✗ 索引變多
  ✗ 維護負擔重複"""
    add_text(s, 6.95, 2.2, 5.85, 4.4, plan_b, size=10, color=TEXT, font='Consolas')

    footer(s, n, total)


def slide_bg_permission(total, n):
    s = add_slide()
    header(s, '權限 grant 表 · bom_settings_admin_grant', 'user × bg/bu/factory × scope × view/edit/approve', SKY, '🔒 隔離')

    add_rounded(s, 0.4, 1.3, 12.55, 5.4, BG_SOFT, line=SKY, line_width=1.5, radius=0.03)
    add_rounded(s, 0.4, 1.3, 12.55, 0.4, SKY, radius=0.03)
    add_text(s, 0.55, 1.35, 12.3, 0.3, '🗄️ 權限 grant 表 schema (對齊既有資料權限管理 pattern)', size=12, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)

    schema = """CREATE TABLE bom_settings_admin_grant (
  grant_id        NUMBER GENERATED ALWAYS AS IDENTITY PRIMARY KEY,
  user_id         NUMBER REFERENCES users(id) NOT NULL,

  -- 範圍三維度(NULL = 不限制)
  bg_code         VARCHAR2(40) REFERENCES org_bg(bg_code) NOT NULL,
  bu_code         VARCHAR2(40) REFERENCES org_bu(bu_code),     -- NULL = 該 BG 全 BU
  factory_code    VARCHAR2(10) REFERENCES bom_factory(factory_code),  -- NULL = 全廠

  scope           VARCHAR2(40) NOT NULL,
  --   'BASELINE'               廠 baseline 主表(切版)
  --   'EQUIP_CATEGORY_PRICE'   設備類別單價
  --   'DL_ROLE'                DL 角色費率
  --   'IDL_ROLE'               IDL 角色費率
  --   'CONSUMABLE'             耗材單價
  --   'EQUIP_CATEGORY_CATALOG' 全 BG 類別 enum(super-admin 級)
  --   'CASE_CLEANSHEET'        案級 Cleansheet 編輯
  --   'ALL'                    全部

  can_view        NUMBER(1) DEFAULT 1,
  can_edit        NUMBER(1) DEFAULT 0,
  can_approve     NUMBER(1) DEFAULT 0,    -- 二簽(漲幅 > 5% 強制)

  granted_by      NUMBER REFERENCES users(id),
  granted_at      TIMESTAMP DEFAULT SYSTIMESTAMP,
  expires_at      TIMESTAMP,              -- NULL = 永久
  is_active       NUMBER(1) DEFAULT 1,
  notes           CLOB
);"""
    add_text(s, 0.6, 1.8, 12.1, 4.9, schema, size=9, color=TEXT, font='Consolas')

    footer(s, n, total)


def slide_bg_permission_examples(total, n):
    s = add_slide()
    header(s, '權限授權範例 · 6 種典型場景', '對應 demo MVA_ADMIN_GRANTS 表', SKY, '🔒 隔離')

    grants = [
        ('Andy', 'OPTO', '—', 'CN', 'BASELINE', '✓ view + edit', '專案授權 2027-12-31'),
        ('Andy', 'OPTO', '—', 'CN', 'EQUIP_CATEGORY_PRICE', '✓ view + edit', '同上'),
        ('Andy', 'OPTO', '—', 'CN', 'DL_ROLE', '✓ view + edit', '同上'),
        ('Andy', 'OPTO', '—', 'CN', 'IDL_ROLE', '只可 view', '永久'),
        ('Mike (DPM)', 'OPTO', 'MOUSE', 'all', 'BASELINE', 'view + approve', '可二簽 · 2026-12-31'),
        ('Lin', 'CONSUMER', 'WEARABLE', 'CN', 'EQUIP_CATEGORY_PRICE', '✓ view + edit', '2027-06-30'),
        ('admin', 'ALL', '—', 'all', 'ALL', '✓ all + approve', '永久 · super-admin'),
    ]

    headers = ['User', 'BG', 'BU', 'Factory', 'Scope', '權限', '備註']
    widths = [1.5, 1.2, 1.2, 1.0, 2.4, 2.0, 3.05]
    cx = 0.4
    for h, w in zip(headers, widths):
        add_rect(s, cx, 1.4, w, 0.4, INK)
        add_text(s, cx, 1.4, w, 0.4, h, size=11, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        cx += w

    for i, row in enumerate(grants):
        y = 1.8 + i * 0.55
        bg = BG_SOFT if i % 2 == 0 else WHITE
        cx = 0.4
        for j, (cell, w) in enumerate(zip(row, widths)):
            add_rect(s, cx, y, w, 0.55, bg)
            color = TEXT
            bold = False
            if j == 4:
                color = PURPLE
                bold = True
            if j == 5 and 'approve' in cell:
                color = RED
                bold = True
            add_text(s, cx + 0.1, y, w - 0.2, 0.55, cell, size=9.5, color=color, bold=bold,
                     font='Consolas' if j in (1,2,3,4) else 'Microsoft JhengHei',
                     anchor=MSO_ANCHOR.MIDDLE)
            cx += w

    # 說明
    add_rounded(s, 0.4, 6.0, 12.55, 0.95, RGBColor(0xDB, 0xEA, 0xFE), line=SKY, radius=0.04)
    add_text(s, 0.55, 6.07, 12.3, 0.9,
        """💡 讀取邏輯:
  • 預設 SELECT 都 filter user.bg_code · 跨 BG 看不到
  • admin role (HQ) 或 scope='ALL' 例外
  • 編輯:每個 INSERT/UPDATE 前 check grant 表 · 對 (bg, bu, factory, scope) 完整 match""",
        size=9.5, color=INK, anchor=MSO_ANCHOR.MIDDLE)

    footer(s, n, total)


def slide_bg_case_binding(total, n):
    s = add_slide()
    header(s, '案綁定 BG/BU 流程', 'BPM 開案時自動帶 user BG · case_factory 同步冗存', SKY, '🔒 隔離')

    add_text(s, 0.4, 1.3, 12, 0.3, '🎬 兩個 demo 案的 BG 綁定',
             size=12, bold=True, color=INK)

    # 兩案案例
    add_rounded(s, 0.4, 1.7, 6.2, 5.0, RGBColor(0xF0, 0xF9, 0xFF), line=SKY, line_width=2, radius=0.03)
    add_rounded(s, 0.4, 1.7, 6.2, 0.45, SKY, radius=0.03)
    add_text(s, 0.55, 1.75, 5.9, 0.35, '🖱️ SteelSeries Rival 3+ Wired', size=12, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    steel = """projects:
  id: 'QT-2026-0148'
  title: 'SteelSeries Rival 3+ Wired Mouse'
  ⭐ bg_code: 'OPTO'      (光電事業群)
  ⭐ bu_code: 'MOUSE'     (滑鼠 BU)

bom_cs_case_factory:
  case_factory_id: 1
  case_id: QT-2026-0148 ↑
  factory_code: 'CN'
  ⭐ bg_code: 'OPTO'      (冗存方便 query)
  ⭐ bu_code: 'MOUSE'
  baseline_id → CN-OPTO-2026Q2

→ 案內 §Cleansheet 看到的:
   • 廠 baseline:CN-OPTO-2026Q2
   • 類別 catalog:OPTO BG 26 類
   • DL 角色:OPTO BG 7 個
   • 跨 BG 看不到 CONSUMER 數字"""
    add_text(s, 0.55, 2.2, 5.9, 4.4, steel, size=10, color=TEXT, font='Consolas')

    add_rounded(s, 6.8, 1.7, 6.15, 5.0, RGBColor(0xFA, 0xF5, 0xFF), line=PURPLE, line_width=2, radius=0.03)
    add_rounded(s, 6.8, 1.7, 6.15, 0.45, PURPLE, radius=0.03)
    add_text(s, 6.95, 1.75, 5.85, 0.35, '⌚ WHOOP Gen4 MP', size=12, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    whoop = """projects:
  id: 'QT-2026-0167'
  title: 'WHOOP Gen4 MP 穿戴手環'
  ⭐ bg_code: 'CONSUMER'   (消費電子事業群)
  ⭐ bu_code: 'WEARABLE'   (可穿戴 BU)

bom_cs_case_factory:
  case_factory_id: 2
  case_id: QT-2026-0167 ↑
  factory_code: 'CN'
  ⭐ bg_code: 'CONSUMER'
  ⭐ bu_code: 'WEARABLE'
  baseline_id → CN-CONSUMER-2026Q2

→ 案內 §Cleansheet 看到的:
   • 廠 baseline:CN-CONSUMER-2026Q2
   • 類別 catalog:CONSUMER BG · 含 FATP 3 類
     (admin 之前從 OPTO copy 24 類過來)
   • DL 角色:CONSUMER BG · wage 略高
   • 看不到 SteelSeries baseline 數字"""
    add_text(s, 6.95, 2.2, 5.85, 4.4, whoop, size=10, color=TEXT, font='Consolas')

    footer(s, n, total)


# ────────────────────────────────────────────────────────────────────────
# Section 2: 設備類別(5 slides)
# ────────────────────────────────────────────────────────────────────────

def slide_cat_why(total, n):
    s = add_slide()
    header(s, '為什麼設備改類別?', '報價階段 EPM 拍板不了用哪一台 · 改成代表機單價', CYAN, '🏭 類別')

    add_rounded(s, 0.4, 1.3, 6.2, 5.4, RGBColor(0xFE, 0xE2, 0xE2), line=RED, line_width=1.5, radius=0.03)
    add_rounded(s, 0.4, 1.3, 6.2, 0.4, RED, radius=0.03)
    add_text(s, 0.55, 1.35, 5.9, 0.3, '❌ v0.10 痛點:細到每一台設備', size=12, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    before = """v0.10 廠級設備 list(28 件實體):

  DEK Horizon 03i Printer      $84,000
  FUJI NXT M3S×6 + M6S×1      $1,013,384
  SPI TRI-7007                 $68,852
  Heller Reflow Mark5          $158,000
  CYBEROPTICS AOI MX600        $92,000
  ...

問題:
  ✗ 報價階段不知道用哪一台
    (DEK 老款 vs 新款差很多)
  ✗ 折舊成本差異大
  ✗ 試產才能定具體機
  ✗ 跨案動態調度
  ✗ 客戶不在乎你用哪台

業界實務:
  TSMC / Foxconn / 鴻海報價 cleansheet
  都用「機台類別 × qty」格式
  不細到 serial number"""
    add_text(s, 0.55, 1.8, 5.9, 4.8, before, size=9.5, color=TEXT, font='Consolas')

    add_rounded(s, 6.8, 1.3, 6.15, 5.4, RGBColor(0xDC, 0xFC, 0xE7), line=GREEN, line_width=1.5, radius=0.03)
    add_rounded(s, 6.8, 1.3, 6.15, 0.4, GREEN, radius=0.03)
    add_text(s, 6.95, 1.35, 5.85, 0.3, '✓ v0.11 解法:類別 × qty', size=12, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    after = """v0.11 案級設備配置(類別 × qty):

  SMT_LINE_FULL          × 1.0  ($1,200,000)
  SMT_AOI                × 0.5  ($100,000 · shared)
  INJ_HORIZ_100T         × 2.0  ($95,000 each)
  INJ_HORIZ_200T         × 1.0  ($150,000)
  BB_ROBOT               × 1.0  ($50,000)
  ICT                    × 1.0  ($50,000)
  EMI_CHAMBER            × 0.3  ($480,000 · shared 30%)
  ...

優點:
  ✓ 報價階段 EPM 拍類別合理
  ✓ 不糾結具體哪台
  ✓ 試產才細化到具體機(可選)
  ✓ qty 支援小數(跨案共用)
  ✓ 跨 BG / 跨案 utilization 看得清

設計考量:
  廠級資產帳(個別機)該在 ERP fixed_assets
  不該跟 Cleansheet 模組混"""
    add_text(s, 6.95, 1.8, 5.85, 4.8, after, size=9.5, color=TEXT, font='Consolas')

    footer(s, n, total)


def slide_cat_enum(total, n):
    s = add_slide()
    header(s, '設備類別 enum · 26 項(OPTO BG)', '6 個 group · 對應實際製造業 setup', CYAN, '🏭 類別')

    add_text(s, 0.4, 1.3, 12, 0.3, '📋 OPTO 光電 BG 設備類別 catalog · 26 類分 6 大 group',
             size=12, bold=True, color=INK)

    groups = [
        ('SMT', SKY, [
            ('SMT_LINE_FULL',  'SMT 整線(含周邊)',   '$1,200,000', '6yr'),
            ('SMT_PRINTER',    'SMT 印刷機',          '$90,000',    '6yr'),
            ('SMT_PNP_LOW',    'SMT 貼片機 低階',     '$400,000',   '6yr'),
            ('SMT_PNP_HIGH',   'SMT 貼片機 高階',     '$1,100,000', '6yr'),
            ('SMT_REFLOW',     '回焊爐',              '$160,000',   '6yr'),
            ('SMT_AOI',        'AOI 檢測機',          '$100,000',   '6yr'),
        ]),
        ('PCBA', CYAN, [
            ('WAVE_SOLDER',    '波焊機',              '$140,000',   '8yr'),
            ('ROUTER',         '分板機',              '$50,000',    '6yr'),
            ('LASER_ETCH',     '雷射雕刻機',          '$80,000',    '8yr'),
        ]),
        ('INJECTION', GREEN, [
            ('INJ_HORIZ_50T',  '成形機 臥式 50T',     '$60,000',    '8yr'),
            ('INJ_HORIZ_100T', '成形機 臥式 100T',    '$95,000',    '8yr'),
            ('INJ_HORIZ_200T', '成形機 臥式 200T',    '$150,000',   '8yr'),
            ('INJ_HORIZ_350T', '成形機 臥式 350T',    '$260,000',   '8yr'),
            ('INJ_VERT_50T',   '成形機 立式 50T',     '$55,000',    '8yr'),
            ('INJ_VERT_100T',  '成形機 立式 100T',    '$85,000',    '8yr'),
        ]),
        ('ASSEMBLY', PURPLE, [
            ('BB_ROBOT',       '組裝 SCARA 機械手臂', '$50,000',    '6yr'),
            ('BB_AUTO_SCREW',  '自動鎖螺絲機',        '$10,000',    '5yr'),
            ('BB_CONVEYOR',    '組裝輸送帶',          '$20,000',    '8yr'),
            ('BB_GLUE_DISP',   '自動點膠機',          '$20,000',    '5yr'),
            ('BB_HEAT_WELD',   '熱熔焊接',            '$9,000',     '5yr'),
        ]),
        ('TEST + QUALITY', RED, [
            ('ICT',            'ICT 測試機',          '$50,000',    '6yr'),
            ('FCT_RIG',        'Functional 治具',     '$7,000',     '3yr'),
            ('XRAY',           'X-Ray 檢測',          '$250,000',   '8yr'),
            ('EMI_CHAMBER',    'EMI/EMC 室',          '$480,000',   '10yr'),
            ('THERMAL_CHAMBER','熱循環室',            '$85,000',    '10yr'),
            ('VIB_SHAKER',     '振動測試',            '$128,000',   '10yr'),
        ]),
    ]

    # 3 column layout
    col_h = 5.5
    col_w = 4.15
    y_start = 1.75
    items_flat = []
    for g, color, items in groups:
        items_flat.append(('HEADER', g, color, None, None))
        items_flat.extend([(i[0], i[1], color, i[2], i[3]) for i in items])

    # spread across 3 columns
    items_per_col = (len(items_flat) + 2) // 3
    for col in range(3):
        x = 0.4 + col * (col_w + 0.05)
        for i in range(items_per_col):
            idx = col * items_per_col + i
            if idx >= len(items_flat):
                break
            item = items_flat[idx]
            y = y_start + i * 0.265
            if item[0] == 'HEADER':
                add_rect(s, x, y, col_w, 0.25, item[2])
                add_text(s, x + 0.1, y, col_w - 0.2, 0.25, f'  {item[1]}', size=11, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
            else:
                code, name, color, price, life = item
                add_rect(s, x, y, col_w, 0.25, BG_SOFT if i % 2 == 0 else WHITE)
                add_text(s, x + 0.1, y, 1.7, 0.25, code, size=8, color=color, font='Consolas', anchor=MSO_ANCHOR.MIDDLE, bold=True)
                add_text(s, x + 1.85, y, 1.4, 0.25, name, size=8.5, color=TEXT, anchor=MSO_ANCHOR.MIDDLE)
                add_text(s, x + 3.2, y, 0.75, 0.25, price, size=8, color=NAVY, font='Consolas', anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT)
                add_text(s, x + 3.85, y, 0.3, 0.25, life, size=7.5, color=MUTED, anchor=MSO_ANCHOR.MIDDLE)

    footer(s, n, total)


def slide_cat_copy(total, n):
    s = add_slide()
    header(s, 'admin 跨 BG copy 類別流程', '加速新 BG onboarding · 從 OPTO copy 24 類到 CONSUMER', CYAN, '🏭 類別')

    add_text(s, 0.4, 1.3, 12, 0.3, '🎬 操作流程(以 CONSUMER BG 初始化為例)',
             size=12, bold=True, color=INK)

    steps = [
        ('1️⃣', 'admin 進「設定 master UI」→ 設備類別 catalog',
            '從 §Cleansheet > Step 7 設定 master UI > 子 tab「🏭 設備類別 catalog」'),
        ('2️⃣', '右上「📋 從來源 BG copy 類別」按鈕',
            '只有 admin role (HQ + IT) 看得到 · 一般 EPM 不會看到'),
        ('3️⃣', 'Modal 開啟 · 選來源 BG + 目標 BG',
            '來源:OPTO 光電(已有 26 類)· 目標:CONSUMER 消費電子(空 / 部分)'),
        ('4️⃣', '勾選要 copy 的類別(全勾 / 部分勾)',
            'CONSUMER 沒用大型成形機 → 略過 INJ_HORIZ_200T/350T\n保留 SMT/PCBA/小型 INJ/Assembly/Test/Quality 共 24 類'),
        ('5️⃣', '系統 INSERT × N rows(複製 + 標 copied_from)',
            "INSERT INTO bom_equip_category_catalog\n  (bg_code='CONSUMER', category_code, display_name, group,\n   copied_from_bg_code='OPTO', copied_from_category_id, ...)"),
        ('6️⃣', 'CONSUMER EPM 加 wearable-only 類別',
            '手動新建 3 類:FATP_PCBA_TEST · FATP_HR_CAL · FATP_FINAL_TEST\n標 restrict_bg_codes=\'CONSUMER\' 軟限制其他 BG 看不到'),
        ('7️⃣', '廠單價:CONSUMER × CN 廠 baseline 配置',
            'INSERT bom_factory_equip_category_price × 24 + 3 = 27 rows\nrep_acq_cost / life / mro_pct 各設'),
    ]
    for i, (n_, t, d) in enumerate(steps):
        y = 1.75 + i * 0.72
        add_rounded(s, 0.4, y, 12.55, 0.66, BG_SOFT, line=CYAN, line_width=0.5, radius=0.04)
        add_text(s, 0.5, y + 0.16, 0.5, 0.35, n_, size=18, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, 1.1, y + 0.05, 5.5, 0.27, t, size=11, bold=True, color=INK, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, 6.7, y + 0.05, 6.2, 0.56, d, size=9, color=TEXT, anchor=MSO_ANCHOR.TOP, font='Consolas')

    footer(s, n, total)


def slide_cat_case_binding(total, n):
    s = add_slide()
    header(s, '案級 binding 操作 · SteelSeries 案 SMT 製程', '從廠 catalog 挑類別 × qty · 設 override', CYAN, '🏭 類別')

    add_text(s, 0.4, 1.3, 12, 0.3, '🎬 操作流程(SteelSeries SMT 製程)',
             size=12, bold=True, color=INK)

    mockup = """┌─ §Cleansheet → Step 3 · 設備類別 binding ──────────────────────────────────────┐
│   案內 14 類 · 對應 bom_cs_case_equip_category · BG: 光電事業群(OPTO)           │
│                                                                                │
│   ⚡ v0.11 規格修正:設備從個別 DEK / FUJI 等實體改為類別 + 代表機單價          │
│                                                                                │
│   [📋 案級 binding (14 類)]  [🏭 BG catalog (26 類)]  [📋 admin 跨 BG copy]    │
│                                                                                │
│   類別                              Group  Process    qty   代表單價    年折舊+MRO│
│   ──────────────────────────────────────────────────────────────────────────  │
│   SMT_LINE_FULL  SMT 整線(含周邊)  SMT    SMT_MAIN  1.0   $1,200,000   $260k │
│   SMT_AOI        AOI 檢測機         SMT    SMT_MAIN  0.5   $100,000     $10.8k│
│   WAVE_SOLDER    波焊機             PCBA   WAVE      1.0   $140,000     $24.5k│
│   ROUTER         分板機             PCBA   ROUTER    1.0   $50,000      $10.8k│
│   LASER_ETCH     雷射雕刻機         PCBA   LASER     1.0   $80,000      $14k  │
│   INJ_HORIZ_100T 成形機 臥式 100T   INJ    BB_ASSY   2.0   $95,000      $33.3k│ ← 2 台
│   INJ_HORIZ_200T 成形機 臥式 200T   INJ    BB_ASSY   1.0   $150,000     $26.3k│
│   BB_ROBOT       SCARA 機械手臂     ASSY   BB_ASSY   1.0   $50,000      $10.8k│
│   ICT            ICT 測試機         TEST   BB_TEST   1.0   $50,000      $10.8k│
│   EMI_CHAMBER    EMI/EMC 室         TEST   BB_TEST   0.3   $480,000     $48.6k│ ← 跨案 30%
│   XRAY           X-Ray 檢測         QUAL   Q_SMT     0.5   $250,000     $25.6k│ ← 跨案 50%
│   ...                                                                          │
│                                                                                │
│   ↳ 寫 schema: INSERT bom_cs_case_equip_category × 14 rows                     │
└────────────────────────────────────────────────────────────────────────────────┘"""
    add_rounded(s, 0.4, 1.7, 12.55, 5.0, WHITE, line=CYAN, line_width=1.5, radius=0.02)
    add_text(s, 0.55, 1.85, 12.3, 4.85, mockup, size=8, color=TEXT, font='Consolas')

    footer(s, n, total)


def slide_cat_schema(total, n):
    s = add_slide()
    header(s, '設備類別 schema · 三層架構', 'catalog enum + 廠級單價 + 案級 binding', CYAN, '🏭 類別')

    add_rounded(s, 0.4, 1.3, 12.55, 5.4, BG_SOFT, line=CYAN, line_width=1.5, radius=0.03)
    add_rounded(s, 0.4, 1.3, 12.55, 0.4, CYAN, radius=0.03)
    add_text(s, 0.55, 1.35, 12.3, 0.3, '🗄️ schema (3 表 · v0.11 新增)', size=12, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)

    schema = """-- (1) BG 級 catalog enum(BG 各自一份 + admin 可跨 BG copy)
CREATE TABLE bom_equip_category_catalog (
  category_id            NUMBER GENERATED ALWAYS AS IDENTITY PRIMARY KEY,
  ⭐ bg_code              VARCHAR2(40) REFERENCES org_bg(bg_code) NOT NULL,
  bu_code                VARCHAR2(40) REFERENCES org_bu(bu_code),
  category_code          VARCHAR2(80),                            -- 'SMT_LINE_FULL' / 'INJ_HORIZ_100T'
  display_name_zh/en/vi  VARCHAR2(200),
  category_group         VARCHAR2(40),                            -- SMT/PCBA/INJECTION/...
  default_process_code   VARCHAR2(40),
  restrict_bg_codes      VARCHAR2(500),                           -- 軟過濾('CONSUMER,OPTO')
  ⭐ copied_from_bg_code  VARCHAR2(40),                            -- admin copy 來源
  ⭐ copied_from_category_id NUMBER,
  ⭐ created_by/by_at, updated_by/_at NUMBER, TIMESTAMP,
  CONSTRAINT becc_uk UNIQUE (bg_code, bu_code, category_code)
);

-- (2) 廠級類別代表單價(透過 baseline_id 自動 BG 隔離)
CREATE TABLE bom_factory_equip_category_price (
  price_id               NUMBER PRIMARY KEY,
  baseline_id            NUMBER REFERENCES bom_factory_baseline(baseline_id) ON DELETE CASCADE,
  category_code          VARCHAR2(80) REFERENCES bom_equip_category_catalog(category_code),
  rep_acq_cost_usd       NUMBER(15,2),                            -- 該廠該類別「代表機」單價
  rep_useful_life_yrs    NUMBER(5,2),
  rep_mro_pct            NUMBER(5,4) DEFAULT 0.05,
  created_by/updated_by NUMBER, TIMESTAMP, ...,
  CONSTRAINT bfecp_uk UNIQUE (baseline_id, category_code)
);

-- (3) 案級 binding(case_factory × category × process × qty)
CREATE TABLE bom_cs_case_equip_category (
  case_factory_id        NUMBER REFERENCES bom_cs_case_factory(case_factory_id) ON DELETE CASCADE,
  category_code          VARCHAR2(80) REFERENCES bom_equip_category_catalog(category_code),
  process_code           VARCHAR2(40) REFERENCES bom_process_catalog(process_code),
  qty                    NUMBER(5,2),                             -- 1.0 / 0.5 跨案共用
  acq_cost_override_usd  NUMBER(15,2),                            -- 議價 / 客戶提供 / 客供 = 0
  useful_life_override_yrs NUMBER(5,2),
  created_by/at NUMBER, TIMESTAMP, ...
  PRIMARY KEY (case_factory_id, category_code, process_code)
);"""
    add_text(s, 0.55, 1.78, 12.3, 4.9, schema, size=8, color=TEXT, font='Consolas')

    footer(s, n, total)


# ────────────────────────────────────────────────────────────────────────
# Section 3: DL 角色費率(3 slides)
# ────────────────────────────────────────────────────────────────────────

def slide_dl_why(total, n):
    s = add_slide()
    header(s, '為什麼 DL 改角色費率?', '取代廠級單一 $4.95/hr · 對齊 IDL 角色 pattern', RED, '👷 DL 角色')

    add_rounded(s, 0.4, 1.3, 6.2, 5.4, RGBColor(0xFE, 0xE2, 0xE2), line=RED, line_width=1.5, radius=0.03)
    add_rounded(s, 0.4, 1.3, 6.2, 0.4, RED, radius=0.03)
    add_text(s, 0.55, 1.35, 5.9, 0.3, '❌ v0.10 設計遺漏', size=12, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    before = """v0.10 DL 設計:
  bom_factory_baseline.dl_wage_per_hr_usd = $4.95

  → 所有製程的所有 DL 都用同一個 $4.95/hr
  → 跟 IDL 17 角色完全不對等

實務問題:
  • SMT 一般操作員 $4.95/hr ✓
  • BB Test Debug 操作員 $8.50/hr(懂電路) ✗
  • Functional 校正員 $12.00/hr(資深) ✗
  • 無塵室 SMT $6.50/hr(工時嚴) ✗
  • 模具操作員(成形機)$5.50/hr ✗

v0.10 強制都用 $4.95 → MVA 算出來不準

對比 IDL(17 角色):
  OPS_MGR $60,320
  SEC_MGR $42,500
  ENGINEER $24,440
  ...
  → IDL 分得細 · DL 卻只有一個 wage
  → 架構不對稱"""
    add_text(s, 0.55, 1.8, 5.9, 4.8, before, size=10, color=TEXT, font='Consolas')

    add_rounded(s, 6.8, 1.3, 6.15, 5.4, RGBColor(0xDC, 0xFC, 0xE7), line=GREEN, line_width=1.5, radius=0.03)
    add_rounded(s, 6.8, 1.3, 6.15, 0.4, GREEN, radius=0.03)
    add_text(s, 6.95, 1.35, 5.85, 0.3, '✓ v0.11 解法:7 個 DL 角色', size=12, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    after = """v0.11 DL 角色 master(對應 bom_factory_dl_role):

  role_code         name              hourly_usd
  ─────────────────────────────────────────────
  OP_GENERAL        一般操作員        $4.95
  OP_DEBUG          Debug 操作員      $8.50
  OP_FUNCT          Functional 操作員 $12.00
  OP_CLEANROOM      無塵室操作員      $6.50
  OP_MOLD           模具操作員        $5.50
  OP_TEST           測試操作員        $7.00
  OP_AUTO_MON       自動化監控員      $6.00

完全對齊 IDL pattern · BG 各自一份 enum
  • OPTO BG 7 個 default
  • CONSUMER BG copy 後 5 個(wage 略高)

優點:
  ✓ 跟 IDL 17 角色架構一致
  ✓ 跨廠 wage 對比(VN $2.50 vs TW $9.20)
  ✓ 月度漲薪按工種(無塵室漲 5%)
  ✓ 計算引擎更精準

舊 dl_wage_per_hr_usd 保留作 fallback(deprecated)"""
    add_text(s, 6.95, 1.8, 5.85, 4.8, after, size=10, color=TEXT, font='Consolas')

    footer(s, n, total)


def slide_dl_case_config(total, n):
    s = add_slide()
    header(s, '案級 DL 配置 · per process × role × count', '取代 v0.10 三個固定欄位(dl_per_shift 等)', RED, '👷 DL 角色')

    add_text(s, 0.4, 1.3, 12, 0.3, '🎬 SteelSeries 案 · CN 廠 · 各製程 DL 配置(role × count)',
             size=12, bold=True, color=INK)

    config = """┌─────────────────────────────────────────────────────────────────────────┐
│  Process       Role             count_per_shift  rate/hr  /day cost     │
│  ─────────────────────────────────────────────────────────────────────  │
│  SMT_MAIN     OP_GENERAL          8 人           $4.95    $79.20        │
│               OP_CLEANROOM        2 人           $6.50    $26.00        │  ← 無塵段
│               OP_DEBUG            3 人           $8.50    $51.00        │
│               OP_FUNCT            2 人           $12.00   $48.00        │
│                                                          ─────────       │
│                                                          $204.20 /shift  │
│                                                                          │
│  BB_ASSY      OP_GENERAL         18 人           $4.95    $178.20        │
│               OP_MOLD             4 人           $5.50    $44.00        │  ← 成形機
│               OP_DEBUG            5 人           $8.50    $85.00        │
│               OP_FUNCT            3 人           $12.00   $72.00        │
│                                                          ─────────       │
│                                                          $379.20 /shift  │
│                                                                          │
│  BB_TEST      OP_TEST             8 人           $7.00    $112.00        │
│               OP_DEBUG            2 人           $8.50    $34.00        │
│               OP_FUNCT            1 人           $12.00   $24.00        │
│                                                                          │
│  WAVE_SOLDER  OP_GENERAL          3 人           $4.95    $29.70        │
│               OP_FUNCT            1 人           $12.00   $24.00        │
│  ...                                                                     │
└─────────────────────────────────────────────────────────────────────────┘

→ 寫入 schema:
  INSERT INTO bom_cs_case_process_dl
    (case_factory_id, process_code, dl_role_code, count_per_shift) × N rows"""
    add_rounded(s, 0.4, 1.7, 12.55, 5.0, WHITE, line=RED, line_width=1.5, radius=0.02)
    add_text(s, 0.55, 1.85, 12.3, 4.85, config, size=8.5, color=TEXT, font='Consolas')

    footer(s, n, total)


def slide_dl_schema(total, n):
    s = add_slide()
    header(s, 'DL 角色 schema · 對齊 IDL pattern', '廠級 master + 案級 binding 兩層', RED, '👷 DL 角色')

    add_rounded(s, 0.4, 1.3, 12.55, 5.4, BG_SOFT, line=RED, line_width=1.5, radius=0.03)
    add_rounded(s, 0.4, 1.3, 12.55, 0.4, RED, radius=0.03)
    add_text(s, 0.55, 1.35, 12.3, 0.3, '🗄️ schema (2 表 · v0.11 新增)', size=12, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)

    schema = """-- (1) 廠級 DL 角色 master(對應 bom_factory_idl_role · 對等架構)
CREATE TABLE bom_factory_dl_role (
  role_id              NUMBER GENERATED ALWAYS AS IDENTITY PRIMARY KEY,
  baseline_id          NUMBER REFERENCES bom_factory_baseline(baseline_id) ON DELETE CASCADE,
  role_code            VARCHAR2(40),                  -- 'OP_GENERAL' / 'OP_DEBUG' / ...
  display_name_zh/en   VARCHAR2(100),
  category             VARCHAR2(40),                  -- GENERAL/DEBUG/FUNCT/CLEANROOM/MOLD/TEST/AUTO_MON
  hourly_rate_usd      NUMBER(10,4),                  -- $4.95 / $8.50 / $12.00
  copied_from          VARCHAR2(40),                  -- admin copy 來源 BG (e.g. 'OPTO')
  notes                CLOB,
  ⭐ created_by/at, updated_by/at  -- audit metadata 5 件套
  CONSTRAINT bfdr_uk UNIQUE (baseline_id, role_code)
);

-- (2) 案級 DL 配置(取代 v0.10 case_process 三欄 dl_per_shift / debug_dl_per_shift / functional_dl_per_shift)
CREATE TABLE bom_cs_case_process_dl (
  case_factory_id      NUMBER REFERENCES bom_cs_case_factory(case_factory_id) ON DELETE CASCADE,
  process_code         VARCHAR2(40) REFERENCES bom_process_catalog(process_code),
  dl_role_code         VARCHAR2(40) NOT NULL,         -- 對 bom_factory_dl_role.role_code
  count_per_shift      NUMBER(5,2),                   -- 人數(可小數 跨製程共用)
  count_per_debug_line NUMBER(5,2) DEFAULT 0,         -- debug 線額外人數
  is_functional        NUMBER(1) DEFAULT 0,
  created_by/at NUMBER, TIMESTAMP,
  PRIMARY KEY (case_factory_id, process_code, dl_role_code)
);

-- (3) 計算引擎 Step 5 區 A · DL cost 改 role-based
SELECT SUM(
  COALESCE(role.hourly_rate_usd, baseline.dl_wage_per_hr_usd /* fallback */) ×
  case_dl.count_per_shift × lines × shifts × multipliers × SEA_hr
) AS dl_cost_per_week
FROM bom_cs_case_process_dl case_dl
JOIN bom_factory_dl_role role ON role.role_code = case_dl.dl_role_code
                             AND role.baseline_id = case_factory.baseline_id
JOIN bom_cs_case_factory case_factory ON case_factory.case_factory_id = case_dl.case_factory_id
WHERE case_dl.case_factory_id = ? AND case_dl.process_code = ?;
-- → / weekly_output → DL cost per unit"""
    add_text(s, 0.55, 1.78, 12.3, 4.9, schema, size=8, color=TEXT, font='Consolas')

    footer(s, n, total)


# ────────────────────────────────────────────────────────────────────────
# Section 4: 設定 master UI(5 slides)
# ────────────────────────────────────────────────────────────────────────

def slide_settings_intro(total, n):
    s = add_slide()
    header(s, '設定 master UI · 7 個子 tab', '從 §Cleansheet > Step 7 進入 · admin/EPM 各看自己 BG', PURPLE, '🔧 設定')

    add_text(s, 0.4, 1.3, 12, 0.3, '🎬 進入路徑 + 7 個子 tab 一覽',
             size=12, bold=True, color=INK)

    subtabs = [
        ('📅 廠 baseline 切版',  'SCD Type 2 · 月度更新 · status: active/superseded/draft'),
        ('🏭 設備類別 catalog',  'BG 各管 + admin copy(從 OPTO copy 到 CONSUMER)'),
        ('💰 類別代表單價',      'per baseline × 類別 · vs 上版漲跌幅 · 編 / 二簽'),
        ('👷 DL 角色費率',       '7 個 default + BG 各自費率 · 月度漲薪入口'),
        ('👔 IDL 角色費率',      '17 個 default · 沿用 v0.10 schema + audit 5 件套'),
        ('📦 耗材單價',          '17 件耗材 · per baseline · 漲價即時 reprice'),
        ('🔐 權限 grant 管理',   '對齊既有資料權限管理 pattern · scope × bg × bu × factory'),
    ]
    for i, (label, desc) in enumerate(subtabs):
        y = 1.75 + i * 0.7
        add_rounded(s, 0.4, y, 12.55, 0.64, BG_SOFT, line=PURPLE, line_width=0.5, radius=0.04)
        add_rounded(s, 0.55, y + 0.1, 0.45, 0.44, PURPLE, radius=0.5)
        add_text(s, 0.55, y + 0.1, 0.45, 0.44, str(i+1), size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, 1.15, y + 0.05, 4, 0.27, label, size=12, bold=True, color=INK, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, 5.2, y + 0.05, 7.5, 0.54, desc, size=10, color=TEXT, anchor=MSO_ANCHOR.MIDDLE)

    footer(s, n, total)


def slide_settings_baseline(total, n):
    s = add_slide()
    header(s, '子 tab 1 · 廠 baseline 切版', '同廠不同 BG 各有自己 baseline · SCD Type 2', PURPLE, '🔧 設定')

    mockup = """┌─ 📅 廠 baseline 切版歷史 · 光電事業群 × CN 廠 ─────────────────────────────────┐
│                                                                               │
│   Version              Status        Effective              DL wage    SGA/Profit   Created by   Actions │
│   ────────────────────────────────────────────────────────────────────────────  │
│   CN-OPTO-2026Q2      ✓ ACTIVE      2026-04-01 ~ NULL      $4.95/hr   2.0%/14.0%   Andy (A5566) [查][編]  │
│   CN-OPTO-2026Q1      SUPERSEDED    2026-01-01 ~ 2026-03   $4.85/hr   2.0%/14.0%   Andy           已封存   │
│   CN-OPTO-2025Q4      SUPERSEDED    2025-10-01 ~ 2025-12   $4.75/hr   2.0%/14.0%   Andy           已封存   │
│                                                                               │
│   ↳ schema: bom_factory_baseline                                             │
│            (factory_code='CN', ⭐ bg_code='OPTO', version_label='...',         │
│             status, effective_from, effective_to, created_by, updated_at, ...)│
└───────────────────────────────────────────────────────────────────────────────┘

切版動作(權限 BASELINE × can_edit + can_approve):

   1. Andy 編 CN-OPTO-2026Q2 → 改 DL wage $4.95 → $5.20(+5.05%)
   2. 漲幅 > 5% → 系統 block · 需 DPM Mike 二簽(can_approve)
   3. Mike approve → 系統:
      ・舊版 CN-OPTO-2026Q1.status = superseded, effective_to = today
      ・新版 CN-OPTO-2026Q2.status = active, effective_from = tomorrow
      ・user_notifications 推給該 BG 所有 draft case 的 EPM
        「baseline 已切版,點此 reprice」"""
    add_rounded(s, 0.4, 1.3, 12.55, 5.6, WHITE, line=PURPLE, line_width=1.5, radius=0.02)
    add_text(s, 0.55, 1.45, 12.3, 5.4, mockup, size=8.5, color=TEXT, font='Consolas')

    footer(s, n, total)


def slide_settings_equip_price(total, n):
    s = add_slide()
    header(s, '子 tab 3 · 類別代表單價', '光電 BG × CN 廠 × 當前 active baseline · vs 上版漲跌幅', PURPLE, '🔧 設定')

    mockup = """┌─ 💰 OPTO × CN × CN-OPTO-2026Q2 baseline · 類別代表單價 ──────────────────────┐
│                                                                             │
│   類別               rep_acq_cost   life   mro_pct   vs 上版    Action     │
│   ─────────────────────────────────────────────────────────────────────    │
│   SMT_LINE_FULL      $1,200,000     6yr    5.0%      ↑ 2.5%     [編]      │
│   SMT_PRINTER        $90,000        6yr    5.0%      ↑ 1.1%     [編]      │
│   INJ_HORIZ_100T     $95,000        8yr    5.0%      ↑ 3.3%     [編]      │
│   INJ_HORIZ_200T     $150,000       8yr    5.0%      —          [編]      │
│   EMI_CHAMBER        $480,000       10yr   3.0%      ↓ 1.0%     [編]      │
│   XRAY               $250,000       8yr    4.0%      ↑ 0.5%     [編]      │
│   ...                                                                       │
│                                                                             │
│   ↳ schema: bom_factory_equip_category_price                               │
│            (baseline_id → 對應 BG · 自動隔離,                              │
│             category_code, rep_acq_cost_usd, rep_useful_life_yrs,          │
│             rep_mro_pct, updated_by, updated_at, ...)                      │
└─────────────────────────────────────────────────────────────────────────────┘

編輯流程:

  1. CN EPM Andy 編 SMT_LINE_FULL → $1,200,000 → $1,250,000(+4.2%)
  2. 系統 check grant 表:
     SELECT * FROM bom_settings_admin_grant
     WHERE user_id = Andy
       AND scope IN ('EQUIP_CATEGORY_PRICE', 'BASELINE', 'ALL')
       AND bg_code = 'OPTO' AND factory_code = 'CN'
       AND is_active = 1 AND (expires_at IS NULL OR expires_at > now())
       AND can_edit = 1
  3. PASS → UPDATE bom_factory_equip_category_price
  4. 寫 audit log: bom_settings_audit_log
     (table='bom_factory_equip_category_price', record_id, old_value, new_value, changed_by, ...)
  5. 漲幅 ≥ 5% → 觸發 DPM 二簽流程"""
    add_rounded(s, 0.4, 1.3, 12.55, 5.6, WHITE, line=PURPLE, line_width=1.5, radius=0.02)
    add_text(s, 0.55, 1.45, 12.3, 5.4, mockup, size=8.5, color=TEXT, font='Consolas')

    footer(s, n, total)


def slide_settings_dl_role(total, n):
    s = add_slide()
    header(s, '子 tab 4 · DL 角色費率', 'OPTO BG 7 個 default · 月度漲薪入口', PURPLE, '🔧 設定')

    mockup = """┌─ 👷 OPTO × CN × CN-OPTO-2026Q2 · DL 角色費率 ─────────────────────────────┐
│                                                                          │
│   role_code         name                 category    hourly_usd  來源    │
│   ──────────────────────────────────────────────────────────────────────  │
│   OP_GENERAL        一般操作員           GENERAL     $4.95       原生     │
│   OP_DEBUG          Debug 操作員         DEBUG       $8.50       原生     │
│   OP_FUNCT          Functional 操作員    FUNCT       $12.00      原生     │
│   OP_CLEANROOM      無塵室操作員         CLEANROOM   $6.50       原生     │
│   OP_MOLD           模具操作員           MOLD        $5.50       原生     │
│   OP_TEST           測試操作員           TEST        $7.00       原生     │
│   OP_AUTO_MON       自動化監控員         AUTO_MON    $6.00       原生     │
│                                                                          │
│   ↳ schema: bom_factory_dl_role                                          │
│            (baseline_id, role_code, display_name_zh,                     │
│             category, hourly_rate_usd, copied_from,                      │
│             created_by, updated_by, ...)                                 │
└──────────────────────────────────────────────────────────────────────────┘

對比 CONSUMER BG · CN 廠(copied 過來 · wage 略高):

   OP_GENERAL    $5.20   (vs OPTO $4.95 · +5.1%)
   OP_DEBUG      $8.90   (vs OPTO $8.50 · +4.7%)
   OP_FUNCT      $12.50  (vs OPTO $12.00 · +4.2%)
   OP_CLEANROOM  $6.80   (vs OPTO $6.50 · +4.6%)
   OP_TEST       $7.40   (vs OPTO $7.00 · +5.7%)
   ❌ 無 OP_MOLD (消費電子不用大量成形機 · 不設此 role)
   ❌ 無 OP_AUTO_MON

→ 同廠不同 BG 議價量 / wage 算法略異 · 這就是 BG 隔離的價值"""
    add_rounded(s, 0.4, 1.3, 12.55, 5.6, WHITE, line=PURPLE, line_width=1.5, radius=0.02)
    add_text(s, 0.55, 1.45, 12.3, 5.4, mockup, size=8.5, color=TEXT, font='Consolas')

    footer(s, n, total)


def slide_settings_permission(total, n):
    s = add_slide()
    header(s, '子 tab 7 · 權限 grant 管理', '對齊既有「資料權限管理」3 層架構', PURPLE, '🔧 設定')

    mockup = """┌─ 🔐 bom_settings_admin_grant · 維護權限授權表 ──────────────────────────────────────┐
│                                                                                 │
│   user           BG         BU         factory   scope                  view edit approve   expires    │
│   ──────────────────────────────────────────────────────────────────────────────────────── │
│   Andy           OPTO       (全 BU)    CN        BASELINE                ✓    ✓    —       2027-12-31│
│   Andy           OPTO       (全 BU)    CN        EQUIP_CATEGORY_PRICE    ✓    ✓    —       同上       │
│   Andy           OPTO       (全 BU)    CN        DL_ROLE                 ✓    ✓    —       同上       │
│   Andy           OPTO       (全 BU)    CN        IDL_ROLE                ✓    —    —       永久       │
│   Mike (DPM)     OPTO       MOUSE      (全廠)    BASELINE                ✓    —    ✓       2026-12-31 │
│   Lin            CONSUMER   WEARABLE   CN        EQUIP_CATEGORY_PRICE    ✓    ✓    —       2027-06-30 │
│   admin          ALL        —          (all)     ALL                     ✓    ✓    ✓       永久       │
│                                                                                 │
└─────────────────────────────────────────────────────────────────────────────────┘

對比既有「資料權限管理」(從截圖):

  既有政策(包含/排除):              v0.11 grant 對應:
  • 預設使用者                         user_id × scope='ALL'
  • 合理庫存專用權限                   user_id × scope='SPECIFIC'
  • IT超級使用者                       admin × ALL × all
  • 光電經營資料權限                   bg='OPTO' × scope='ALL'
  • LW利潤中心權限                     利潤中心 → bu_code

  既有層級:                            v0.11 對應:
  • 部門 / 利潤中心 / 事業處 / 事業群    bg_code (從 user.事業群名稱 對)
  • 組織代碼 (ORG_CODE)                bu_code (從 user.利潤中心代碼 對)
  • 包含 / 排除規則                    is_active = 1/0(or 加 exclude 旗標)"""
    add_rounded(s, 0.4, 1.3, 12.55, 5.6, WHITE, line=PURPLE, line_width=1.5, radius=0.02)
    add_text(s, 0.55, 1.45, 12.3, 5.4, mockup, size=8, color=TEXT, font='Consolas')

    footer(s, n, total)


# ────────────────────────────────────────────────────────────────────────
# Section 5: 公式變更(2 slides)
# ────────────────────────────────────────────────────────────────────────

def slide_formula_dl(total, n):
    s = add_slide()
    header(s, 'Step 5 區 A DL 公式 · role-based 改寫', '取代 v0.10 單一 DL_wage × Total_DL_day', GREEN, '📊 公式')

    add_rounded(s, 0.4, 1.3, 6.2, 5.4, RGBColor(0xFE, 0xE2, 0xE2), line=RED, line_width=1.5, radius=0.03)
    add_rounded(s, 0.4, 1.3, 6.2, 0.4, RED, radius=0.03)
    add_text(s, 0.55, 1.35, 5.9, 0.3, '❌ v0.10 公式', size=12, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    before = """Step 6: Total DL/day
  = (dl_per_shift × lines × shifts)
  + (debug_dl × debug_lines × shifts)
  + (functional_dl × shifts)
  = (10 × 1 × 2) + (3 × 0 × 2) + (2 × 2)
  = 24 (人)

Step 8: DL cost / week
  = DL_wage × SEA_wk × Total_DL × mult2 × mult1
  = $4.95 × 60 × 24 × 1.112 × 1.150
  = $9,116 / week

Step 9: DL cost / unit
  = (DL/wk + IDL_line_dep) / weekly_output
  = ($9,116 + $645) / 59,531
  = $0.1640 / unit

問題:所有人都用同一個 $4.95
  → 不準(Debug 應該算 $8.50)
  → 不公平(無塵室應該算 $6.50)"""
    add_text(s, 0.55, 1.8, 5.9, 4.8, before, size=10, color=TEXT, font='Consolas')

    add_rounded(s, 6.8, 1.3, 6.15, 5.4, RGBColor(0xDC, 0xFC, 0xE7), line=GREEN, line_width=1.5, radius=0.03)
    add_rounded(s, 6.8, 1.3, 6.15, 0.4, GREEN, radius=0.03)
    add_text(s, 6.95, 1.35, 5.85, 0.3, '✓ v0.11 公式', size=12, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    after = """Step 6 改:per shift cost(SUMPRODUCT)
  = Σ_role (role.hourly_rate × count_per_shift × SEA_day)
  = OP_GENERAL    $4.95 × 8  × 10 = $396.0
  + OP_CLEANROOM  $6.50 × 2  × 10 = $130.0  ← 改了!
  + OP_DEBUG      $8.50 × 3  × 10 = $255.0  ← 改了!
  + OP_FUNCT      $12.00 × 2 × 10 = $240.0  ← 改了!
  = $1,021.0 / shift

Step 8 改:DL cost / week
  = Σ_role per shift × lines × shifts × mult1 × mult2 × Wday
  = $1,021 × 1 × 2 × 1.15 × 1.112 × 6.67
  = $17,371 / week

Step 9: DL cost / unit
  = ($17,371 + $645) / 59,531
  = $0.3026 / unit  ← 比舊 $0.164 高 ↑ 84%
  (反映實際 wage 結構 · 比舊版精準)

優點:
  ✓ 反映實際 wage 結構
  ✓ 跨製程 role 各自獨立
  ✓ 月度漲薪可單一 role 調"""
    add_text(s, 6.95, 1.8, 5.85, 4.8, after, size=10, color=TEXT, font='Consolas')

    footer(s, n, total)


def slide_formula_equip(total, n):
    s = add_slide()
    header(s, 'Step 5 區 C Equipment 公式 · 類別 rep_acq 改寫', '取代 v0.10 line_qty × qty_per_line × acq_cost', GREEN, '📊 公式')

    add_rounded(s, 0.4, 1.3, 6.2, 5.4, RGBColor(0xFE, 0xE2, 0xE2), line=RED, line_width=1.5, radius=0.03)
    add_rounded(s, 0.4, 1.3, 6.2, 0.4, RED, radius=0.03)
    add_text(s, 0.55, 1.35, 5.9, 0.3, '❌ v0.10 公式', size=12, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    before = """Step 1: 抽 process 名下的設備(SMT_MAIN 7 件)
  DEK / FUJI NXT / SPI / Reflow / AOI /
  ESD Tray / Stencil

Step 2: 總 acquisition cost
  = Σ (line_qty × qty_per_line × acq_cost)
  = DEK    1 × 1 × $84,000 = $84,000
  + FUJI   1 × 1 × $1,013k = $1,013,384
  + SPI    1 × 1 × $68,852 = $68,852
  + Reflow 1 × 1 × $158k   = $158,000
  + AOI    1 × 1 × $92,000 = $92,000
  + ESD    1 × 752 × $18.5 = $13,912
  + Stencil 1 × 6 × $550   = $3,300
  = $1,433,448

Step 3: Annual Depreciation
  = Σ (acq / useful_life)
  ≈ $238,883 / yr

Step 4: Annual MRO
  = $1,433k × 5% = $71,650 / yr"""
    add_text(s, 0.55, 1.8, 5.9, 4.8, before, size=9.5, color=TEXT, font='Consolas')

    add_rounded(s, 6.8, 1.3, 6.15, 5.4, RGBColor(0xDC, 0xFC, 0xE7), line=GREEN, line_width=1.5, radius=0.03)
    add_rounded(s, 6.8, 1.3, 6.15, 0.4, GREEN, radius=0.03)
    add_text(s, 6.95, 1.35, 5.85, 0.3, '✓ v0.11 公式', size=12, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    after = """Step 1 改:抽 process 名下的類別 binding
  (SMT_MAIN 2 類):
  SMT_LINE_FULL × 1.0
  SMT_AOI × 0.5

Step 2: 總 acquisition cost(類別代表單價)
  = Σ (qty × rep_acq_cost_usd)
  = SMT_LINE_FULL  1.0 × $1,200,000 = $1,200,000
  + SMT_AOI        0.5 × $100,000   = $50,000
  = $1,250,000  (代表機 · 試產才細到具體機)

Step 3: Annual Depreciation
  = Σ (qty × rep_acq / rep_useful_life)
  = 1.0 × $1,200,000 / 6 + 0.5 × $100,000 / 6
  ≈ $208,333 / yr

Step 4: Annual MRO
  = Σ (qty × rep_acq × rep_mro_pct)
  = 1.0 × $1,200,000 × 5% + 0.5 × $100,000 × 5%
  = $62,500 / yr

優點:
  ✓ 報價階段算得出 cost
  ✓ 不糾結具體機種(DEK vs FUJI)
  ✓ 試產可選性細化(override)
  ✓ 跨案 utilization 看類別總和"""
    add_text(s, 6.95, 1.8, 5.85, 4.8, after, size=9.5, color=TEXT, font='Consolas')

    footer(s, n, total)


# ────────────────────────────────────────────────────────────────────────
# Section 6: Checklist / Close(2 slides)
# ────────────────────────────────────────────────────────────────────────

def slide_checklist(total, n):
    s = add_slide()
    header(s, '👥 各角色 Checklist', 'EPM / DPM / admin / BPM 各自查看點', ORANGE, '🎯 速查')

    roles = [
        ('🛠️ EPM(廠 EPM)', SKY, [
            '☐ 自家 BG baseline 月度更新',
            '☐ 廠類別代表單價 month-end check',
            '☐ DL 角色費率漲薪通報',
            '☐ 案級類別 binding 配置',
            '☐ Compute regression test PASS',
        ]),
        ('🔒 DPM(總監)', PURPLE, [
            '☐ baseline 切版二簽(漲幅 > 5%)',
            '☐ 案 Cleansheet lock 前審查',
            '☐ 24 cells matrix 看 margin',
            '☐ cross-BG 不該看的別看',
            '☐ 跨案 utilization 異常排查',
        ]),
        ('⭐ admin / IT', RED, [
            '☐ 新 BG 上線 catalog copy',
            '☐ user.bg_code 從 LDAP 同步',
            '☐ grant 表月度 review',
            '☐ super-admin 跨 BG 操作 audit',
            '☐ schema migration v0.11',
        ]),
        ('💼 BPM(業務 PM)', GREEN, [
            '☐ 案開立綁正確 BG/BU',
            '☐ 看 quote(不看 true cost)',
            '☐ 客戶議價時 reprice 觸發',
            '☐ 跨 BG 比價 (admin 介面)',
            '☐ 報價交期 vs 月度 baseline 切版時機',
        ]),
    ]
    for i, (name, color, items) in enumerate(roles):
        col = i % 2
        row = i // 2
        x = 0.4 + col * 6.27
        y = 1.45 + row * 2.7
        add_rounded(s, x, y, 6.2, 2.55, BG_SOFT, line=color, line_width=1.5, radius=0.03)
        add_rounded(s, x, y, 6.2, 0.42, color, radius=0.03)
        add_text(s, x + 0.15, y + 0.05, 5.8, 0.32, name, size=12, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
        for j, line in enumerate(items):
            add_text(s, x + 0.25, y + 0.5 + j * 0.4, 5.8, 0.35, line, size=10, color=TEXT, anchor=MSO_ANCHOR.TOP)

    footer(s, n, total)


def slide_close(total, n):
    s = add_slide()
    add_rect(s, 0, 0, 13.333, 7.5, INK)
    add_rect(s, 0, 2.8, 13.333, 0.06, OCEAN)
    add_text(s, 0.7, 1.2, 12, 0.6, '完。', size=44, bold=True, color=WHITE)
    add_text(s, 0.7, 1.9, 12, 0.4, 'Q&A · 等待 schema migration 拍板', size=18, color=OCEAN)

    add_rounded(s, 0.7, 3.4, 12.0, 3.5, RGBColor(0x14, 0x2F, 0x4E), radius=0.04)
    add_text(s, 1.0, 3.6, 11.5, 0.4, '🔗 v3 後續行動', size=16, bold=True, color=OCEAN)
    add_text(s, 1.0, 4.1, 11.5, 2.7, """• 互動 demo:Cortex_互動Demo_v0.11.html(§Cleansheet Step 1-7 全套 + BG 隔離)
• SD 更新:cleansheet-mva-sd v0.4(類別 + DL 角色 + BG/BU 隔離 + 權限 grant)
• schema migration:13 個表(7 個新 + 6 個 ALTER)+ 1 個權限 grant 表
• 與既有「資料權限管理」整合:對齊 user 主檔欄位 + 包含/排除政策模式
• 兩案 BG 綁定:SteelSeries → OPTO 光電 · WHOOP → CONSUMER 消費電子
• Phase 1 開工建議:先做 BG 隔離 schema → catalog/類別單價 → DL 角色 → 設定 UI → 案級遷移""",
        size=12, color=RGBColor(0xCB, 0xD5, 0xE1))

    add_text(s, 0.7, 7.05, 12, 0.3, 'Cortex MVA 操作手冊 v3 · BG/BU 權限隔離版 · 2026-06 · FOXLINK GPT',
             size=10, color=MUTED, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════
# RUN
# ════════════════════════════════════════════════════════════════════════
total = 2 + 5 + 5 + 3 + 5 + 2 + 2   # cover/overview + 5 sections + close
print(f'Building {total} slides...')

slide_cover()
slide_overview(total)

page = 3
slide_bg_intro(total, page); page += 1
slide_bg_isolation(total, page); page += 1
slide_bg_permission(total, page); page += 1
slide_bg_permission_examples(total, page); page += 1
slide_bg_case_binding(total, page); page += 1

slide_cat_why(total, page); page += 1
slide_cat_enum(total, page); page += 1
slide_cat_copy(total, page); page += 1
slide_cat_case_binding(total, page); page += 1
slide_cat_schema(total, page); page += 1

slide_dl_why(total, page); page += 1
slide_dl_case_config(total, page); page += 1
slide_dl_schema(total, page); page += 1

slide_settings_intro(total, page); page += 1
slide_settings_baseline(total, page); page += 1
slide_settings_equip_price(total, page); page += 1
slide_settings_dl_role(total, page); page += 1
slide_settings_permission(total, page); page += 1

slide_formula_dl(total, page); page += 1
slide_formula_equip(total, page); page += 1

slide_checklist(total, page); page += 1
slide_close(total, page); page += 1

OUT = 'd:/vibe_coding/foxlink_gpt/docs/Cortex_MVA操作流程說明手冊_v3_BG權限版.pptx'
prs.save(OUT)
print(f'[OK] Saved: {OUT}')
print(f'     Total slides: {len(prs.slides)}')
