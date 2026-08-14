# -*- coding: utf-8 -*-
"""
Cortex MVA 操作詳細手冊 v2 — 對齊 Cortex_互動Demo_v0.10.html §Cleansheet
6 大 Step × 多子步驟 · 每 step 含【目的】【UI】【操作】【結果】【schema】
~50 slides · 16:9 wide screen
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Palette
INK    = RGBColor(0x0B, 0x1F, 0x3A)
TEXT   = RGBColor(0x37, 0x47, 0x55)
MUTED  = RGBColor(0x6B, 0x72, 0x80)
OCEAN  = RGBColor(0x02, 0xC3, 0x9A)
NAVY   = RGBColor(0x1C, 0x72, 0x93)
LINE   = RGBColor(0xE5, 0xE7, 0xEB)
BG_SOFT= RGBColor(0xFA, 0xFB, 0xFC)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GOLD   = RGBColor(0xCA, 0x8A, 0x04)
RED    = RGBColor(0xDC, 0x26, 0x26)
GREEN  = RGBColor(0x16, 0xA3, 0x4A)
PURPLE = RGBColor(0x7C, 0x3A, 0xED)
CYAN   = RGBColor(0x08, 0x91, 0xB2)
SKY    = RGBColor(0x0E, 0xA5, 0xE9)

STEP_COLORS = {
    1: SKY,     # 製程輸入
    2: PURPLE,  # IDL matrix
    3: CYAN,    # 設備
    4: GOLD,    # 耗材
    5: RED,     # Compute
    6: NAVY,    # 結果矩陣
}


def add_slide(): return prs.slides.add_slide(prs.slide_layouts[6])


def add_text(s, x, y, w, h, text, size=12, bold=False, color=None,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font='Microsoft JhengHei'):
    tx = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tx.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    if color:
        r.font.color.rgb = color
    return tx


def add_rect(s, x, y, w, h, fill, line=None, line_width=0.5):
    rect = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    rect.fill.solid()
    rect.fill.fore_color.rgb = fill
    if line is None:
        rect.line.fill.background()
    else:
        rect.line.color.rgb = line
        rect.line.width = Pt(line_width)
    rect.shadow.inherit = False
    return rect


def add_rounded(s, x, y, w, h, fill, line=None, line_width=0.5, radius=0.05):
    rect = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    rect.adjustments[0] = radius
    rect.fill.solid()
    rect.fill.fore_color.rgb = fill
    if line is None:
        rect.line.fill.background()
    else:
        rect.line.color.rgb = line
        rect.line.width = Pt(line_width)
    rect.shadow.inherit = False
    return rect


def add_pill(s, x, y, text, fg, bg=None, size=9, w=None):
    pill_bg = bg or fg
    fg_text = WHITE if bg is None else fg
    pw = w if w is not None else (len(text) * 0.085 + 0.2)
    add_rounded(s, x, y, pw, 0.24, pill_bg, radius=0.5)
    add_text(s, x, y + 0.03, pw, 0.2, text, size=size, bold=True, color=fg_text,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return pw


def header(s, title, subtitle=None, step_n=None, step_label=None):
    color = STEP_COLORS.get(step_n, NAVY)
    add_rect(s, 0, 0, 13.333, 0.16, color)
    if step_n is not None:
        add_rounded(s, 0.4, 0.3, 1.5, 0.45, color, radius=0.3)
        add_text(s, 0.4, 0.32, 1.5, 0.4, f'STEP {step_n}', size=12, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, 2.05, 0.25, 11, 0.5, title, size=22, bold=True, color=INK)
        if subtitle:
            add_text(s, 2.05, 0.78, 11, 0.32, subtitle, size=11, color=MUTED)
    else:
        add_text(s, 0.4, 0.25, 11, 0.5, title, size=24, bold=True, color=INK)
        if subtitle:
            add_text(s, 0.4, 0.8, 11, 0.32, subtitle, size=12, color=MUTED)


def footer(s, n, total):
    add_rect(s, 0.4, 7.08, 12.55, 0.02, LINE)
    add_text(s, 0.4, 7.12, 9, 0.22, 'Cortex MVA 操作手冊 v2 · 對齊 Cortex_互動Demo_v0.10.html §Cleansheet',
             size=8.5, color=MUTED)
    add_text(s, 12.0, 7.12, 1, 0.22, f'{n} / {total}', size=8.5, color=MUTED, align=PP_ALIGN.RIGHT)


def info_card(s, x, y, w, h, ico, title, body, color=NAVY):
    """4-corner cards used heavily for 【目的】【操作】etc."""
    add_rounded(s, x, y, w, h, RGBColor(0xFA, 0xFB, 0xFC), line=color, line_width=1, radius=0.03)
    add_rounded(s, x, y, w, 0.32, color, radius=0.03)
    add_text(s, x + 0.1, y + 0.04, w - 0.2, 0.26, f'{ico}  {title}', size=11, bold=True, color=WHITE,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x + 0.15, y + 0.4, w - 0.3, h - 0.5, body, size=10, color=TEXT, anchor=MSO_ANCHOR.TOP)


def mockup_box(s, x, y, w, h, title, lines, color=NAVY):
    """UI mockup style box · 給操作示意用"""
    add_rounded(s, x, y, w, h, WHITE, line=color, line_width=1.5, radius=0.03)
    add_rounded(s, x, y, w, 0.3, color, radius=0.03)
    # 3 dots(modal corner)
    for i, c in enumerate([RGBColor(0xEF,0x44,0x44), RGBColor(0xF5,0xA5,0x24), RGBColor(0x10,0xB9,0x81)]):
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x+0.08+i*0.14), Inches(y+0.09), Inches(0.12), Inches(0.12))
        dot.fill.solid(); dot.fill.fore_color.rgb = c
        dot.line.fill.background(); dot.shadow.inherit = False
    add_text(s, x + 0.6, y + 0.04, w - 0.7, 0.24, title, size=10, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x + 0.15, y + 0.4, w - 0.3, h - 0.5, lines, size=9, color=TEXT, font='Consolas', anchor=MSO_ANCHOR.TOP)


# ════════════════════════════════════════════════════════════════════════
# COVER + TOC + OVERVIEW
# ════════════════════════════════════════════════════════════════════════

def slide_cover():
    s = add_slide()
    add_rect(s, 0, 0, 13.333, 7.5, INK)
    add_rect(s, 0, 2.7, 13.333, 0.06, OCEAN)
    add_text(s, 0.7, 1.2, 12, 0.4, 'FOXLINK 正崴 · Cortex BOM × MVA', size=12, color=OCEAN)
    add_text(s, 0.7, 1.6, 12, 0.95, 'MVA 操作詳細手冊', size=48, bold=True, color=WHITE)
    add_text(s, 0.7, 3.0, 12, 0.5, 'Step-by-step 帶 EPM / DPM 完成 MVA 全流程', size=20, color=WHITE)

    add_rounded(s, 0.7, 3.85, 12, 1.8, RGBColor(0x14, 0x2F, 0x4E), radius=0.04)
    add_text(s, 1.0, 4.0, 11.5, 0.4, '本手冊涵蓋的 6 大 Step', size=14, bold=True, color=OCEAN)
    info_steps = [
        ('1', 'Step 1', '案級製程輸入', '9 製程 × 22 欄位', SKY),
        ('2', 'Step 2', 'IDL multiplier matrix', '17 角色 × 9 製程', PURPLE),
        ('3', 'Step 3', '案級設備 binding', '廠 master 58 → 案級 28', CYAN),
        ('4', 'Step 4', '案級耗材 binding', '17 件耗材', GOLD),
        ('5', 'Step 5', 'Compute Trace 5 區公式', 'DL+IDL+Equip+Fac+Others+Total', RED),
        ('6', 'Step 6', '9×9 結果矩陣', 'MVA + SG&A + Profit → TC', NAVY),
    ]
    for i, (n, st, t, sub, c) in enumerate(info_steps):
        x = 1.0 + (i % 3) * 4.0
        y = 4.5 + (i // 3) * 0.5
        add_rounded(s, x, y, 0.32, 0.32, c, radius=0.5)
        add_text(s, x, y, 0.32, 0.32, n, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + 0.4, y, 3.5, 0.32, f'{st} · {t}', size=10, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)

    add_text(s, 0.7, 6.0, 12, 0.4, '對齊 SteelSeries Rival 3+ Wired Mouse 真實案例 · 9 製程 · 24 cells', size=12, color=OCEAN)
    add_text(s, 0.7, 6.45, 12, 0.4, 'v2 詳細版 · 2026-06 · 對應 Cortex_互動Demo_v0.10.html §Cleansheet', size=11, color=MUTED)


def slide_toc(total_pages):
    s = add_slide()
    header(s, '操作流程全景', '從進 §Cleansheet section → 6 Step 完整跑完 · 預估 30 分鐘')

    # 上方 6 step pipeline (chevron)
    add_text(s, 0.4, 1.3, 12, 0.3, '🎬 操作 Pipeline (進 SteelSeries 案 → §Cleansheet section)',
             size=12, bold=True, color=INK)
    chev_w = 1.95
    chev_x_start = 0.4
    chev_y = 1.85
    steps = [
        (1, '製程輸入', '~5 min'),
        (2, 'IDL Matrix', '~3 min'),
        (3, '設備 binding', '~4 min'),
        (4, '耗材 binding', '~2 min'),
        (5, 'Compute 5 區', '~10 min'),
        (6, '結果矩陣', '~5 min'),
    ]
    for i, (n, title, dur) in enumerate(steps):
        color = STEP_COLORS[n]
        x = chev_x_start + i * (chev_w + 0.1)
        sh = s.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(x), Inches(chev_y), Inches(chev_w), Inches(0.9))
        sh.fill.solid(); sh.fill.fore_color.rgb = color
        sh.line.fill.background(); sh.shadow.inherit = False
        add_text(s, x, chev_y + 0.05, chev_w - 0.25, 0.3, f'Step {n}', size=12, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x, chev_y + 0.35, chev_w - 0.25, 0.3, title, size=11, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x, chev_y + 0.6, chev_w - 0.25, 0.25, dur, size=9, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # 下方:每 Step slide count + 重點
    add_text(s, 0.4, 3.2, 12, 0.3, '📚 本手冊每 Step 涵蓋', size=12, bold=True, color=INK)
    cards = [
        (1, '製程輸入', 5, ['§A Excel 截圖對應', '可即時編輯', '改 TAKT → DL 自動動']),
        (2, 'IDL Matrix', 4, ['17 角色 × 9 製程', '153 格 multiplier', '初始抄 Excel J64-Z83']),
        (3, '設備 binding', 5, ['廠 master 58 件', '案內 28 件', 'line_qty × qty_per_line']),
        (4, '耗材 binding', 3, ['17 件配置', 'annual_usage', '對應 Consumables sheet']),
        (5, 'Compute', 10, ['5 區公式逐步驟', 'A DL 9 / B IDL 4 / C Equip 6', 'D Fac 6 / E Others 4 / F Total 4']),
        (6, '結果矩陣', 4, ['9×9 cost component', 'SMT/Assy roll-up', 'Total Cost waterfall']),
    ]
    for i, (n, t, cnt, bullets) in enumerate(cards):
        col = i % 3
        row = i // 3
        x = 0.4 + col * 4.25
        y = 3.6 + row * 1.7
        color = STEP_COLORS[n]
        add_rounded(s, x, y, 4.05, 1.55, BG_SOFT, line=color, line_width=1.5, radius=0.03)
        add_rounded(s, x, y, 4.05, 0.4, color, radius=0.03)
        add_text(s, x + 0.15, y + 0.05, 3, 0.3, f'Step {n} · {t}', size=11, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + 3, y + 0.05, 1, 0.3, f'{cnt} slides', size=9, color=WHITE, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT)
        for j, b in enumerate(bullets):
            add_text(s, x + 0.2, y + 0.5 + j * 0.32, 3.7, 0.3, f'• {b}', size=9.5, color=TEXT)

    footer(s, 2, total_pages)


def slide_prereq(total_pages):
    s = add_slide()
    header(s, '開始前 · 必要前提', 'Step 1 之前 · Phase A/B/C 必須完成(本手冊不重複講 · 看 v1 手冊)')

    # 3 大前提 cards
    prereqs = [
        ('🏗️', 'Phase A 廠 master 已建', NAVY,
         '・ CN 廠 baseline 已 import\n  (Rival 3+ Cleansheet xlsx → SupplierBaseInput / IDL 17 / 設備 58 / 耗材 17)\n・ 製程目錄 10 項 enum 已建\n・ MOUSE_STD 模板已建\n\n→ 確認:Admin → Factory Baselines → CN-2024Q4 status=active'),
        ('📂', 'Phase C 案 case_factory 已開', NAVY,
         '・ BPM 已開 SteelSeries Rival 3+ project\n・ case_factory 已建(CN baseline 鎖住)\n・ process_template = MOUSE_STD\n・ qty_scenarios 配 LOW(100K) + HIGH(418K)\n・ pkg_versions 配 RETAIL(商包)+ BULK(工包)\n\n→ 確認:project 開啟 · §customer 已填'),
        ('📥', 'BOM xlsx 已 import', NAVY,
         '・ RD 上傳 BOM xlsx\n・ AI 解析 22 子料 + ERP 比對完\n・ §variants 已標(EE BOM shared / ME BOM per_variant)\n・ §BOM 採購策略已選定\n\n→ 確認:§BOM section status = done'),
    ]
    for i, (ico, t, c, body) in enumerate(prereqs):
        x = 0.4 + i * 4.25
        y = 1.5
        add_rounded(s, x, y, 4.05, 5.3, BG_SOFT, line=c, line_width=1.5, radius=0.03)
        add_rounded(s, x, y, 4.05, 0.55, c, radius=0.03)
        add_text(s, x + 0.15, y + 0.1, 3.8, 0.4, f'{ico}  {t}', size=12, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + 0.2, y + 0.7, 3.8, 4.5, body, size=10, color=TEXT)

    add_rounded(s, 0.4, 6.55, 12.55, 0.4, RGBColor(0xFE, 0xF3, 0xC7), line=GOLD, radius=0.04)
    add_text(s, 0.55, 6.6, 12.3, 0.3, '⚠️  若以上未完成 · 請先回 §workflow_checklist Phase A-C · 完成後再回 §Cleansheet 進 Step 1',
             size=10, bold=True, color=RGBColor(0x92, 0x40, 0x0E), anchor=MSO_ANCHOR.MIDDLE)

    footer(s, 3, total_pages)


# ════════════════════════════════════════════════════════════════════════
# STEP 1 — 案級製程輸入(5 slides)
# ════════════════════════════════════════════════════════════════════════

def slide_step1_intro(total_pages, page_n):
    s = add_slide()
    header(s, '案級製程輸入 · 概念', '9 製程 × 22 欄位 · 對應 Rival 3+ Excel Cleansheet sheet §A 上半', 1, '案級製程輸入')

    # 兩欄:左 概念 + 右 對應 Excel 截圖示意
    add_rounded(s, 0.4, 1.3, 6.2, 5.4, BG_SOFT, line=SKY, line_width=1.5, radius=0.03)
    add_rounded(s, 0.4, 1.3, 6.2, 0.4, SKY, radius=0.03)
    add_text(s, 0.55, 1.35, 5.9, 0.3, '💡 為什麼要這步?', size=12, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, 0.6, 1.85, 5.9, 4.8, """製程輸入是 MVA 計算的 source of truth — 系統靠這些
變數跑 DL/IDL/Equipment/Facility 4 區公式。

每個製程要填 22 個欄位,分 6 群組:

⏱  Staffed Hours (4 欄)
    working hr/day · days/wk · hrs/wk(計算)· shifts

⚡  Throughput (4 欄)
    TAKT · UPH(計算)· Yield · Efficiency

📊  Volume (4 欄,3 個計算)
    Annual demand · avg/wk · max/wk · 安全係數

👷  DL Config (4 欄)
    DL/shift · debug DL · functional DL · lines

👔  IDL Line-dep (4 欄)
    Line Leader · Tech · Supervisor · TAKT 來源

🕒  SEA Hours (4 欄,2 derived)
    SEA day/wk · 註記 · baseline 引用""", size=10, color=TEXT)

    # 右側:Excel 截圖式 mockup
    add_rounded(s, 6.8, 1.3, 6.15, 5.4, WHITE, line=NAVY, line_width=1.5, radius=0.03)
    add_rounded(s, 6.8, 1.3, 6.15, 0.4, NAVY, radius=0.03)
    add_text(s, 6.95, 1.35, 5.85, 0.3, '📊 對應 Excel Cleansheet sheet §A', size=12, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    excel_lines = """A. Direct Labor Cost: PCBA + Box Build + Debug
                       SMT   Wave  Router Laser BB-Assy
─── Staffed Hours ──────────────────────────────
Working hr/day         23    20    20     20    20
Working days/week      6.67  6.67  6.67   6.67  6.00
Working hr/week (calc) 153   133   133    133   120
Shifts per day         2     2     2      2     1

─── Throughput ─────────────────────────────────
TAKT time (s)          8.64  12.00 24.00  4.00  19.00
UPH (calc)             388   279   140    838   176
Yield %                98%   98%   98%    98%   98%
Efficiency %           95%   95%   95%    95%   95%

─── Volume ─────────────────────────────────────
Annual Demand          418,000 (per case_factory)
Avg weekly             8,360
Max weekly demand      10,032   ← Avg × 1.2"""
    add_text(s, 6.95, 1.8, 5.85, 4.8, excel_lines, size=8.5, color=TEXT, font='Consolas')

    footer(s, page_n, total_pages)


def slide_step1_ui_entry(total_pages, page_n):
    s = add_slide()
    header(s, '案級製程輸入 · 進入 UI', '在 §Cleansheet section 內 · 找 [Step 1 製程輸入] step pill', 1)

    # 操作示意 mockup
    add_text(s, 0.4, 1.3, 12, 0.3, '🎯 點什麼:Cleansheet section → Step pill navigator → 「📋 Step 1 製程輸入」',
             size=12, bold=True, color=INK)

    mockup = """┌────── §Cleansheet (MVA 計算 · v0.7 qty-aware) ──────────────────────────┐
│  [CN][VN][TW]  [📥 Re-import] [🔄 Compute] [🔒 Lock]                     │
│                                                                          │
│  CN Factory: CN-2024Q4 · Effective 2024-10-11 · DL wage $4.95/hr · ...   │
│                                                                          │
│  📊 Qty Scenario: [低批量 100K][高批量 418K ●]                            │
│                                                                          │
│  [MVA Total $1.8558]  [SGA+Profit $1.6855]  [SMT MVA]  [Assy MVA]        │
│                                                                          │
│  🎬 MVA 操作每一步展開 · 點 step 看該階段 EPM 實際輸入 / 系統實際算什麼  │
│                                                                          │
│  ╔════════════════════╗ ┌──────────┐ ┌──────────┐ ┌──────────┐ ┌──────┐ │
│  ║ 📋 Step 1 製程輸入 ║ │ 👥 Step 2│ │ 🏭 Step 3│ │ 📦 Step 4│ │⚙Step 5│ │
│  ║ ← 點這個            ║ │ IDL matrx│ │  設備    │ │  耗材    │ │Compute│ │
│  ╚════════════════════╝ └──────────┘ └──────────┘ └──────────┘ └──────┘ │
│                                                                          │
│  ━━━━━━━━━━ Panel 1 內容展開 ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━      │
└──────────────────────────────────────────────────────────────────────────┘"""
    add_rounded(s, 0.4, 1.7, 12.55, 4.5, WHITE, line=SKY, line_width=1.5, radius=0.02)
    add_text(s, 0.55, 1.85, 12.3, 4.3, mockup, size=10, color=TEXT, font='Consolas')

    # 提示
    add_rounded(s, 0.4, 6.35, 12.55, 0.6, RGBColor(0xDB, 0xEA, 0xFE), line=SKY, radius=0.04)
    add_text(s, 0.55, 6.4, 12.3, 0.5,
             '✏️ Step pill 為 v0.10 新加 · 點任一 pill 切換顯示對應 panel · 預設停在 Step 1\n💡 若 case 非 CN · Step pill 也會展示 · 但內容是唯讀(等真實 xlsx 後可編輯)',
             size=10, color=INK, anchor=MSO_ANCHOR.MIDDLE)

    footer(s, page_n, total_pages)


def slide_step1_groups(total_pages, page_n):
    s = add_slide()
    header(s, '案級製程輸入 · 6 群組 22 欄位', 'Step 1 panel 內主體 · 每群組 4 欄 · derived 欄(紫色)系統自算', 1)

    # 6 grid 卡(對應 6 groups)
    groups = [
        ('⏱', 'Staffed Hours', SKY, [
            ('Working hours per day', '23', False),
            ('Working days per week', '6.67', False),
            ('Working hours per week', '153.41', True),
            ('Shifts per day', '2', False),
        ]),
        ('⚡', 'Throughput', GREEN, [
            ('TAKT time (s)', '8.64', False),
            ('UPH (units/hour)', '388.06', True),
            ('Yield %', '0.98', False),
            ('Efficiency rate', '0.95', False),
        ]),
        ('📊', 'Volume', GOLD, [
            ('Annual demand', '418,000', True),
            ('Avg weekly demand', '8,360', True),
            ('Max weekly demand', '10,032', True),
            ('(Max = Avg × 1.2)', '安全係數', True),
        ]),
        ('👷', 'DL Config', RED, [
            ('DL / shift', '10', False),
            ('Debug DL / shift', '3', False),
            ('Functional DL / shift', '2', False),
            ('Lines / debug lines', '1 / 0', True),
        ]),
        ('👔', 'IDL Line-dep', PURPLE, [
            ('Line Leader / shift', '0.5', False),
            ('Technician / shift', '0.5', False),
            ('Supervisor / day', '0.25', False),
            ('TAKT 來源', 'SMT Flow!I3', True),
        ]),
        ('🕒', 'SEA Hours', RGBColor(0x6B,0x72,0x80), [
            ('SEA hours / day', '10', False),
            ('SEA hours / week', '60', False),
            ('(SEA 常數)', '註記', True),
            ('baseline 引用', '✓ active', True),
        ]),
    ]
    for i, (ico, name, color, fields) in enumerate(groups):
        col = i % 3
        row = i // 3
        x = 0.4 + col * 4.2
        y = 1.45 + row * 2.6
        add_rounded(s, x, y, 4.1, 2.4, WHITE, line=color, line_width=1.2, radius=0.03)
        add_rounded(s, x, y, 4.1, 0.35, color, radius=0.03)
        add_text(s, x + 0.15, y + 0.04, 3.7, 0.27, f'{ico}  {name}', size=11, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
        for j, (k, v, derived) in enumerate(fields):
            ry = y + 0.45 + j * 0.46
            add_rect(s, x + 0.15, ry, 3.8, 0.4, BG_SOFT, line=LINE)
            add_text(s, x + 0.25, ry + 0.04, 2.3, 0.3, k, size=9, color=TEXT, anchor=MSO_ANCHOR.MIDDLE)
            val_color = PURPLE if derived else INK
            add_text(s, x + 2.5, ry + 0.04, 1.4, 0.3, v, size=10, bold=True, color=val_color,
                     align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE, font='Consolas')
            if derived:
                add_text(s, x + 2.5, ry + 0.22, 1.4, 0.16, '(derived)', size=7, color=PURPLE,
                         align=PP_ALIGN.RIGHT)

    footer(s, page_n, total_pages)


def slide_step1_interactive(total_pages, page_n):
    s = add_slide()
    header(s, '案級製程輸入 · ✏️ 即時編輯 demo', '改 TAKT / Yield → 自動重算 DL_CPU → 影響 §Step 6 矩陣 · v0.10 新加', 1)

    # 左:操作步驟 · 右:即時看 effect
    add_rounded(s, 0.4, 1.3, 6.2, 5.4, BG_SOFT, line=GOLD, line_width=1.5, radius=0.03)
    add_rounded(s, 0.4, 1.3, 6.2, 0.4, GOLD, radius=0.03)
    add_text(s, 0.55, 1.35, 5.9, 0.3, '🎬 操作步驟(以 SMT_MAIN 為例)', size=12, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    steps_md = """1️⃣  點製程 tab「SMT 主板貼裝」(預設已選)

2️⃣  在「⚡ Throughput」群組找「TAKT time (s)」
    當前值 8.64 → 點 input 改成 10.00

3️⃣  系統立即:
    • UPH 自動重算 388 → 335 (derived 紫色)
    • 製程 tab「SMT」右上角出現橘色小圓點
    • 頂部出現「✏ 已改」徽章

4️⃣  下方「🎯 即時重算預覽」橫條:
    baseline DL/unit: $0.148
    → 新算:        $0.171 ↑ 15.5%

5️⃣  切到 Step 6 結果矩陣 ↘
    DL_CPU 列 SMT 欄位 自動同步 $0.148 → $0.171
    MVA Total 自動更新

6️⃣  不滿意 → 點「⟲ Reset 本製程」或
            「⟲ Reset ALL」清掉所有 override"""
    add_text(s, 0.6, 1.85, 5.9, 4.8, steps_md, size=10, color=TEXT)

    add_rounded(s, 6.8, 1.3, 6.15, 5.4, WHITE, line=SKY, line_width=1.5, radius=0.03)
    add_rounded(s, 6.8, 1.3, 6.15, 0.4, SKY, radius=0.03)
    add_text(s, 6.95, 1.35, 5.85, 0.3, '🎯 即時重算預覽 UI mockup', size=12, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)

    preview = """┌─ Throughput ─────────────────────────────┐
│ TAKT time (s)        [10.00  ▾ ] s       │  ← 改了 · 底色變黃
│ UPH (units/hour)     335.31    (derived) │  ← 紫色 · 自動算
│ Yield %              0.98                │
│ Efficiency rate      0.95                │
└──────────────────────────────────────────┘

┌─ 🎯 即時重算預覽 ────────────────────────┐
│ baseline DL/unit: $0.1480                │
│ → 新算: $0.1710  ↑ 15.5%                 │
│                          [⟲ Reset 本製程]│
└──────────────────────────────────────────┘

╔══ Step 6 · 9×9 結果矩陣 ═════════════════╗
║                  SMT      Wave   ...     ║
║ DL_CPU         $0.1710  $0.349  ...      ║  ← 自動同步
║              (改前 0.148)                ║
║ IDL_CPU        $0.044  ...               ║
║ ...                                      ║
║ MVA Σ          $1.879   ↑ vs $1.856      ║
╚══════════════════════════════════════════╝

✅ schema 邏輯:
   UPDATE bom_cs_case_process
   SET takt_seconds = 10.00
   WHERE case_factory_id = ? AND process_code = 'SMT_MAIN'
   → cs_run.dirty = 1 (等下次 Compute 重算)"""
    add_text(s, 6.95, 1.8, 5.85, 4.8, preview, size=8, color=TEXT, font='Consolas')

    footer(s, page_n, total_pages)


def slide_step1_schema(total_pages, page_n):
    s = add_slide()
    header(s, '案級製程輸入 · 寫入 DB', '對應 bom_cs_case_process schema · UPDATE 觸發 cs_run.dirty=1', 1)

    add_rounded(s, 0.4, 1.3, 12.55, 5.4, BG_SOFT, line=SKY, line_width=1.5, radius=0.03)
    add_rounded(s, 0.4, 1.3, 12.55, 0.4, SKY, radius=0.03)
    add_text(s, 0.55, 1.35, 12.3, 0.3, '🗄️ bom_cs_case_process — 案級製程設定主表', size=12, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)

    schema_md = """-- Oracle 23 AI · NUMBER GENERATED ALWAYS AS IDENTITY
CREATE TABLE bom_cs_case_process (
  case_factory_id          NUMBER REFERENCES bom_cs_case_factory(case_factory_id) ON DELETE CASCADE,
  process_code             VARCHAR2(40) REFERENCES bom_process_catalog(process_code),
  step_order               NUMBER NOT NULL,

  -- ⏱ Staffed Hours (4 欄)
  working_hours_per_day    NUMBER(5,2),   -- 23  (SMT 23hr · BB 20hr)
  working_days_per_week    NUMBER(5,2),   -- 6.67 (SMT 6.67d · BB 6.00d)
  shifts_per_day           NUMBER(2),     -- 2   (SMT 2 shifts · BB 1 shift)

  -- ⚡ Throughput (3 欄,UPH derived)
  takt_seconds             NUMBER(10,4),  -- 8.64
  yield_pct                NUMBER(5,4),   -- 0.98
  efficiency_pct           NUMBER(5,4),   -- 0.95

  -- 👷 DL Config (5 欄)
  dl_per_shift             NUMBER(5,2),   -- 10
  debug_dl_per_shift       NUMBER(5,2),   -- 3
  functional_dl_per_shift  NUMBER(5,2),   -- 2
  lines_installed          NUMBER(3) DEFAULT 1,
  debug_lines_installed    NUMBER(3) DEFAULT 0,

  -- 👔 IDL Line-dep (3 欄)
  line_leader_per_shift    NUMBER(5,2),   -- 0.5
  technician_per_shift     NUMBER(5,2),   -- 0.5
  supervisor_per_day       NUMBER(5,2),   -- 0.25

  -- 🕒 SEA Hours (2 欄)
  sea_hours_per_day        NUMBER(5,2) DEFAULT 10,
  sea_hours_per_week       NUMBER(5,2) DEFAULT 60,

  -- TAKT 來源(文字註記)
  takt_source              VARCHAR2(80),  -- 'SMT Process Flow!I3'
  notes                    CLOB,

  PRIMARY KEY (case_factory_id, process_code)
);

-- Annual Demand 是 case_factory level(不在製程表)
-- → bom_cs_case_factory.annual_demand · 影響所有製程的 demand 計算"""
    add_text(s, 0.6, 1.8, 12.1, 5.0, schema_md, size=9, color=TEXT, font='Consolas')

    footer(s, page_n, total_pages)


# ════════════════════════════════════════════════════════════════════════
# STEP 2 — IDL Matrix(4 slides)
# ════════════════════════════════════════════════════════════════════════

def slide_step2_intro(total_pages, page_n):
    s = add_slide()
    header(s, 'IDL × 製程 multiplier matrix · 概念', '17 角色 × 9 製程 = 153 個 multiplier · 案級配置', 2, 'IDL matrix')

    add_rounded(s, 0.4, 1.3, 6.2, 5.4, BG_SOFT, line=PURPLE, line_width=1.5, radius=0.03)
    add_rounded(s, 0.4, 1.3, 6.2, 0.4, PURPLE, radius=0.03)
    add_text(s, 0.55, 1.35, 5.9, 0.3, '💡 為什麼要 matrix?', size=12, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, 0.6, 1.85, 5.9, 4.8, """IDL(Indirect Labor)= 間接人力 · 不直接動產線。

問題:廠 OPS 經理 1 個人,工廠所有製程都吃他的薪水
       怎麼分攤?
解法:設「multiplier」表示該角色花多少比例
       服務該製程。

範例:
  OPS_MGR(廠 OPS 經理 $60,320/年)
    × SMT 0.08 (花 8% 時間管 SMT)
    × Wave 0.04
    × BB Assy 0.125
    × ...
  → SMT 製程吃他 $60,320 × 0.08 = $4,826/年

17 個角色:
  • CENTRAL service(廠級共用)
    OPS_MGR · HR · FIN · IT · ADMIN · SECURITY ·
    JANITOR · CANTEEN · DRIVER · NURSE
  • MFG IDL(製造相關)
    SEC_MGR_SMT · SEC_MGR_BB · ENGINEER ·
    PE_ENG · QE_ENG · PLANNER · STOREKEEPER

每角色每製程一個 multiplier · 共 153 格""", size=10, color=TEXT)

    # 右側:matrix 截圖示意
    add_rounded(s, 6.8, 1.3, 6.15, 5.4, WHITE, line=NAVY, line_width=1.5, radius=0.03)
    add_rounded(s, 6.8, 1.3, 6.15, 0.4, NAVY, radius=0.03)
    add_text(s, 6.95, 1.35, 5.85, 0.3, '📊 matrix UI 示意(顏色越深 multiplier 越大)', size=12, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    matrix = """IDL Role     年薪     SMT   Wave  Router Laser BB-Assy
─────────────────────────────────────────────────────
OPS_MGR     $60,320   0.080 0.040 0.040  0.040 0.125
SEC_MGR_SMT $42,500   ███   0.200 0.150  0.100 —
SEC_MGR_BB  $42,500   —     —     —      —     ███
ENGINEER    $24,440   0.300 0.100 0.080  0.040 0.200
PE_ENG      $24,440   0.250 0.080 0.050  0.030 0.150
QE_ENG      $24,440   0.100 0.030 0.020  0.010 0.080
PLANNER     $19,000   0.060 0.030 0.020  0.010 0.080
STOREKEEPER $14,500   0.050 0.030 0.020  0.010 0.100
HR_MGR      $28,000   0.040 0.020 0.020  0.020 0.060
FIN_MGR     $32,000   0.025 0.015 0.015  0.015 0.040
IT_MGR      $30,000   0.030 0.015 0.015  0.015 0.050
ADMIN       $13,500   0.030 0.015 0.015  0.015 0.045
SECURITY    $9,800    0.020 0.015 0.015  0.015 0.030
JANITOR     $8,800    0.040 0.030 0.030  0.030 0.050
CANTEEN     $9,800    0.035 0.020 0.020  0.020 0.050
DRIVER      $12,000   0.020 0.010 0.010  0.010 0.030
NURSE       $14,000   0.015 0.010 0.010  0.010 0.025
─────────────────────────────────────────────────────
Col Σ                  1.10  0.63  0.49   0.39  1.32"""
    add_text(s, 6.95, 1.8, 5.85, 4.8, matrix, size=7, color=TEXT, font='Consolas')

    footer(s, page_n, total_pages)


def slide_step2_ui(total_pages, page_n):
    s = add_slide()
    header(s, 'IDL matrix · UI 操作', '進 Step 2 panel · 看 153 格 matrix · 顏色漸層輔助判讀', 2)

    add_text(s, 0.4, 1.3, 12, 0.3, '🎯 點什麼:Step 1 → Step pill「👥 Step 2 IDL Matrix」',
             size=12, bold=True, color=INK)

    mockup = """╔════════════════════════════════════════════════════════════════════════════════╗
║ 👥 Step 2 · IDL × 製程 Multiplier Matrix                                       ║
║                            17 角色 × 9 製程 · 對應 Excel rows 64-83 · 案級配置 ║
║                                              17 角色 × 9 製程 = 153 格        ║
╠════════════════════════════════════════════════════════════════════════════════╣
║ IDL 角色          年薪 USD  SMT  Wave Router Laser BB-Assy BB-Test Mat Q-SMT Q-BB Row Σ║
║ ───────────────────────────────────────────────────────────────────────────────║
║ 廠 OPS 經理      $60,320   0.080 0.040 0.040 0.040 0.125  0.080 0.080 0.020 0.020 0.523║
║ SMT 段經理       $42,500   0.500 0.200 0.150 0.100 0      0     0     0.050 0     1.000║
║ ...              ...       ...                                              ...║
║                                                                                ║
║ Col Σ(per process)         1.10  0.63  0.49  0.39  1.32  ...                  ║
╚════════════════════════════════════════════════════════════════════════════════╝

📌 顏色漸層:multiplier 越大 → 紫色越深 · 0 → 灰色「—」
   → 一眼看出哪角色重押哪製程(如 SMT 段經理重押 SMT_MAIN)"""
    add_rounded(s, 0.4, 1.7, 12.55, 4.0, WHITE, line=PURPLE, line_width=1.5, radius=0.02)
    add_text(s, 0.55, 1.85, 12.3, 3.8, mockup, size=8, color=TEXT, font='Consolas')

    # 提示
    add_rounded(s, 0.4, 5.85, 12.55, 1.1, RGBColor(0xF3, 0xE8, 0xFF), line=PURPLE, radius=0.04)
    add_text(s, 0.55, 5.92, 12.3, 1, """💡 判讀技巧:
  • 每列 Row Σ = 該角色「總負載」· 不該 > 1.0 (代表他超過自己 100% 工時 → 異常)
  • 每欄 Col Σ = 該製程吃多少 IDL · 跨案比較 / SMT 重複用同樣 IDL → 利用率好
  • 0(灰色)= 該角色完全不碰該製程(例 SMT 段經理不碰 BB Test)""",
             size=10, color=TEXT, anchor=MSO_ANCHOR.TOP)

    footer(s, page_n, total_pages)


def slide_step2_init(total_pages, page_n):
    s = add_slide()
    header(s, 'IDL matrix · 初始化策略', '首次開案 · 抄 Excel J64-Z83 · 後續 EPM 在 UI 編', 2)

    add_text(s, 0.4, 1.3, 12, 0.3, '🎯 153 格 multiplier 不可能從 0 填 · 系統提供 2 種 seed 策略',
             size=12, bold=True, color=INK)

    # 兩個策略 card
    add_rounded(s, 0.4, 1.7, 6.2, 5.0, BG_SOFT, line=GREEN, line_width=1.5, radius=0.03)
    add_rounded(s, 0.4, 1.7, 6.2, 0.4, GREEN, radius=0.03)
    add_text(s, 0.55, 1.75, 5.9, 0.3, '✅ 策略 A:Excel import (首次)', size=12, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, 0.6, 2.2, 5.9, 4.5, """• 首次開案,廠 EPM 的歷史 Cleansheet xlsx
  J64-Z83 已有完整 multiplier matrix
  (該廠之前所有案累積的經驗值)

• 系統 import 時自動抄這 153 格進
  bom_cs_case_idl_alloc 表
  (case_factory_id, process_code, role_code, multiplier)

• 後續若同廠開類似案 → 沿用相同 seed

• 適用:
  ✓ 廠級首份 baseline 已 import
  ✓ 案產品類型相近(都做滑鼠 / 都做鍵盤)
  ✓ EPM 信任歷史值

• 後續 EPM 在 UI 改 · 不回 Excel""", size=10, color=TEXT)

    add_rounded(s, 6.8, 1.7, 6.15, 5.0, BG_SOFT, line=GOLD, line_width=1.5, radius=0.03)
    add_rounded(s, 6.8, 1.7, 6.15, 0.4, GOLD, radius=0.03)
    add_text(s, 6.95, 1.75, 5.85, 0.3, '⚠️ 策略 B:模板 default (新類型)', size=12, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, 7.0, 2.2, 5.85, 4.5, """• 廠沒做過這類產品 / 案產品差異大

• 系統用 process_template 預設值:
  MOUSE_STD 模板 → SMT 重押 0.5,BB 輕押
  KEYBOARD_STD → 倒過來
  WHOOP_WEARABLE → 模組多 + FATP 重押

• EPM 需做更多手工調整

• 適用:
  ✗ 廠首次做這類產品
  ✗ 沒歷史可參考
  ✓ 新廠 import 後第一個案

• 警告:multiplier 可能太離譜 →
  Phase E Compute 後跑 regression test
  vs 廠平均 MVA · > 30% 差異需 EPM 解釋""", size=10, color=TEXT)

    footer(s, page_n, total_pages)


def slide_step2_schema(total_pages, page_n):
    s = add_slide()
    header(s, 'IDL matrix · 寫入 DB', '案級 153 row · 廠級角色定義 + 案級 multiplier 兩層', 2)

    add_rounded(s, 0.4, 1.3, 12.55, 5.4, BG_SOFT, line=PURPLE, line_width=1.5, radius=0.03)
    add_rounded(s, 0.4, 1.3, 12.55, 0.4, PURPLE, radius=0.03)
    add_text(s, 0.55, 1.35, 12.3, 0.3, '🗄️ bom_factory_idl_role(廠級)+ bom_cs_case_idl_alloc(案級)', size=12, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)

    schema = """-- (1) 廠級 IDL 角色定義 · 17 列 / 廠 / baseline 版本
CREATE TABLE bom_factory_idl_role (
  baseline_id        NUMBER REFERENCES bom_factory_baseline(baseline_id) ON DELETE CASCADE,
  role_code          VARCHAR2(40),                  -- 'OPS_MGR', 'SEC_MGR_SMT', ...
  display_name_en    VARCHAR2(100),
  category           VARCHAR2(40),                  -- MFG_IDL / CENTRALIZED_SERVICE
  annual_rate_usd    NUMBER(15,2),                  -- $60,320 ← 廠 EPM 年度更新
  weekly_rate_usd    NUMBER(15,2),                  -- annual / 50
  PRIMARY KEY (baseline_id, role_code)
);

-- (2) 案級 multiplier matrix · 一案 153 列 (17 × 9)
CREATE TABLE bom_cs_case_idl_alloc (
  case_factory_id      NUMBER REFERENCES bom_cs_case_factory(case_factory_id) ON DELETE CASCADE,
  process_code         VARCHAR2(40) REFERENCES bom_process_catalog(process_code),
  role_code            VARCHAR2(40),                -- 對應 bom_factory_idl_role.role_code
  multiplier           NUMBER(10,6),                -- 0.08, 0.125, 0.005 ...
  PRIMARY KEY (case_factory_id, process_code, role_code)
);

-- Step 2 完成後 · 案級 153 列寫入
-- 計算引擎(Step 5 區 B) JOIN 兩表:
--   SELECT SUM(r.annual_rate_usd × a.multiplier) AS idl_annual_for_process
--   FROM bom_cs_case_idl_alloc a
--   JOIN bom_factory_idl_role r ON r.role_code = a.role_code
--   WHERE a.case_factory_id = ? AND a.process_code = ?

-- 範例:SMT_MAIN 製程的 IDL 年成本
--   = OPS_MGR $60,320 × 0.08
--   + SEC_MGR_SMT $42,500 × 0.500
--   + ENGINEER $24,440 × 0.300
--   + ... (17 角色 SUMPRODUCT)
--   = $51,234 / 年 → / 50 / weekly_output = $0.0436 / unit"""
    add_text(s, 0.6, 1.8, 12.1, 5.0, schema, size=8.5, color=TEXT, font='Consolas')

    footer(s, page_n, total_pages)


# ════════════════════════════════════════════════════════════════════════
# STEP 3 — 設備 binding(5 slides)
# ════════════════════════════════════════════════════════════════════════

def slide_step3_intro(total_pages, page_n):
    s = add_slide()
    header(s, '案級設備 binding · 概念', '廠 master 58 件 → 案級挑 28 件 · 每件設 line_qty × qty_per_line', 3, '設備 binding')

    add_rounded(s, 0.4, 1.3, 6.2, 5.4, BG_SOFT, line=CYAN, line_width=1.5, radius=0.03)
    add_rounded(s, 0.4, 1.3, 6.2, 0.4, CYAN, radius=0.03)
    add_text(s, 0.55, 1.35, 5.9, 0.3, '💡 雙層架構 · 為什麼?', size=12, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, 0.6, 1.85, 5.9, 4.8, """廠 master(bom_factory_equipment)= 廠內所有設備
案級 binding(bom_cs_case_equipment)= 本案用多少

問題:DEK 印刷機是同一台 · SteelSeries 案吃 50%
       Razer 案吃 50% · capacity 怎麼算?
解法:廠級存設備本身 · 案級存「我這案佔用 line_qty」

範例(DEK 印刷機):
  廠 master:
    desc='DEK Horizon 03i Printer'
    acq_cost=$84,000 · life=6 yrs
  SteelSeries 案 binding:
    line_qty=1 · qty_per_line=1
  Razer 案 binding:
    line_qty=0.5 · qty_per_line=1 (共用)

→ 折舊算法:
  SteelSeries 案吃 $84k / 6 × 1×1 = $14k/年
  Razer 案吃 $84k / 6 × 0.5×1 = $7k/年
  合理分攤

→ 跨案分析:
  同款 DEK 跨多案 utilization 一秒可查""", size=10, color=TEXT)

    add_rounded(s, 6.8, 1.3, 6.15, 5.4, WHITE, line=NAVY, line_width=1.5, radius=0.03)
    add_rounded(s, 6.8, 1.3, 6.15, 0.4, NAVY, radius=0.03)
    add_text(s, 6.95, 1.35, 5.85, 0.3, '📊 設備分類(6 type)', size=12, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)

    types = [
        ('EQUIPMENT', '主機設備', '高額 · DEK / FUJI NXT / Reflow', RGBColor(0x08,0x91,0xB2), 28),
        ('FIXTURE',   '治具',     '工裝 · ESD Tray / 模具',          GREEN, 8),
        ('TOOLING',   '工具',     'Test fixture / Microscope',       GOLD, 5),
        ('CARRIER',   '載具',     'ESD Tray / Carrier (752 個/line)', PURPLE, 4),
        ('STENCIL',   '鋼網',     'SMT 印刷模板(壽命短)',           RED, 3),
        ('MRO',       'MRO 物料', 'UPS / 壓縮空氣罐 / 廠級共用',     RGBColor(0x6B,0x72,0x80), 10),
    ]
    for i, (code, name, ex, c, cnt) in enumerate(types):
        y = 1.85 + i * 0.78
        add_rounded(s, 6.95, y, 5.85, 0.7, BG_SOFT, line=c, line_width=1, radius=0.02)
        add_rounded(s, 6.95, y, 1.05, 0.7, c, radius=0.02)
        add_text(s, 7.0, y + 0.05, 0.95, 0.3, code, size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, 7.0, y + 0.35, 0.95, 0.3, f'廠 {cnt} 件', size=8, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, 8.15, y + 0.05, 4.6, 0.3, name, size=11, bold=True, color=INK, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, 8.15, y + 0.35, 4.6, 0.3, ex, size=8.5, color=TEXT, anchor=MSO_ANCHOR.MIDDLE)

    footer(s, page_n, total_pages)


def slide_step3_factory_view(total_pages, page_n):
    s = add_slide()
    header(s, '設備 binding · 廠 master view', 'v0.10 新加 · 點 [🏭 廠 master] toggle 看廠 58 件全清單', 3)

    add_text(s, 0.4, 1.3, 12, 0.3, '🎯 點什麼:Step 3 panel → toggle bar 切 [🏭 廠 master (58 件 · 可選用)]',
             size=12, bold=True, color=INK)

    mockup = """┌───────────────────────────────────────────────────────────────────────────┐
│  🏭 Step 3 · 設備 binding                                                 │
│   廠 master 58 件 · 對應 bom_factory_equipment(本案未用的灰色)            │
│   廠 58 件 · 廠 acq Σ $3.85M                                               │
│                                                                           │
│  [📋 案級 (28 件用於本案)]  [🏭 廠 master (58 件 · 可選用) ●]              │
│                                                                           │
│  ┌───────────────────────────────────────────────────────────────────────┐│
│  │ Equipment 描述              Type   Process    案內  Line  Qty  Life  ││
│  │ ───────────────────────────────────────────────────────────────────  ││
│  │ DEK Horizon 03i Printer    EQUIP  SMT_MAIN    ✓     1     1   6yr   ││ ← 案內(亮)
│  │ FUJI NXT M3S×6 + M6S×1     EQUIP  SMT_MAIN    ✓     1     1   6yr   ││
│  │ SPI TRI-7007               EQUIP  SMT_MAIN    ✓     1     1   6yr   ││
│  │ ...                                                                  ││
│  │ DEK ELA Auto Printer       EQUIP  SMT_MAIN    —     —     —   6yr   ││ ← 案外(灰)
│  │ Universal Genesis P&P      EQUIP  SMT_MAIN    —     —     —   6yr   ││   廠 master · 本案未選
│  │ Koh-Young SPI 8030         EQUIP  SMT_MAIN    —     —     —   6yr   ││
│  │ ...                                                                  ││
│  └───────────────────────────────────────────────────────────────────────┘│
│                                                                           │
│  ↳ 廠 master schema: bom_factory_equipment 全廠 58 件                     │
└───────────────────────────────────────────────────────────────────────────┘"""
    add_rounded(s, 0.4, 1.7, 12.55, 4.8, WHITE, line=CYAN, line_width=1.5, radius=0.02)
    add_text(s, 0.55, 1.85, 12.3, 4.6, mockup, size=8, color=TEXT, font='Consolas')

    add_rounded(s, 0.4, 6.6, 12.55, 0.35, RGBColor(0xE0, 0xF2, 0xFE), line=CYAN, radius=0.04)
    add_text(s, 0.55, 6.62, 12.3, 0.32,
             '💡 用途:看到「廠級有 30 件本案沒用」· 可能本案漏選 · 或證明本案精簡',
             size=10, color=INK, anchor=MSO_ANCHOR.MIDDLE)

    footer(s, page_n, total_pages)


def slide_step3_case_view(total_pages, page_n):
    s = add_slide()
    header(s, '設備 binding · 案級 view', '預設視角 · 案內 28 件 · 每件設 line_qty / qty_per_line / 案內小計', 3)

    mockup = """┌──────────────────────────────────────────────────────────────────────────────────┐
│  🏭 Step 3 · 設備 binding · 案內 28 件 · 案內 acq Σ $1.6M                          │
│                                                                                  │
│  [📋 案級 (28 件用於本案) ●]  [🏭 廠 master (58 件 · 可選用)]                     │
│                                                                                  │
│  Equipment 描述               Type      Process     Line qty  Qty/line Life  Acq Cost   案內小計│
│  ──────────────────────────────────────────────────────────────────────────────  │
│  DEK Horizon 03i Printer     EQUIP     SMT_MAIN     1         1        6     $84,000   $84,000│
│  FUJI NXT M3S×6 + M6S×1      EQUIP     SMT_MAIN     1         1        6     $1,013k   $1,013k│
│  SPI TRI-7007 (Solder Paste) EQUIP     SMT_MAIN     1         1        6     $68,852   $68,852│
│  Heller Reflow Oven Mark5    EQUIP     SMT_MAIN     1         1        6     $158,000  $158k  │
│  ESD Tray (carrier)          CARRIER   SMT_MAIN     1         752      3     $18.5     $13,912│ ← qty/line=752
│  ...                                                                                          │
│  Wave Solder ERSA VS280-A    EQUIP     WAVE_SOLDER  1         1        8     $135,000  $135k  │
│  ...                                                                                          │
│  EMI/EMC Chamber (shared)    EQUIP     BB_TEST      0.3       1        10    $480,000  $144k  │ ← 跨案共用 30%
│  X-Ray Pass-thru (TR7500S)   EQUIP     Q_SMT        0.5       1        8     $285,000  $143k  │ ← 跨案 50%
└──────────────────────────────────────────────────────────────────────────────────┘"""
    add_rounded(s, 0.4, 1.3, 12.55, 4.0, WHITE, line=CYAN, line_width=1.5, radius=0.02)
    add_text(s, 0.55, 1.42, 12.3, 3.85, mockup, size=7.5, color=TEXT, font='Consolas')

    # 注意事項
    add_rounded(s, 0.4, 5.5, 12.55, 1.45, RGBColor(0xE0, 0xF2, 0xFE), line=CYAN, radius=0.04)
    add_text(s, 0.55, 5.55, 12.3, 1.4, """💡 line_qty / qty_per_line 怎麼填?
  • line_qty:本案吃這台設備多少「線」· 跨案共用設備 → 用小數(0.3 = 吃 30% capacity)
  • qty_per_line:該線需要幾台這設備(ESD Tray 一線需 752 個 carrier · 設備數 ≠ 1)
  • life_yrs:可 override 廠 master default(汰換 / 客戶壓力提早攤提)
  • acq_cost:可 override(汰換 / 議價 / 客供設備 = $0)

→ 折舊算法(Step 5 區 C):depreciation = line_qty × qty_per_line × acq_cost / life_yrs""",
             size=9.5, color=TEXT, anchor=MSO_ANCHOR.TOP)

    footer(s, page_n, total_pages)


def slide_step3_picking(total_pages, page_n):
    s = add_slide()
    header(s, '設備 binding · 怎麼挑 / 怎麼填', '從廠 master 挑 → 設 line_qty/qty_per_line → 可 override', 3)

    # 操作流程
    add_text(s, 0.4, 1.3, 12, 0.3, '🎬 EPM 操作流程(以 SMT_MAIN 製程為例)', size=12, bold=True, color=INK)

    steps = [
        ('1️⃣', '切到廠 master view', '看廠級 58 件 · 過濾 process=SMT_MAIN · 有 12 件 SMT 設備可選'),
        ('2️⃣', '案內必選的「線級」設備', 'DEK / FUJI NXT / SPI / Reflow / AOI · 一條 SMT 線必備 → 全勾'),
        ('3️⃣', '案內必選的「個數型」設備', 'ESD Tray / Stencil → 用 qty_per_line(本案 PCB 一線需 752 個 ESD Tray)'),
        ('4️⃣', '案內可選的「shared」設備', 'X-Ray / EMI Chamber → line_qty < 1(本案佔 50% / 30%)'),
        ('5️⃣', '若廠級無此設備 → 先回廠 master 加', 'Admin → Factory Equipment → Add new equipment'),
        ('6️⃣', '可 override 廠 master default', 'acq_cost / life_yrs · 案內議價或壓力提早攤提時用'),
    ]
    for i, (n, t, d) in enumerate(steps):
        y = 1.75 + i * 0.78
        add_rounded(s, 0.4, y, 12.55, 0.7, BG_SOFT, line=CYAN, line_width=0.5, radius=0.04)
        add_text(s, 0.55, y + 0.18, 0.5, 0.35, n, size=18, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, 1.15, y + 0.05, 4.5, 0.3, t, size=11, bold=True, color=INK, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, 1.15, y + 0.35, 11.7, 0.3, d, size=9.5, color=TEXT, anchor=MSO_ANCHOR.MIDDLE)

    footer(s, page_n, total_pages)


def slide_step3_schema(total_pages, page_n):
    s = add_slide()
    header(s, '設備 binding · 寫入 DB', '廠 master + 案級 binding 兩層 schema', 3)

    add_rounded(s, 0.4, 1.3, 12.55, 5.4, BG_SOFT, line=CYAN, line_width=1.5, radius=0.03)
    add_rounded(s, 0.4, 1.3, 12.55, 0.4, CYAN, radius=0.03)
    add_text(s, 0.55, 1.35, 12.3, 0.3, '🗄️ bom_factory_equipment(廠級)+ bom_cs_case_equipment(案級)', size=12, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)

    schema = """-- (1) 廠級設備 master · 跨案共用 · 58 列 / 廠
CREATE TABLE bom_factory_equipment (
  equipment_id            NUMBER GENERATED ALWAYS AS IDENTITY PRIMARY KEY,
  factory_code            VARCHAR2(10) REFERENCES bom_factory(factory_code),
  equipment_code          VARCHAR2(80),                  -- 廠內編號
  description             VARCHAR2(500),                 -- 'DEK Horizon 03i Printer'
  manufacturer            VARCHAR2(100),                 -- 'DEK'
  model                   VARCHAR2(100),                 -- 'Horizon 03i'
  part_no                 VARCHAR2(80),
  equipment_type          VARCHAR2(40),                  -- EQUIPMENT/FIXTURE/TOOLING/CARRIER/STENCIL/MRO
  custom_or_standard      VARCHAR2(20),                  -- CUSTOM / STANDARD
  default_process_code    VARCHAR2(40),                  -- 'SMT_MAIN'
  acquisition_cost_usd    NUMBER(15,2),                  -- $84,000 ← 廠 EPM 月度更新
  salvage_value_usd       NUMBER(15,2) DEFAULT 0,
  default_useful_life_yrs NUMBER(5,2),                   -- 6
  is_active               NUMBER(1) DEFAULT 1
);

-- (2) 案級設備 binding · 案內挑 + 設量 + 可 override · 28 列 / 案
CREATE TABLE bom_cs_case_equipment (
  case_factory_id              NUMBER REFERENCES bom_cs_case_factory(case_factory_id) ON DELETE CASCADE,
  equipment_id                 NUMBER REFERENCES bom_factory_equipment(equipment_id),
  process_code                 VARCHAR2(40),              -- 可改廠 master default
  line_qty                     NUMBER(5,2),               -- 1.0 / 0.5 / 0.3 ← 跨案共用比例
  qty_per_line                 NUMBER(5),                 -- 752 (ESD Tray)
  useful_life_override_yrs     NUMBER(5,2),               -- NULL → 用廠 default 6 yr
  acquisition_cost_override_usd NUMBER(15,2),             -- NULL → 用廠 default $84,000
  notes                        CLOB,
  PRIMARY KEY (case_factory_id, equipment_id)
);

-- (3) 計算引擎(Step 5 區 C)JOIN 兩表:
SELECT
  SUM(e.line_qty × e.qty_per_line ×
      COALESCE(e.acquisition_cost_override_usd, m.acquisition_cost_usd) ×
      0.05) AS annual_mro,                                              -- MRO 廠政策 5%
  SUM(e.line_qty × e.qty_per_line ×
      COALESCE(e.acquisition_cost_override_usd, m.acquisition_cost_usd) /
      COALESCE(e.useful_life_override_yrs, m.default_useful_life_yrs)) AS annual_depr
FROM bom_cs_case_equipment e
JOIN bom_factory_equipment m ON m.equipment_id = e.equipment_id
WHERE e.case_factory_id = ? AND e.process_code = ?;"""
    add_text(s, 0.6, 1.8, 12.1, 5.0, schema, size=8, color=TEXT, font='Consolas')

    footer(s, page_n, total_pages)


# ════════════════════════════════════════════════════════════════════════
# STEP 4 — 耗材 binding(3 slides)
# ════════════════════════════════════════════════════════════════════════

def slide_step4_intro(total_pages, page_n):
    s = add_slide()
    header(s, '案級耗材 binding · 概念', '17 件耗材 · annual_qty × unit_cost = 年成本 · /unit 計算', 4, '耗材 binding')

    add_rounded(s, 0.4, 1.3, 6.2, 5.4, BG_SOFT, line=GOLD, line_width=1.5, radius=0.03)
    add_rounded(s, 0.4, 1.3, 6.2, 0.4, GOLD, radius=0.03)
    add_text(s, 0.55, 1.35, 5.9, 0.3, '💡 為什麼分耗材跟設備?', size=12, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, 0.6, 1.85, 5.9, 4.8, """設備:資產 · 一次買進 · 折舊攤提 · cost 算 Equipment 區
耗材:消費 · 每年買掉 · 直接吃成本 · cost 算 Facility 區 D
       (Indirect Materials)

耗材特徵:
  • 用量大 (錫膏每年 280 KG)
  • 單價低 (錫膏 $52/KG)
  • 隨年量 scale(年量 × 2 → 耗材 × 2)

對比設備(DEK $84K · 6 年攤提 · 不隨年量變):
  耗材完全 demand-driven · 年量 100K → 耗材年成本 X
  年量 418K → 耗材年成本 4.18X

17 件耗材:
  • 錫膏 / 錫條 / 助焊劑 (SMT 必備)
  • 螺絲 / 熱熔膠 / OPP 袋 (BB Assy)
  • 測試 cable / 校正器材 (BB Test)
  • Pallet / 纏繞膜 (Material Mgmt)
  • Probe / 樣品袋 (Quality)""", size=10, color=TEXT)

    add_rounded(s, 6.8, 1.3, 6.15, 5.4, WHITE, line=NAVY, line_width=1.5, radius=0.03)
    add_rounded(s, 6.8, 1.3, 6.15, 0.4, NAVY, radius=0.03)
    add_text(s, 6.95, 1.35, 5.85, 0.3, '📊 17 件耗材對 SteelSeries 案', size=12, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    table = """Consumable                       UOM  Proc    Annual qty Unit$  年成本
─────────────────────────────────────────────────────────────────────
錫膏 Solder Paste SAC305 (500g) KG   SMT     280       $52.00 $14,560
清潔劑(鋼網)                     L    SMT      85      $8.50    $722
錫條 Wave Solder Bar             KG   Wave    420     $28.00 $11,760
助焊劑 Flux                      L    Wave     65     $12.50    $812
分板膠 / 刀片                    PCS  Router 1,200      $3.20  $3,840
雷射護目鏡                       PCS  Laser    24     $35.00    $840
螺絲 M2×6mm                      PCS  BB-Assy 2,508,000 $0.005 $12,040
熱熔膠條                          KG   BB-Assy   180     $6.20  $1,116
OPP 防靜電袋                     PCS  BB-Assy 418,000  $0.018  $7,524
氣動工具油                       L    BB-Assy    24     $22.00   $528
測試 cable (USB-A)               PCS  BB-Test    60      $3.50   $210
校正器材                          PCS  BB-Test    12     $85.00 $1,020
Pallet 棧板                       PCS  Mat-Mgmt  800      $8.50 $6,800
纏繞膜                            ROLL Mat-Mgmt  240     $12.00 $2,880
ICT Probe                         PCS  Q-SMT     600      $1.20    $720
品保樣品袋                       PCS  Q-BB       60      $0.45     $27
X-Ray 底片                       PCS  Q-SMT     480     $12.00  $5,760

Σ Total 年成本                                                  $71,159
÷ Annual Demand 418,000                                              -
= per unit                                                       $0.0017"""
    add_text(s, 6.95, 1.78, 5.85, 4.9, table, size=6.5, color=TEXT, font='Consolas')

    footer(s, page_n, total_pages)


def slide_step4_operations(total_pages, page_n):
    s = add_slide()
    header(s, '耗材 binding · UI + 操作', 'Step 4 panel · 17 件 list + 年總成本 Σ + per unit', 4)

    add_text(s, 0.4, 1.3, 12, 0.3, '🎯 操作流程:從廠 master 17 件耗材 binding 進案 · 設 annual_qty',
             size=12, bold=True, color=INK)

    # 左:UI mockup
    mockup = """┌─ 📦 Step 4 · 案級耗材 binding ────────────────────────────────────┐
│  17 件耗材 · 年總成本 Σ $71.2k                                    │
│                                                                   │
│  ID  Description           UOM  Proc       Annual qty  Unit cost  │
│  ────────────────────────────────────────────────────────────────  │
│  C01 錫膏 Solder Paste     KG   SMT_MAIN     280       $52.0000   │
│  C02 清潔劑(鋼網)          L    SMT_MAIN      85        $8.5000   │
│  C03 錫條 Wave Solder      KG   WAVE_SOLDER  420       $28.0000   │
│  C04 助焊劑 Flux           L    WAVE_SOLDER   65       $12.5000   │
│  C05 分板膠 / 刀片         PCS  ROUTER...   1,200       $3.2000   │
│  ...                                                              │
│                                                                   │
│  Σ TOTAL                                              $71,159     │
│  ÷ Annual 418K                                       $0.0017/unit │
└───────────────────────────────────────────────────────────────────┘"""
    add_rounded(s, 0.4, 1.7, 6.5, 5.0, WHITE, line=GOLD, line_width=1.5, radius=0.02)
    add_text(s, 0.55, 1.85, 6.3, 4.8, mockup, size=8, color=TEXT, font='Consolas')

    # 右:操作步驟
    add_rounded(s, 7.1, 1.7, 5.85, 5.0, BG_SOFT, line=GOLD, line_width=1.5, radius=0.03)
    add_rounded(s, 7.1, 1.7, 5.85, 0.4, GOLD, radius=0.03)
    add_text(s, 7.25, 1.75, 5.55, 0.3, '🎬 操作步驟', size=12, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    ops_md = """1️⃣  廠 EPM 第一次 import 時
    Consumables sheet → 系統自動建廠 master
    17 件 unit_cost / UOM / annual_usage_default

2️⃣  案級 binding 預設帶入廠 master 17 件
    annual_qty 從 default 帶 · EPM 視案調整
    (年量大 → 耗材按比例增)

3️⃣  動態欄位
    • annual_qty:本案年用量
      (預設 = 廠 default · 可改)
    • unit_cost_override_usd:可改
      (供應商換 / 議價 / 客供 = 0)

4️⃣  系統算
    年成本 = annual_qty × COALESCE(override,
                                    廠 master.unit_cost)
    Σ → / Annual demand → /unit

5️⃣  跑 Step 5 區 D · 寫入 IND_MAT
    cost_per_unit_usd

6️⃣  若耗材廠裡新增 → 先回廠 master 加"""
    add_text(s, 7.25, 2.2, 5.55, 4.5, ops_md, size=10, color=TEXT)

    footer(s, page_n, total_pages)


def slide_step4_schema(total_pages, page_n):
    s = add_slide()
    header(s, '耗材 binding · 寫入 DB', 'bom_factory_consumable + bom_cs_case_consumable', 4)

    add_rounded(s, 0.4, 1.3, 12.55, 5.4, BG_SOFT, line=GOLD, line_width=1.5, radius=0.03)
    add_rounded(s, 0.4, 1.3, 12.55, 0.4, GOLD, radius=0.03)
    add_text(s, 0.55, 1.35, 12.3, 0.3, '🗄️ bom_factory_consumable(廠級)+ bom_cs_case_consumable(案級)', size=12, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)

    schema = """-- (1) 廠級耗材 master · 跨案共用
CREATE TABLE bom_factory_consumable (
  consumable_id          NUMBER GENERATED ALWAYS AS IDENTITY PRIMARY KEY,
  factory_code           VARCHAR2(10) REFERENCES bom_factory(factory_code),
  consumable_code        VARCHAR2(80),
  description            VARCHAR2(500),                 -- '錫膏 Solder Paste SAC305 (500g)'
  default_process_code   VARCHAR2(40),                  -- 'SMT_MAIN'
  unit_cost_usd          NUMBER(15,4),                  -- $52.0000 / KG ← 廠 EPM 月度更新
  unit_of_measure        VARCHAR2(20),                  -- KG/L/PCS/ROLL
  annual_usage_default   NUMBER,                        -- 280 (廠級每年消耗量,參考)
  is_active              NUMBER(1) DEFAULT 1
);

-- (2) 案級耗材 binding · 案內 17 件
CREATE TABLE bom_cs_case_consumable (
  case_factory_id         NUMBER REFERENCES bom_cs_case_factory(case_factory_id) ON DELETE CASCADE,
  consumable_id           NUMBER REFERENCES bom_factory_consumable(consumable_id),
  process_code            VARCHAR2(40),
  annual_usage_qty        NUMBER(15,4),                  -- 280 ← 案級可改 · 年量大 → 多用
  unit_cost_override_usd  NUMBER(15,4),                  -- NULL → 用廠 default · 議價時 override
  notes                   CLOB,
  PRIMARY KEY (case_factory_id, consumable_id)
);

-- (3) 計算引擎(Step 5 區 D · IND_MAT)JOIN:
SELECT
  SUM(c.annual_usage_qty × COALESCE(c.unit_cost_override_usd, m.unit_cost_usd))
  AS annual_indirect_materials
FROM bom_cs_case_consumable c
JOIN bom_factory_consumable m ON m.consumable_id = c.consumable_id
WHERE c.case_factory_id = ? AND c.process_code = ?;

-- → / Annual demand → per unit
-- 範例 SMT_MAIN:錫膏 + 清潔劑 = ($14,560 + $722) / 418,000 = $0.0365/unit"""
    add_text(s, 0.6, 1.8, 12.1, 5.0, schema, size=8.5, color=TEXT, font='Consolas')

    footer(s, page_n, total_pages)


# ════════════════════════════════════════════════════════════════════════
# STEP 5 — Compute trace 5 區公式(10 slides)
# ════════════════════════════════════════════════════════════════════════

def slide_step5_overview(total_pages, page_n):
    s = add_slide()
    header(s, 'Compute Trace · 5 區公式 overview', 'A DL · B IDL · C Equipment · D Facility · E Others · F Total', 5, 'Compute')

    add_text(s, 0.4, 1.3, 12, 0.3, '🧮 6 個 zone tabs · 對應 Cleansheet Excel 5 區公式 + 1 個 Total Roll-up',
             size=12, bold=True, color=INK)

    zones = [
        ('A', '👷 Direct Labor',       9, 'DL_CPU',          'DL_wage × SEA × Total_DL × mult1 × mult2 / weekly_output',           'C58 (對應 Excel)'),
        ('B', '👥 Indirect Labor',     4, 'IDL_CPU',         'Σ(role.annual × multiplier) / 50 / weekly_output',                   'Y70'),
        ('C', '🏭 Equipment',          6, 'EQUIP_MRO + DEPR', '(Σ acq×5% + Σ acq/life) × line_qty × qty/line / annual_output',     'H90/H91'),
        ('D', '🏗️ Facility + Ind Mat', 6, 'FACILITY + IND_MAT', '(sqft × $/sqft + Σ consumable_annual) / annual_output',          'J100'),
        ('E', '📦 Others (Common)',    4, 'FREIGHT + VAT + LOSS', 'freight/demand + motherboard×vat/18 + motherboard×loss · 全廠', 'J107 / J110 / J114'),
        ('F', '💰 Total Roll-up',      4, 'TC',              'motherboard + MVA + motherboard×sga% + (MVA+motherboard)×profit%',   'C125'),
    ]
    for i, (code, label, steps, sch, formula, ref) in enumerate(zones):
        y = 1.75 + i * 0.85
        # zone badge
        add_rounded(s, 0.4, y, 1.2, 0.7, RED, radius=0.04)
        add_text(s, 0.4, y + 0.05, 1.2, 0.3, f'區 {code}', size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, 0.4, y + 0.4, 1.2, 0.3, f'{steps} 步', size=9, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # content card
        add_rounded(s, 1.7, y, 11.25, 0.7, BG_SOFT, line=RED, line_width=0.5, radius=0.04)
        add_text(s, 1.85, y + 0.05, 4, 0.3, label, size=12, bold=True, color=INK, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, 5.9, y + 0.05, 3, 0.3, f'→ schema: {sch}', size=9, color=CYAN, anchor=MSO_ANCHOR.MIDDLE, font='Consolas')
        add_text(s, 9.0, y + 0.05, 4, 0.3, f'Excel ref: {ref}', size=9, color=MUTED, anchor=MSO_ANCHOR.MIDDLE, font='Consolas')
        add_text(s, 1.85, y + 0.38, 11, 0.32, formula, size=8.5, color=TEXT, font='Consolas', anchor=MSO_ANCHOR.MIDDLE)

    footer(s, page_n, total_pages)


def slide_step5_zone_A_dl(total_pages, page_n):
    s = add_slide()
    header(s, '區 A · Direct Labor · 9 步驟', 'compute_dl_cost(SMT_MAIN) · 用 CN $4.95/hr · 對應 SD §7.2', 5)

    add_text(s, 0.4, 1.3, 12, 0.3, '🎯 EPM 點「▶ 下一步」9 次 · 系統逐步揭露中間值',
             size=12, bold=True, color=INK)

    # 9 步驟卡(2 col × 5 row, last row 1 col)
    dl_steps = [
        ('1', 'UPH (units / hour)',       '3600 / TAKT × Yield × Eff',
            '3600 / 8.64 × 98% × 95% = 388.06',          'Cleansheet!C14'),
        ('2', 'Max output / line / week', 'Working hr/wk × UPH',
            '153.41 × 388.06 = 59,531',                  'Cleansheet!C21'),
        ('3', 'Demand max/wk',            'Annual / 50 × 1.2',
            '418,000 / 50 × 1.2 = 10,032',               'Cleansheet!C18-20'),
        ('4', 'Lines installed',          'ROUNDUP(max / max_output_per_line)',
            'ROUNDUP(10,032 / 59,531) = 1',              'Cleansheet!C23'),
        ('5', 'Weekly output',            'max_output_per_line × lines',
            '59,531 × 1 = 59,531',                       'Cleansheet!C24'),
        ('6', 'Total DL / day',           '(DL×lines×shifts) + (debug×debug_lines×shifts) + (functional×shifts)',
            '(10×1×2) + (3×0×2) + (2×2) = 24',           'Cleansheet!C33'),
        ('7', 'Multiplier 1 + 2',         'mult1 = (Whr/2 × 6) / SEA_wk',
            'mult1 = (23/2 × 6)/60 = 1.15  mult2 = (10×6.67)/60 = 1.112',  'C50, C51'),
        ('8', 'DL cost / week',           'DL_wage × SEA_wk × Total_DL × mult2 × mult1',
            '$4.95 × 60 × 24 × 1.112 × 1.15 = $9,116',   'Cleansheet!C56'),
        ('9', 'DL cost / unit',           '(DL/wk + IDL_line_dep) / weekly_output',
            '($9,116 + $645) / 59,531 = $0.1640',        'Cleansheet!C58 (FINAL)'),
    ]
    for i, (n, label, formula, calc, ref) in enumerate(dl_steps):
        col = i % 2
        row = i // 2
        x = 0.4 + col * 6.27
        y = 1.7 + row * 1.0
        add_rounded(s, x, y, 6.2, 0.92, WHITE, line=RED, line_width=0.5, radius=0.03)
        add_rounded(s, x, y + 0.05, 0.45, 0.82, RED, radius=0.5)
        add_text(s, x, y + 0.05, 0.45, 0.82, n, size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + 0.55, y + 0.06, 4.5, 0.25, label, size=10, bold=True, color=INK, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + 5.15, y + 0.06, 1.0, 0.25, ref, size=7.5, color=MUTED, anchor=MSO_ANCHOR.MIDDLE, font='Consolas', align=PP_ALIGN.RIGHT)
        add_text(s, x + 0.55, y + 0.34, 5.6, 0.22, f'公式: {formula}', size=8, color=MUTED, font='Consolas')
        add_text(s, x + 0.55, y + 0.6, 5.6, 0.3, f'↳ {calc}', size=9, color=PURPLE, font='Consolas')

    footer(s, page_n, total_pages)


def slide_step5_zone_B_idl(total_pages, page_n):
    s = add_slide()
    header(s, '區 B · Indirect Labor · 4 步驟', 'compute_idl_cost(SMT_MAIN) · 17 角色 SUMPRODUCT', 5)

    add_text(s, 0.4, 1.3, 12, 0.3, '🎯 切到 zone B → 跑 4 步 → 寫入 IDL_CPU',
             size=12, bold=True, color=INK)

    idl_steps = [
        ('1', '17 角色 × multiplier × 年薪',
            'Σ (role.annual × multiplier(role, process))',
            'Cleansheet rows 64-83',
            'OPS_MGR $60,320 × 0.080 = $4,826  +\nSEC_MGR_SMT $42,500 × 0.500 = $21,250 +\nENGINEER $24,440 × 0.300 = $7,332 + ...\n(17 角色 SUMPRODUCT) = $51,234 / yr'),
        ('2', 'IDL per week',
            'IDL_annual / 50',
            'Cleansheet!Y70',
            '$51,234 / 50 = $1,025 / wk'),
        ('3', 'Weekly output (從 §A.5 帶)',
            'Whr/wk × UPH × lines',
            '§A Step 5',
            '153.41 × 388.06 × 1 = 59,531 u/wk'),
        ('4', 'IDL cost / unit (FINAL)',
            'IDL_per_week / weekly_output',
            'Cleansheet!Y70/C24',
            '$1,025 / 59,531 = $0.0172 / unit'),
    ]
    for i, (n, label, formula, ref, calc) in enumerate(idl_steps):
        y = 1.7 + i * 1.25
        add_rounded(s, 0.4, y, 12.55, 1.15, WHITE, line=RED, line_width=0.5, radius=0.03)
        add_rounded(s, 0.4, y + 0.05, 0.65, 1.05, RED, radius=0.5)
        add_text(s, 0.4, y + 0.05, 0.65, 1.05, n, size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, 1.2, y + 0.05, 7, 0.3, label, size=12, bold=True, color=INK)
        add_text(s, 11.4, y + 0.05, 1.5, 0.3, ref, size=9, color=MUTED, font='Consolas', align=PP_ALIGN.RIGHT)
        add_text(s, 1.2, y + 0.35, 11.5, 0.25, f'公式: {formula}', size=9, color=MUTED, font='Consolas')
        add_text(s, 1.2, y + 0.62, 11.5, 0.5, f'↳ {calc}', size=9, color=PURPLE, font='Consolas')

    footer(s, page_n, total_pages)


def slide_step5_zone_C_equip(total_pages, page_n):
    s = add_slide()
    header(s, '區 C · Equipment · 6 步驟', 'compute_equipment_cost(SMT_MAIN) · MRO 5% + Depreciation', 5)

    equip_steps = [
        ('1', '抽 process 名下的設備',
            'WHERE proc = SMT_MAIN', 'bom_cs_case_equipment',
            'DEK + FUJI NXT + SPI + Reflow + AOI + ESD Tray + Stencil = 7 件'),
        ('2', '總 acquisition cost',
            'Σ (line_qty × qty_per_line × acq_cost)', 'Equipment List col P',
            'DEK $84k×1×1 + FUJI $1,013k + SPI $69k + Reflow $158k + AOI $92k\n+ ESD Tray $18.5×752 + Stencil $550×6 = $1,433k'),
        ('3', 'Annual Depreciation',
            'Σ (acq / useful_life)', 'Cleansheet!H91',
            '$1,433k / 6 yrs ≈ $238,883 / yr (按各設備 life 加權)'),
        ('4', 'Annual MRO',
            'acq × 5% (廠 maintenance policy)', 'Cleansheet!H90',
            '$1,433k × 5% = $71,650 / yr'),
        ('5', 'Annual output ratio',
            'weekly_output × 50 (從 §A.5)', 'Cleansheet!C20/C24',
            '59,531 × 50 = 2,976,550 u/yr (>>418K 因 capacity 大)'),
        ('6', 'MRO + Depr / unit (FINAL)',
            '(MRO + Depr) / annual_output', 'Cleansheet!H90+H91',
            '($71,650 + $238,883) / 2,976,550 = $0.1043 / unit'),
    ]
    for i, (n, label, formula, ref, calc) in enumerate(equip_steps):
        y = 1.45 + i * 0.92
        add_rounded(s, 0.4, y, 12.55, 0.83, WHITE, line=RED, line_width=0.5, radius=0.03)
        add_rounded(s, 0.4, y + 0.05, 0.55, 0.73, RED, radius=0.5)
        add_text(s, 0.4, y + 0.05, 0.55, 0.73, n, size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, 1.05, y + 0.04, 7, 0.27, label, size=11, bold=True, color=INK)
        add_text(s, 11.0, y + 0.04, 2.0, 0.27, ref, size=9, color=MUTED, font='Consolas', align=PP_ALIGN.RIGHT)
        add_text(s, 1.05, y + 0.3, 11.5, 0.22, f'公式: {formula}', size=8.5, color=MUTED, font='Consolas')
        add_text(s, 1.05, y + 0.52, 11.5, 0.28, f'↳ {calc}', size=8.5, color=PURPLE, font='Consolas')

    footer(s, page_n, total_pages)


def slide_step5_zone_D_facility(total_pages, page_n):
    s = add_slide()
    header(s, '區 D · Facility + Indirect Materials · 6 步驟', 'compute_facility_cost(SMT_MAIN) · sqft + 耗材', 5)

    fac_steps = [
        ('1', '廠 floor $/sqft (baseline)',
            'baseline.floor_sqft_prod', 'SupplierBaseInput!B12',
            'CN baseline = $14.4276 / sqft'),
        ('2', 'Process 佔用 sqft',
            'allocation_per_process', 'Cleansheet rows 95-103',
            'SMT_MAIN = 800 sqft (SMT 主板貼裝佔大)'),
        ('3', 'Annual Facility cost',
            'sqft × $/sqft', 'Cleansheet!J100',
            '800 × $14.4276 = $11,542 / yr'),
        ('4', 'Annual Indirect Materials',
            'Σ (consumable.annual_qty × unit_cost) WHERE proc = SMT_MAIN', 'Consumables sheet',
            '錫膏 $14,560 + 清潔劑 $722 = $15,282 / yr'),
        ('5', 'Annual output (從 §A.5)',
            'weekly_output × 50', 'Cleansheet!C24',
            '59,531 × 50 = 2,976,550 u/yr'),
        ('6', 'Facility + IM / unit (FINAL)',
            '(facility + ind_mat) / annual_output', 'Cleansheet!J100',
            '($11,542 + $15,282) / 2,976,550 = $0.0090 / unit'),
    ]
    for i, (n, label, formula, ref, calc) in enumerate(fac_steps):
        y = 1.45 + i * 0.92
        add_rounded(s, 0.4, y, 12.55, 0.83, WHITE, line=RED, line_width=0.5, radius=0.03)
        add_rounded(s, 0.4, y + 0.05, 0.55, 0.73, RED, radius=0.5)
        add_text(s, 0.4, y + 0.05, 0.55, 0.73, n, size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, 1.05, y + 0.04, 7, 0.27, label, size=11, bold=True, color=INK)
        add_text(s, 11.0, y + 0.04, 2.0, 0.27, ref, size=9, color=MUTED, font='Consolas', align=PP_ALIGN.RIGHT)
        add_text(s, 1.05, y + 0.3, 11.5, 0.22, f'公式: {formula}', size=8.5, color=MUTED, font='Consolas')
        add_text(s, 1.05, y + 0.52, 11.5, 0.28, f'↳ {calc}', size=8.5, color=PURPLE, font='Consolas')

    footer(s, page_n, total_pages)


def slide_step5_zone_E_others(total_pages, page_n):
    s = add_slide()
    header(s, '區 E · Others (Common) · 4 步驟', 'compute_others_cost · Freight + VAT + Loss · 全廠通用,不分製程', 5)

    others_steps = [
        ('1', 'Inbound Freight / unit',
            'inbound_freight_annual / annual_demand', 'Cleansheet!J107',
            '$2,500 / 418,000 = $0.00598 / unit'),
        ('2', 'VAT / unit',
            'motherboard × VAT_rate / 18', 'Cleansheet!J114',
            '$8.683 × 17% / 18 = $0.082 / unit  (CN 17% VAT)'),
        ('3', 'Loss factor / unit',
            'motherboard × loss_factor_pct', 'Cleansheet!J110',
            '$8.683 × 0.08% = $0.00695 / unit'),
        ('4', 'Others Total (FINAL)',
            'Freight + VAT + Loss', 'Cleansheet!Common 欄',
            '$0.00598 + $0.082 + $0.00695 = $0.09493 / unit'),
    ]
    for i, (n, label, formula, ref, calc) in enumerate(others_steps):
        y = 1.7 + i * 1.25
        add_rounded(s, 0.4, y, 12.55, 1.15, WHITE, line=RED, line_width=0.5, radius=0.03)
        add_rounded(s, 0.4, y + 0.05, 0.65, 1.05, RED, radius=0.5)
        add_text(s, 0.4, y + 0.05, 0.65, 1.05, n, size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, 1.2, y + 0.05, 7, 0.3, label, size=12, bold=True, color=INK)
        add_text(s, 11.4, y + 0.05, 1.5, 0.3, ref, size=9, color=MUTED, font='Consolas', align=PP_ALIGN.RIGHT)
        add_text(s, 1.2, y + 0.35, 11.5, 0.25, f'公式: {formula}', size=9, color=MUTED, font='Consolas')
        add_text(s, 1.2, y + 0.62, 11.5, 0.5, f'↳ {calc}', size=9, color=PURPLE, font='Consolas')

    footer(s, page_n, total_pages)


def slide_step5_zone_F_total(total_pages, page_n):
    s = add_slide()
    header(s, '區 F · Total Roll-up · 4 步驟', 'compute_total · MVA + SGA + Profit → TC · 全廠通用 · 對應 C125', 5)

    total_steps = [
        ('1', 'MVA Σ (9 製程 × 9 component)',
            'Σ cs_run_cell.cost_per_unit_usd', 'Cleansheet!C118',
            '9×9 matrix SUM = $1.8558 / unit (CN 案 baseline)'),
        ('2', 'SG&A',
            'motherboard × baseline.sga_pct', 'Cleansheet!J121',
            '$8.683 × 2.0% = $0.1737 / unit'),
        ('3', 'Profit',
            '(MVA + motherboard) × baseline.profit_pct', 'Cleansheet!J121',
            '($1.8558 + $8.683) × 14.0% = $1.4754 / unit'),
        ('4', 'Total Cost / unit (TC) FINAL',
            'Material + MVA + SGA + Profit', 'Cleansheet!C125',
            '$8.683 + $1.8558 + $0.1737 + $1.4754 = $12.1879 / unit'),
    ]
    for i, (n, label, formula, ref, calc) in enumerate(total_steps):
        y = 1.7 + i * 1.25
        add_rounded(s, 0.4, y, 12.55, 1.15, WHITE, line=RED, line_width=0.5, radius=0.03)
        add_rounded(s, 0.4, y + 0.05, 0.65, 1.05, RED, radius=0.5)
        add_text(s, 0.4, y + 0.05, 0.65, 1.05, n, size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, 1.2, y + 0.05, 7, 0.3, label, size=12, bold=True, color=INK)
        add_text(s, 11.4, y + 0.05, 1.5, 0.3, ref, size=9, color=MUTED, font='Consolas', align=PP_ALIGN.RIGHT)
        add_text(s, 1.2, y + 0.35, 11.5, 0.25, f'公式: {formula}', size=9, color=MUTED, font='Consolas')
        add_text(s, 1.2, y + 0.62, 11.5, 0.5, f'↳ {calc}', size=9, color=PURPLE, font='Consolas')

    footer(s, page_n, total_pages)


def slide_step5_ui(total_pages, page_n):
    s = add_slide()
    header(s, 'Compute Trace · UI 操作', 'zone tabs 切換 + ▶ 下一步 動畫 + 底部全 5 區 final 並列', 5)

    add_text(s, 0.4, 1.3, 12, 0.3, '🎯 完整操作流程 demo',
             size=12, bold=True, color=INK)

    mockup = """┌─ ⚙️ Step 5 · Compute Trace — 5 區公式 (SMT 主板貼裝) ───────────────────────┐
│                                                                              │
│  ┌──────────────────── zone tabs(6 區並列) ─────────────────────────────┐  │
│  │ [👷 區 A · DL 9 步●] [👥 區 B · IDL 4] [🏭 區 C · Equip 6] ...        │  │
│  │ [🏗️ 區 D · Fac 6]   [📦 區 E · Others 4] [💰 區 F · Total 4]          │  │
│  └─────────────────────────────────────────────────────────────────────────┘ │
│                                                                              │
│  ┌─ ▶ 區 A · A. Direct Labor ─────────────────────────────────────────────┐ │
│  │ 公式:DL_wage × SEA × Total_DL × mult1 × mult2 / weekly_output  3/9 步 │ │
│  │            [⟲ Reset]  [▶ 下一步]  [⏭ 全部跑完]                          │ │
│  └─────────────────────────────────────────────────────────────────────────┘ │
│                                                                              │
│  ┌─ 步驟逐步揭露區(紅底) ────────────────────────────────────────────────┐ │
│  │ ① UPH         3600/8.64×98%×95% = 388.06                              │ │
│  │ ② Max output  153.41 × 388.06 = 59,531        ← 已揭露                │ │
│  │ ③ Demand      418k/50×1.2 = 10,032            ← 已揭露(current)      │ │
│  │ ④ Lines       (尚未揭露,淡灰)                                          │ │
│  │ ⑤-⑨...        (尚未揭露)                                                │ │
│  └─────────────────────────────────────────────────────────────────────────┘ │
│                                                                              │
│  📊 5 區公式狀態(全廠)                                                    │
│  ┌── 區 A ── 區 B ── 區 C ── 區 D ── 區 E ── 區 F ──┐                     │
│  │ $0.164  $0.017  $0.104  $0.009  $0.095  $12.188 │  ← 跑完每區同步     │
│  └──────────────────────────────────────────────────┘                     │
└──────────────────────────────────────────────────────────────────────────────┘"""
    add_rounded(s, 0.4, 1.7, 12.55, 5.25, WHITE, line=RED, line_width=1.5, radius=0.02)
    add_text(s, 0.55, 1.85, 12.3, 5.0, mockup, size=8, color=TEXT, font='Consolas')

    footer(s, page_n, total_pages)


def slide_step5_schema(total_pages, page_n):
    s = add_slide()
    header(s, 'Compute · 寫入 DB', 'bom_cs_run + bom_cs_run_cell · 一案多 run · 一 run 多 cell', 5)

    add_rounded(s, 0.4, 1.3, 12.55, 5.4, BG_SOFT, line=RED, line_width=1.5, radius=0.03)
    add_rounded(s, 0.4, 1.3, 12.55, 0.4, RED, radius=0.03)
    add_text(s, 0.55, 1.35, 12.3, 0.3, '🗄️ bom_cs_run(snapshot 主表)+ bom_cs_run_cell(matrix 細胞)', size=12, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)

    schema = """-- (1) Run snapshot · 一案一廠多 run(可多版本 reprice)
CREATE TABLE bom_cs_run (
  run_id                  NUMBER GENERATED ALWAYS AS IDENTITY PRIMARY KEY,
  case_factory_id         NUMBER REFERENCES bom_cs_case_factory(case_factory_id) ON DELETE CASCADE,
  run_label               VARCHAR2(80),                  -- 'v1' / '客戶議價 #2'
  status                  VARCHAR2(20) DEFAULT 'computing', -- computing/ready/archived/failed
  compute_engine          VARCHAR2(20) DEFAULT 'db_v1',
  computed_at             TIMESTAMP,
  computed_by             NUMBER REFERENCES users(id),
  motherboard_cost_usd    NUMBER(15,6),                  -- from BOM main board
  errors_json             CLOB,
  warnings_json           CLOB,
  raw_inputs_json         CLOB,                          -- 整份 input snapshot 備查
  created_at              TIMESTAMP DEFAULT SYSTIMESTAMP
);

-- (2) cs_run_cell · 9 製程 × 9 component matrix · 多 qty_scenario 寫多份
CREATE TABLE bom_cs_run_cell (
  run_id                  NUMBER REFERENCES bom_cs_run(run_id) ON DELETE CASCADE,
  qty_scenario_code       VARCHAR2(40) DEFAULT 'BASELINE',
  process_code            VARCHAR2(40),                  -- 'SMT_MAIN' / 'COMMON'
  component_code          VARCHAR2(40),                  -- 'DL_CPU' / 'IDL_CPU' / 'EQUIP_MRO' / ...
  cost_per_unit_usd       NUMBER(15,6),                  -- 每格的 $ / unit
  formula_text            VARCHAR2(500),                 -- '(C56+C55)/C24'
  source_cell_ref         VARCHAR2(80),                  -- 'Cleansheet!C58'
  intermediate_json       CLOB,                          -- 步驟 1-9 中間值備查
  PRIMARY KEY (run_id, qty_scenario_code, process_code, component_code)
);

-- Step 5 寫入範例(區 A SMT_MAIN DL):
INSERT INTO bom_cs_run_cell VALUES (
  run_id=1, qty_scenario_code='HIGH', process_code='SMT_MAIN',
  component_code='DL_CPU', cost_per_unit_usd=0.1640,
  formula_text='(DL_wk + IDL_line_dep_wk) / weekly_output',
  source_cell_ref='Cleansheet!C58',
  intermediate_json='{"uph":388.06,"max_output":59531,"lines":1,"dl_wk":9116,"idl_line_dep":645}'
);"""
    add_text(s, 0.6, 1.8, 12.1, 5.0, schema, size=8, color=TEXT, font='Consolas')

    footer(s, page_n, total_pages)


# ════════════════════════════════════════════════════════════════════════
# STEP 6 — 結果矩陣(4 slides)
# ════════════════════════════════════════════════════════════════════════

def slide_step6_matrix(total_pages, page_n):
    s = add_slide()
    header(s, '結果矩陣 · 9×9 view', '對應 Excel Price Summary table · 一眼看 MVA 構成', 6, '結果矩陣')

    add_text(s, 0.4, 1.3, 12, 0.3, '📊 Step 6 panel · 9 cost component × 9 process matrix · MVA 整體拆解',
             size=12, bold=True, color=INK)

    matrix = """┌────────────────────────────────────────────────────────────────────────────────────┐
│ Cost Component × Process Matrix · 9 × 9 · MATRIX                                  │
├────────────────────────────────────────────────────────────────────────────────────┤
│ Cost Component         SMT    Wave   Router Laser  BB-Assy BB-Test Mat   Q-SMT Q-BB COMMON │
│ ───────────────────────────────────────────────────────────────────────────────── │
│ DL_CPU                0.164  0.349  0.082  0.025  0.362  0.131  0.082  0.020 0.029 —     │
│ IDL_CPU               0.017  0.005  0.003  0.002  0.028  0.012  0.006  0.001 0.002 —     │
│ EQUIP_MRO             0.005  0.001  0.001  0.001  0.005  0.001  0.001  0.001 0.001 —     │
│ EQUIP_DEPR            0.099  0.025  0.009  0.012  0.005  0.005  0.001  0.001 0.001 —     │
│ IND_MAT               0.363  0.082  0.006  0.002  0.018  0.001  0.020  0.002 —     —     │
│ FACILITY              0.008  0.003  0.001  0.001  0.027  0.011  0.039  0.005 0.008 —     │
│ FREIGHT (Common only) —      —      —      —      —      —      —      —     —     0.006 │
│ VAT (Common only)     —      —      —      —      —      —      —      —     —     0.000 │
│ LOSS (Common only)    —      —      —      —      —      —      —      —     —     0.007 │
│ ───────────────────────────────────────────────────────────────────────────────── │
│ MVA Subtotal          0.656  0.465  0.102  0.043  0.445  0.161  0.149  0.030 0.041 0.013 │ → $1.8558 (Total)
└────────────────────────────────────────────────────────────────────────────────────┘

📊 Total MVA = $1.8558 / unit  ✓ 對齊 Excel Price Summary (ε < 0.01)"""
    add_rounded(s, 0.4, 1.7, 12.55, 4.3, WHITE, line=NAVY, line_width=1.5, radius=0.02)
    add_text(s, 0.55, 1.85, 12.3, 4.1, matrix, size=7.5, color=TEXT, font='Consolas')

    # 判讀提示
    add_rounded(s, 0.4, 6.15, 12.55, 0.8, RGBColor(0xE0, 0xE7, 0xFF), line=NAVY, radius=0.04)
    add_text(s, 0.55, 6.2, 12.3, 0.75, """💡 判讀技巧:
  • SMT 列 IND_MAT = $0.363 最大 → 錫膏耗材成本高 · 可砍?
  • BB-Assy DL_CPU = $0.362 最大 → 人工密集 · 可自動化?
  • Common 欄(VAT 0.000)只在 Common only 行有值 · 不分製程""", size=9, color=TEXT, anchor=MSO_ANCHOR.TOP)

    footer(s, page_n, total_pages)


def slide_step6_smt_assy(total_pages, page_n):
    s = add_slide()
    header(s, '結果 · SMT / Assy roll-up', '兩大群組 MVA 分布 · 哪邊吃 cost 大?', 6)

    add_text(s, 0.4, 1.3, 12, 0.3, '📊 MVA $1.8558 = SMT $1.279 (68.9%) + Assy + Common $0.577 (31.1%)',
             size=12, bold=True, color=INK)

    # 左:SMT
    add_rounded(s, 0.4, 1.7, 6.2, 5.0, RGBColor(0xEE, 0xF2, 0xFF), line=RGBColor(0x63,0x66,0xF1), line_width=2, radius=0.03)
    add_rounded(s, 0.4, 1.7, 6.2, 0.45, RGBColor(0x63,0x66,0xF1), radius=0.03)
    add_text(s, 0.55, 1.75, 5.9, 0.35, '📦 SMT-side MVA  $1.2791  ·  68.9% of MVA', size=13, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    smt_list = """SMT_MAIN          $0.656 · 35.4%  ███████████████████
WAVE_SOLDER       $0.465 · 25.0%  █████████████
ROUTER_OFFLINE    $0.102 ·  5.5%  ███
LASER_ETCH        $0.043 ·  2.3%  █
Q_SMT             $0.013 ·  0.7%

→ SMT_MAIN + WAVE_SOLDER = 60.4% (重押)
→ ROUTER + LASER + Q-SMT = 8.5% (相對輕)
→ EPM 改善焦點:SMT 主板貼裝 IND_MAT $0.363
   (錫膏耗材成本最大 · 可議錫膏單價或換 SAC)"""
    add_text(s, 0.6, 2.3, 5.9, 4.3, smt_list, size=9.5, color=TEXT, font='Consolas')

    # 右:Assy + Common
    add_rounded(s, 6.8, 1.7, 6.15, 5.0, RGBColor(0xF0, 0xFD, 0xF4), line=GREEN, line_width=2, radius=0.03)
    add_rounded(s, 6.8, 1.7, 6.15, 0.45, GREEN, radius=0.03)
    add_text(s, 6.95, 1.75, 5.85, 0.35, '🔧 Assy + Common MVA  $0.5767  ·  31.1% of MVA', size=13, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    assy_list = """BB_ASSY           $0.445 · 24.0%  ████████████
BB_TEST           $0.161 ·  8.7%  █████
MAT_MGMT          $0.149 ·  8.0%  ████
Q_BB              $0.041 ·  2.2%  █
COMMON            $0.013 ·  0.7%   (Freight + VAT + Loss)

→ BB_ASSY 一支獨秀 24% · 主因 DL_CPU $0.362
   (BB Assy 22 DL/shift × 1 shift = 22 人/日)
→ EPM 改善焦點:
   ✓ 自動鎖螺絲機 → 砍 5-10 DL
   ✓ AGV / Conveyor → 砍 Material Mgmt
   ✓ 客戶 buy-in:壓減 Test 站數
→ BB Test 8.7% 不算高(其實 standalone test 不貴)"""
    add_text(s, 6.95, 2.3, 5.85, 4.3, assy_list, size=9.5, color=TEXT, font='Consolas')

    footer(s, page_n, total_pages)


def slide_step6_waterfall(total_pages, page_n):
    s = add_slide()
    header(s, '結果 · Total Cost waterfall', 'Material + MVA + SGA + Profit → TC · per unit · CN baseline', 6)

    add_text(s, 0.4, 1.3, 12, 0.3, '💰 Total Cost per Unit (Ex-Factory · 計算明細)',
             size=12, bold=True, color=INK)

    # waterfall 表
    add_rounded(s, 0.4, 1.7, 12.55, 4.0, WHITE, line=OCEAN, line_width=2, radius=0.03)

    items = [
        ('Motherboard Material (from BOM)', '$8.6830', '✓ Material 從 BOM 採購策略總和帶入', INK),
        ('+ MVA (9 製程合計)',                '$1.8558', '✓ 從 Step 5 區 A-E 寫入 cs_run_cell · Σ',  OCEAN),
        ('+ SG&A (2.0% × motherboard)',       '$0.1737', '✓ 從 baseline.sga_pct',                      INK),
        ('+ Profit (14.0% × (MVA+motherboard))','$1.4754', '✓ 從 baseline.profit_pct',                  INK),
        ('═══ Transformation Cost (TC) ═══',  '$12.1879', '🎯 客戶報價基底 (CN 案 baseline)',         OCEAN),
    ]
    for i, (label, val, note, color) in enumerate(items):
        y = 1.9 + i * 0.70
        if i == 4:
            add_rect(s, 0.55, y - 0.05, 12.25, 0.7, RGBColor(0xCF, 0xFA, 0xF1))
            add_text(s, 0.7, y, 7, 0.4, label, size=14, bold=True, color=color, anchor=MSO_ANCHOR.MIDDLE)
            add_text(s, 8.0, y, 4.7, 0.4, val, size=16, bold=True, color=color, font='Consolas', anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT)
        else:
            add_text(s, 0.7, y, 6, 0.4, label, size=12, color=TEXT, anchor=MSO_ANCHOR.MIDDLE)
            add_text(s, 6.8, y, 2.5, 0.4, val, size=13, bold=True, color=color, font='Consolas', anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT)
            add_text(s, 9.5, y, 3.3, 0.4, note, size=10, color=MUTED, anchor=MSO_ANCHOR.MIDDLE)

    # 提示
    add_rounded(s, 0.4, 5.85, 12.55, 1.1, RGBColor(0xFE, 0xF3, 0xC7), line=GOLD, radius=0.04)
    add_text(s, 0.55, 5.92, 12.3, 1.0, """💡 注意:
  • TC ≠ 客戶報價(報價 = TC × markup% by mfg + tier · v0.7 雙價結構)
  • 改 Qty scenario(High → Low):lines util 下降 → MVA × 1.41 → TC 跳 $13.5(看 §factory_matrix)
  • CN $12.19 vs VN $11.76(wage 便宜)vs TW $13.91(自動化高但 labor 貴)""",
             size=9, color=TEXT, anchor=MSO_ANCHOR.TOP)

    footer(s, page_n, total_pages)


def slide_step6_lock(total_pages, page_n):
    s = add_slide()
    header(s, '結果 · Lock & Propagate', '滿意 → DPM Lock → 寫 cs_run_result fact table · 進 §factory_matrix', 6)

    add_text(s, 0.4, 1.3, 12, 0.3, '🔒 DPM 滿意計算結果 → 按「🔒 Lock」進報價階段',
             size=12, bold=True, color=INK)

    # 4 step
    lock_steps = [
        ('1', '🔍 DPM review §Cleansheet 結果',
            '看 9×9 matrix · SMT/Assy roll-up · TC waterfall · 對齊預期'),
        ('2', '🔒 DPM 按 Step 6 區頂部「Lock」按鈕',
            'cs_run.status = ready → 開 BOM lock 預檢查(料齊 / cs_run 齊 / 雙價齊)'),
        ('3', '⚡ 系統自動 propagate',
            '寫 cs_run_result × 24 列(3 factor × 2 variant × 2 qty × 2 pkg)\n'
            'Material × MVA × SGA × Profit 全部寫入 fact table'),
        ('4', '📈 進 §factory_matrix 看 24 cells pivot',
            'BPM / 採購主管 / finance 切 qty/pkg/variant toggle 看 6 cells 主表\n'
            'Margin Analysis heatmap 顯示 Top Markup Items / per PKG'),
    ]
    for i, (n, title, desc) in enumerate(lock_steps):
        y = 1.7 + i * 1.25
        add_rounded(s, 0.4, y, 12.55, 1.15, BG_SOFT, line=NAVY, line_width=1, radius=0.04)
        add_rounded(s, 0.4, y + 0.05, 0.75, 1.05, NAVY, radius=0.5)
        add_text(s, 0.4, y + 0.05, 0.75, 1.05, n, size=24, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, 1.3, y + 0.1, 11, 0.35, title, size=13, bold=True, color=INK)
        add_text(s, 1.3, y + 0.5, 11.5, 0.7, desc, size=10, color=TEXT)

    footer(s, page_n, total_pages)


# ════════════════════════════════════════════════════════════════════════
# CLOSING — Checklist + Timeline + Q&A
# ════════════════════════════════════════════════════════════════════════

def slide_checklist(total_pages, page_n):
    s = add_slide()
    header(s, '速查 Checklist — 6 Step 一頁帶走', '備忘:跑完一個案的 §Cleansheet 全流程 · 6 Step × 細項')

    # 大表
    items = [
        (1, '製程輸入', [
            ('☐ 9 製程 tab 依序進'),
            ('☐ 6 群組(Staffed/Throughput/Volume/DL/IDL/SEA)'),
            ('☐ 改 takt/yield 即時看 DL 動'),
            ('☐ 機密欄 (DL/IDL config) 確認'),
        ], SKY),
        (2, 'IDL Matrix', [
            ('☐ 17 角色 × 9 製程 multiplier'),
            ('☐ 從 Excel J64-Z83 seed'),
            ('☐ Row Σ 檢查 ≤ 1.0'),
            ('☐ Col Σ 檢查跨製程合理'),
        ], PURPLE),
        (3, '設備 binding', [
            ('☐ 廠 master view 看 58 件'),
            ('☐ 案級 view 挑 28 件'),
            ('☐ line_qty / qty_per_line 填'),
            ('☐ shared 設備(<1)記得標'),
        ], CYAN),
        (4, '耗材 binding', [
            ('☐ 17 件耗材 annual_qty'),
            ('☐ unit_cost override (議價時)'),
            ('☐ 年總成本 Σ 合理性'),
            ('☐ proc 分類正確'),
        ], GOLD),
        (5, 'Compute Trace', [
            ('☐ 區 A DL 9 步'),
            ('☐ 區 B IDL 4 步'),
            ('☐ 區 C Equip 6 步'),
            ('☐ 區 D Fac 6 步 / E Others 4 / F Total 4'),
        ], RED),
        (6, '結果 + Lock', [
            ('☐ 9×9 matrix 對齊預期'),
            ('☐ SMT/Assy roll-up % 合理'),
            ('☐ TC waterfall 看 Total'),
            ('☐ DPM Lock + Propagate'),
        ], NAVY),
    ]
    for i, (n, title, lines, c) in enumerate(items):
        col = i % 3
        row = i // 3
        x = 0.4 + col * 4.25
        y = 1.45 + row * 2.7
        add_rounded(s, x, y, 4.1, 2.55, BG_SOFT, line=c, line_width=1.5, radius=0.03)
        add_rounded(s, x, y, 4.1, 0.42, c, radius=0.03)
        add_text(s, x + 0.15, y + 0.05, 3.8, 0.32, f'Step {n} · {title}', size=12, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
        for j, line in enumerate(lines):
            add_text(s, x + 0.2, y + 0.55 + j * 0.45, 3.8, 0.4, line, size=10, color=TEXT, anchor=MSO_ANCHOR.TOP)

    footer(s, page_n, total_pages)


def slide_pitfalls(total_pages, page_n):
    s = add_slide()
    header(s, '常見問題 · 注意事項', 'EPM / DPM 操作 §Cleansheet 時的雷區')

    pitfalls = [
        ('1', '🚨 改了 Step 1 製程輸入,矩陣不動?',
            'Step 1 可即時 preview DL_CPU 變化,但 Step 6 矩陣其他 component(IDL/Equip/Fac/Others)不會即時動 · 需要按 Step 5 「Compute」按鈕重跑各區公式才會完全 sync · 否則只有 DL_CPU 即時'),
        ('2', '⚠️ Row Σ > 1.0 在 IDL Matrix?',
            'OPS_MGR 總和 0.523 < 1.0 → 合理 · 若某角色 Row Σ > 1.0 = 該角色超時 · 公司沒這人 · 改 multiplier 或 split 角色'),
        ('3', '⚠️ 設備 acq cost 跨案重複計算?',
            'DEK 1 台跨多案共用 → 各案 line_qty 加總應 = 1.0 · 系統不會幫你檢查 · DPM 跨案查 utilization 看' ),
        ('4', '⚠️ qty scenario LOW MVA 變高 +40%?',
            '正常 · 低量 → lines util ↓ → DL fixed cost 攤少量 → /unit 高 · v0.7 已 demo · 業務報低量單時要警告'),
        ('5', '🚨 Step 5 區 A-D 看到負值 / 0?',
            'csComputeXxxSteps 用 case_process 輸入 · 若 takt = 0 或 yield = 0 會 division by zero · UI 應 validate · 沒 validate 就會 NaN'),
        ('6', '⚠️ Compute 完跑 regression test fail?',
            'cs_run.status=failed · warnings_json 記每 cell vs Excel 差距 · ε > 0.01 表示某公式抄錯 · EPM 要回去 trace 哪個 component 不對'),
    ]
    for i, (n, title, desc) in enumerate(pitfalls):
        col = i % 2
        row = i // 2
        x = 0.4 + col * 6.27
        y = 1.45 + row * 1.85
        add_rounded(s, x, y, 6.2, 1.75, BG_SOFT, line=GOLD, line_width=1, radius=0.04)
        add_rounded(s, x, y + 0.05, 0.45, 1.65, GOLD, radius=0.5)
        add_text(s, x, y + 0.05, 0.45, 1.65, n, size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + 0.55, y + 0.05, 5.6, 0.3, title, size=10.5, bold=True, color=INK, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + 0.55, y + 0.45, 5.6, 1.25, desc, size=9, color=TEXT)

    footer(s, page_n, total_pages)


def slide_close(total_pages, page_n):
    s = add_slide()
    add_rect(s, 0, 0, 13.333, 7.5, INK)
    add_rect(s, 0, 2.8, 13.333, 0.06, OCEAN)
    add_text(s, 0.7, 1.2, 12, 0.6, '完。', size=44, bold=True, color=WHITE)
    add_text(s, 0.7, 1.9, 12, 0.4, 'Q&A · Internal Operation Discussion', size=18, color=OCEAN)

    add_rounded(s, 0.7, 3.4, 12.0, 3.5, RGBColor(0x14, 0x2F, 0x4E), radius=0.04)
    add_text(s, 1.0, 3.6, 11.5, 0.4, '🔗 相關資源', size=16, bold=True, color=OCEAN)
    add_text(s, 1.0, 4.1, 11.5, 2.7, """• 互動 demo:Cortex_互動Demo_v0.10.html(§Cleansheet 6-Step 完整操作)
• 上層 v1 手冊:Cortex_MVA操作流程說明手冊_v1.pptx(Phase A-G 端到端宏觀)
• 董事長簡報:Cortex_BOM_MVA_董事長簡報_v1.pptx(單頁 4:3 ROI 重點)
• 設計文件:cleansheet-mva-sd v0.3(14 表 schema · 5 Layer · 7 Phase · 100+ 公式對應)
• 兩案範本:
    - Rival 3+ Wired Mouse Cleansheet_China-20241011-19S.xlsx (本手冊主例)
    - WHOOP_Gen4 MP quotation book_ for AI.xlsx (FATP 結構不同)
• 開發 Roadmap:Phase 3-4 重點是讓 EPM 不再回 Excel · 全在系統 UI 完成 §Cleansheet""", size=12, color=RGBColor(0xCB, 0xD5, 0xE1))

    add_text(s, 0.7, 7.05, 12, 0.3, 'Cortex MVA 操作詳細手冊 v2 · 2026-06 · FOXLINK GPT',
             size=10, color=MUTED, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════
# RUN
# ════════════════════════════════════════════════════════════════════════
total = 0
total += 1  # cover
total += 1  # toc
total += 1  # prereq
# Step 1: 5
# Step 2: 4
# Step 3: 5
# Step 4: 3
# Step 5: 10 (overview + 6 zones + UI + schema + intro is overview already)
# wait recount:
# step5 — overview (1) + zone A (1) + zone B (1) + zone C (1) + zone D (1) + zone E (1) + zone F (1) + UI (1) + schema (1) = 9
# Step 6: 4
total += (5 + 4 + 5 + 3 + 9 + 4)  # = 30
total += 3  # checklist + pitfalls + close
# = 36 slides

print(f'Building {total} slides...')

slide_cover()
slide_toc(total)
slide_prereq(total)

page = 4
# Step 1 — 5 slides
slide_step1_intro(total, page);        page += 1
slide_step1_ui_entry(total, page);     page += 1
slide_step1_groups(total, page);       page += 1
slide_step1_interactive(total, page);  page += 1
slide_step1_schema(total, page);       page += 1

# Step 2 — 4 slides
slide_step2_intro(total, page);   page += 1
slide_step2_ui(total, page);      page += 1
slide_step2_init(total, page);    page += 1
slide_step2_schema(total, page);  page += 1

# Step 3 — 5 slides
slide_step3_intro(total, page);         page += 1
slide_step3_factory_view(total, page);  page += 1
slide_step3_case_view(total, page);     page += 1
slide_step3_picking(total, page);       page += 1
slide_step3_schema(total, page);        page += 1

# Step 4 — 3 slides
slide_step4_intro(total, page);       page += 1
slide_step4_operations(total, page);  page += 1
slide_step4_schema(total, page);      page += 1

# Step 5 — 9 slides
slide_step5_overview(total, page);       page += 1
slide_step5_zone_A_dl(total, page);      page += 1
slide_step5_zone_B_idl(total, page);     page += 1
slide_step5_zone_C_equip(total, page);   page += 1
slide_step5_zone_D_facility(total, page);page += 1
slide_step5_zone_E_others(total, page);  page += 1
slide_step5_zone_F_total(total, page);   page += 1
slide_step5_ui(total, page);             page += 1
slide_step5_schema(total, page);         page += 1

# Step 6 — 4 slides
slide_step6_matrix(total, page);    page += 1
slide_step6_smt_assy(total, page);  page += 1
slide_step6_waterfall(total, page); page += 1
slide_step6_lock(total, page);      page += 1

# Closing
slide_checklist(total, page);  page += 1
slide_pitfalls(total, page);   page += 1
slide_close(total, page);      page += 1

OUT = 'd:/vibe_coding/foxlink_gpt/docs/Cortex_MVA操作流程說明手冊_v2_詳細版.pptx'
prs.save(OUT)
print(f'[OK] Saved: {OUT}')
print(f'     Total slides: {len(prs.slides)}')
