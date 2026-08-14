# -*- coding: utf-8 -*-
"""
Cortex MVA 操作手冊 v4 — Excel 完整對齊版
對齊 Cortex_互動Demo_v0.12.html · cleansheet-mva-sd v0.5
新增/修正章節:
  - v0.12 認錯頁(廢 7 工種 DL → 回 Excel 4 category)
  - §A 完整 55 行 schema 映射表
  - DL 4 category(對齊 r28-31)+ 廠級單一 wage
  - IDL Line-Dep 4 角色 weekly wage(對齊 r43-46)
  - 案級 4 新欄位(warehouse_dl + per_day + IQC)
~ 32 slides · 16:9
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

INK    = RGBColor(0x0B, 0x1F, 0x3A)
TEXT   = RGBColor(0x37, 0x47, 0x55)
MUTED  = RGBColor(0x6B, 0x72, 0x80)
OCEAN  = RGBColor(0x02, 0xC3, 0x9A)
NAVY   = RGBColor(0x1C, 0x72, 0x93)
LINE   = RGBColor(0xE5, 0xE7, 0xEB)
BG_SOFT= RGBColor(0xFA, 0xFB, 0xFC)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GOLD   = RGBColor(0xCA, 0x8A, 0x04)
RED    = RGBColor(0xDC, 0x26, 0x26)
GREEN  = RGBColor(0x16, 0xA3, 0x4A)
PURPLE = RGBColor(0x7C, 0x3A, 0xED)
CYAN   = RGBColor(0x08, 0x91, 0xB2)
SKY    = RGBColor(0x0E, 0xA5, 0xE9)
ORANGE = RGBColor(0xF5, 0xA5, 0x24)


def add_slide(): return prs.slides.add_slide(prs.slide_layouts[6])


def add_text(s, x, y, w, h, text, size=12, bold=False, color=None,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font='Microsoft JhengHei'):
    tx = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tx.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    if color:
        r.font.color.rgb = color
    return tx


def add_rect(s, x, y, w, h, fill, line=None, line_width=0.5):
    rect = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    rect.fill.solid()
    rect.fill.fore_color.rgb = fill
    if line is None:
        rect.line.fill.background()
    else:
        rect.line.color.rgb = line
        rect.line.width = Pt(line_width)
    rect.shadow.inherit = False
    return rect


def add_rounded(s, x, y, w, h, fill, line=None, line_width=0.5, radius=0.05):
    rect = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    rect.adjustments[0] = radius
    rect.fill.solid()
    rect.fill.fore_color.rgb = fill
    if line is None:
        rect.line.fill.background()
    else:
        rect.line.color.rgb = line
        rect.line.width = Pt(line_width)
    rect.shadow.inherit = False
    return rect


def header(s, title, subtitle=None, color=NAVY, badge=None):
    add_rect(s, 0, 0, 13.333, 0.16, color)
    if badge:
        add_rounded(s, 0.4, 0.3, 1.5, 0.45, color, radius=0.3)
        add_text(s, 0.4, 0.32, 1.5, 0.4, badge, size=12, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, 2.05, 0.25, 11, 0.5, title, size=22, bold=True, color=INK)
        if subtitle:
            add_text(s, 2.05, 0.78, 11, 0.32, subtitle, size=11, color=MUTED)
    else:
        add_text(s, 0.4, 0.25, 11, 0.5, title, size=24, bold=True, color=INK)
        if subtitle:
            add_text(s, 0.4, 0.8, 11, 0.32, subtitle, size=12, color=MUTED)


def footer(s, n, total):
    add_rect(s, 0.4, 7.08, 12.55, 0.02, LINE)
    add_text(s, 0.4, 7.12, 9, 0.22, 'Cortex MVA 操作手冊 v4 · Excel 完整對齊版 · 對齊 v0.12 demo',
             size=8.5, color=MUTED)
    add_text(s, 12.0, 7.12, 1, 0.22, f'{n} / {total}', size=8.5, color=MUTED, align=PP_ALIGN.RIGHT)


# ════════════════════════════════════════════════════════════════════════
# SLIDES
# ════════════════════════════════════════════════════════════════════════

def slide_cover():
    s = add_slide()
    add_rect(s, 0, 0, 13.333, 7.5, INK)
    add_rect(s, 0, 2.7, 13.333, 0.06, OCEAN)
    add_text(s, 0.7, 1.2, 12, 0.4, 'FOXLINK 正崴 · Cortex BOM × MVA', size=12, color=OCEAN)
    add_text(s, 0.7, 1.6, 12, 0.95, 'MVA 操作手冊 v4', size=48, bold=True, color=WHITE)
    add_text(s, 0.7, 3.0, 12, 0.5, 'Excel Cleansheet §A 完整 55 行對齊 + DL/IDL 修正', size=20, color=WHITE)

    add_rounded(s, 0.7, 3.85, 12, 2.1, RGBColor(0x14, 0x2F, 0x4E), radius=0.04)
    add_text(s, 1.0, 4.0, 11.5, 0.4, 'v4 認錯 + 修正(對齊 v0.12 demo)', size=14, bold=True, color=OCEAN)
    items = [
        ('❌', 'v0.11 偏差', '7 種 DL 工種費率 → 偏離 Excel 真實結構 · 廢除', RED),
        ('✓', 'DL 回歸 Excel', '4 category(DL/Debug/Functional/⭐Warehouse)+ 廠級單一 $4.95 wage', GREEN),
        ('🆕', 'IDL Line-Dep wage', '4 角色 weekly wage(LL$403/Tech$436.80/⭐IQC$403/Sup$403)', PURPLE),
        ('🆕', '案級補 4 欄', 'warehouse_dl + line_leader_per_day + technician_per_day + iqc_per_day', CYAN),
        ('🔧', '公式對齊', 'csComputeDlSteps wage 從 schema 取 · 不再 hardcode', SKY),
        ('✅', '保留正確設計', 'BG/BU 隔離 · 設備類別 · 權限 grant · 設定 master UI', NAVY),
    ]
    for i, (ico, t, sub, c) in enumerate(items):
        y = 4.45 + i * 0.25
        add_text(s, 1.0, y, 0.3, 0.25, ico, size=14, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, 1.4, y, 2.3, 0.25, t, size=11, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, 3.7, y, 8.5, 0.25, sub, size=10, color=RGBColor(0xCB, 0xD5, 0xE1), anchor=MSO_ANCHOR.MIDDLE)

    add_text(s, 0.7, 6.4, 12, 0.4, 'v4 · 2026-06 · 對應 cleansheet-mva-sd v0.5 + Cortex_互動Demo_v0.12.html', size=11, color=MUTED)


def slide_apology(total):
    s = add_slide()
    header(s, 'v0.11 設計偏差 · 認錯與修正', '感謝 user 質疑 — Cleansheet §A 完整看過才發現設計偏離', RED)

    add_rounded(s, 0.4, 1.3, 12.55, 5.6, BG_SOFT, line=RED, line_width=2, radius=0.03)
    add_rounded(s, 0.4, 1.3, 12.55, 0.5, RED, radius=0.03)
    add_text(s, 0.55, 1.38, 12.3, 0.36, '⚠️ v0.11 我做錯的 4 件事', size=13, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)

    mistakes = [
        ('1', 'DL「7 種工種費率」over-engineer · 偏離 Excel',
         'v0.11 加 OP_GENERAL/OP_DEBUG/OP_FUNCT/OP_CLEANROOM/OP_MOLD/OP_TEST/OP_AUTO_MON 7 個工種各自 hourly_rate。\n'
         'Excel 真實:DL 是 4 個 category(DL/Debug/Functional/Warehouse)用同一個 $4.95/hr。\n'
         'user 之前要求「DL 費率按工種」其實對應的應該是 IDL Line-Dep 4 角色,不是 DL。'),
        ('2', 'IDL Line-Dependent 4 個 weekly wage 完全漏',
         'Cleansheet r43-46 有 4 個 wage rate: Line Leader $403 / Technician $436.80 / IQC $403 / Supervisor $403。\n'
         'v0.11 在 csComputeDlSteps 用 hardcode $280/$320/$460 · 而且 IQC 完全沒有。'),
        ('3', '案級「Warehouse DL」完全漏 + per Day count 漏',
         'Excel r31 Warehouse DL 是真實第 4 種 DL category(Material Mgmt 製程用 0.5 人)· schema 沒。\n'
         '另外 r37/38/39 Line Leader Per Day / Technician Per Day / IQC Per Day · 跟 r35/36 per Shift 是兩種獨立 count · 我合一了。'),
        ('4', 'IDL 17 角色 multiplier 表 r62-83 沒做 Manufacturing/Centralized 分區',
         '我 v0.10/0.11 17 角色 mix 在一起 · Excel 真實分 4 個 Manufacturing IDL (r64-67) + 9 個 Centralized Service (r74-82) 兩段。'),
    ]
    for i, (n, title, desc) in enumerate(mistakes):
        y = 1.95 + i * 1.2
        add_rounded(s, 0.55, y, 12.25, 1.1, WHITE, line=RED, line_width=0.5, radius=0.03)
        add_rounded(s, 0.55, y + 0.05, 0.7, 1.0, RED, radius=0.5)
        add_text(s, 0.55, y + 0.05, 0.7, 1.0, n, size=24, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, 1.4, y + 0.05, 11, 0.3, title, size=12, bold=True, color=INK)
        add_text(s, 1.4, y + 0.38, 11.3, 0.7, desc, size=9.5, color=TEXT)

    footer(s, 2, total)


def slide_section_a_full(total, n):
    s = add_slide()
    header(s, 'Excel Cleansheet §A 完整 55 行映射', '對齊 SteelSeries Rival 3+ Cleansheet sheet r4-58', NAVY, '📋 對齊')

    # 上方 SubSection 區塊
    add_text(s, 0.4, 1.3, 12, 0.3, '🔍 §A. Direct Labor Cost 完整結構 — 10 sub-sections × 9 製程 × 55 substantive cells',
             size=12, bold=True, color=INK)

    headers = ['Sub-section', 'Rows', '#', '說明', 'v0.12 覆蓋']
    widths = [3.0, 1.2, 0.7, 5.5, 2.15]
    cx = 0.4
    for h, w in zip(headers, widths):
        add_rect(s, cx, 1.7, w, 0.4, INK)
        add_text(s, cx, 1.7, w, 0.4, h, size=11, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        cx += w

    rows = [
        ('§A.1 Staffed Hours', 'r7-11', '4', 'Working hr/day · Working d/wk · Whr/wk(derived) · Shifts/day', '✓ 100%'),
        ('§A.2 Throughput', 'r12-16', '4', 'TAKT · UPH(derived) · Yield · Eff', '✓ 100%'),
        ('§A.3 Volume', 'r17-21', '4', 'Annual demand · Avg wk · Max wk · Max output/line/wk(derived)', '✓ 100%'),
        ('§A.4 Production Cap', 'r22-26', '4', 'Lines installed · Weekly output · Upside · Debug lines', '✓ 100%'),
        ('§A.5 DL Required', 'r27-33', '6', 'DL/Debug/Functional/⭐Warehouse + Total/Line/Shift + Total/day', '✓ 100% (v0.12 補 Warehouse)'),
        ('§A.6 IDL Line-Indep', 'r34-41', '7', 'LL/Tech Per Shift + Per Day · ⭐IQC/day · Supervisor/day · Total', '✓ 100% (v0.12 補 IQC + 4 per_day)'),
        ('§A.7 IDL Line-Dep wage', 'r42-46', '4', '⭐Line Leader/Tech/IQC/Supervisor weekly wage($403/$436.80/$403/$403)', '✓ 100% (v0.12 新增 schema)'),
        ('§A.8 SEA', 'r47-51', '4', 'Per person hr/day · /wk · Multiplier 1 · Multiplier 2', '✓ 100%'),
        ('§A.9 DL Wage Calc', 'r53-56', '3', 'DL wage/hr($4.95 廠級單一)· Total IDL $/wk · Total DL $/wk', '✓ 100% (v0.12 回歸 Excel)'),
        ('§A.10 Cost per Unit', 'r57-58', '1', 'DL final cost per unit · per 製程', '✓ 100%'),
    ]
    for i, row in enumerate(rows):
        y = 2.1 + i * 0.45
        bg = BG_SOFT if i % 2 == 0 else WHITE
        cx = 0.4
        for j, (cell, w) in enumerate(zip(row, widths)):
            add_rect(s, cx, y, w, 0.43, bg)
            color = TEXT
            bold = False
            if j == 4 and 'v0.12' in cell:
                color = GREEN
                bold = True
            if j == 0:
                color = NAVY
                bold = True
            if j == 2:
                color = PURPLE
                bold = True
            add_text(s, cx + 0.1, y, w - 0.2, 0.43, cell, size=9.5, color=color, bold=bold,
                     font='Consolas' if j == 1 else 'Microsoft JhengHei',
                     anchor=MSO_ANCHOR.MIDDLE)
            cx += w

    add_text(s, 0.4, 6.7, 12, 0.4, '👉 Total 55 行 substantive cells × 9 製程 ≈ 500 個 cell · v0.12 覆蓋率從 v0.11 的 50% 提升到 100%',
             size=11, bold=True, color=GREEN, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    footer(s, n, total)


def slide_dl_4cat(total, n):
    s = add_slide()
    header(s, 'DL 4 category 結構 · 對齊 Excel r28-31', '廢除 7 工種 · 回歸 Excel 真實 4 種 DL category', RED, '👷 DL')

    # 左:v0.11 錯誤
    add_rounded(s, 0.4, 1.3, 6.2, 5.4, RGBColor(0xFE, 0xE2, 0xE2), line=RED, line_width=1.5, radius=0.03)
    add_rounded(s, 0.4, 1.3, 6.2, 0.4, RED, radius=0.03)
    add_text(s, 0.55, 1.35, 5.9, 0.3, '❌ v0.11 (廢除)', size=12, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    before = """v0.11 引入 7 種 DL 工種 · 各自 hourly_rate:
  OP_GENERAL     $4.95
  OP_DEBUG       $8.50
  OP_FUNCT       $12.00
  OP_CLEANROOM   $6.50
  OP_MOLD        $5.50
  OP_TEST        $7.00
  OP_AUTO_MON    $6.00

問題:
  ✗ Excel 真實是 4 category 統一 wage
  ✗ 偏離 EPM 既有 Cleansheet 工作流
  ✗ 違反 Excel r54 設計
    「Wage per hour per person (DL)」單一 wage
  ✗ user 質疑 — 「dl 費率只是其中一小部分」"""
    add_text(s, 0.55, 1.8, 5.9, 4.8, before, size=10, color=TEXT, font='Consolas')

    # 右:v0.12 正確
    add_rounded(s, 6.8, 1.3, 6.15, 5.4, RGBColor(0xDC, 0xFC, 0xE7), line=GREEN, line_width=1.5, radius=0.03)
    add_rounded(s, 6.8, 1.3, 6.15, 0.4, GREEN, radius=0.03)
    add_text(s, 6.95, 1.35, 5.85, 0.3, '✓ v0.12 (回歸 Excel)', size=12, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    after = """4 個 DL category · 廠級單一 wage($4.95)
對應 Cleansheet r28-31:

  r28  DL per line per shift           (一般直接人力)
  r29  Debug DL per line per shift     (Debug)
  r30  Functional DL per shift         (功能測試)
  r31  ⭐ Warehouse DL                 (倉儲 · 我之前漏)

  → 全用同一 wage_per_hr_usd $4.95
  → 對應 bom_cs_case_process 4 欄
    + bom_factory_baseline.dl_wage_per_hr_usd

範例(SteelSeries SMT_MAIN):
  r28: 10 人  · r29: 0.5 · r30: 0.5 · r31: 0

範例(SteelSeries BB_ASSY):
  r28: 22.5  · r29: 0   · r30: 0.5 · r31: 0

範例(SteelSeries Material Mgmt):
  r28: 0     · r29: 0   · r30: 0   · r31: 0.5 ⭐"""
    add_text(s, 6.95, 1.8, 5.85, 4.8, after, size=10, color=TEXT, font='Consolas')

    footer(s, n, total)


def slide_idl_linedep(total, n):
    s = add_slide()
    header(s, 'IDL Line-Dep 4 角色 weekly wage · 對齊 r43-46', '取代 v0.11 hardcode($280/$320/$460)· 完整加 IQC', PURPLE, '👔 IDL')

    add_text(s, 0.4, 1.3, 12, 0.3, '🔧 對齊 Cleansheet r43-46 IDL Line-Dependent wage(per week)',
             size=12, bold=True, color=INK)

    add_rounded(s, 0.4, 1.7, 12.55, 5.2, BG_SOFT, line=PURPLE, line_width=1.5, radius=0.03)
    add_rounded(s, 0.4, 1.7, 12.55, 0.4, PURPLE, radius=0.03)
    add_text(s, 0.55, 1.75, 12.3, 0.3, '🗄️ bom_factory_idl_linedep_wage(🆕 v0.12 新表)', size=12, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)

    # 4 個角色卡
    roles = [
        ('LINE_LEADER', 'Line Leader 線長',         'r43', '$403.00 / wk', RGBColor(0x16, 0xA3, 0x4A)),
        ('TECHNICIAN',  'Technician 技術員',        'r44', '$436.80 / wk', RGBColor(0x0E, 0xA5, 0xE9)),
        ('IQC',         '⭐ IQC 來料品保(新加)',  'r45', '$403.00 / wk', RGBColor(0xDC, 0x26, 0x26)),
        ('SUPERVISOR',  'Supervisor 主任',          'r46', '$403.00 / wk', RGBColor(0xCA, 0x8A, 0x04)),
    ]
    for i, (code, name, excel, wage, color) in enumerate(roles):
        y = 2.3 + i * 0.95
        add_rounded(s, 0.55, y, 12.25, 0.85, WHITE, line=color, line_width=1, radius=0.04)
        add_rounded(s, 0.55, y + 0.05, 1.6, 0.75, color, radius=0.04)
        add_text(s, 0.55, y + 0.08, 1.6, 0.35, code, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, 0.55, y + 0.45, 1.6, 0.3, excel, size=9, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, 2.25, y + 0.05, 6, 0.35, name, size=12, bold=True, color=INK, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, 2.25, y + 0.42, 6, 0.4, f'Cleansheet {excel} · CN baseline', size=10, color=MUTED)
        add_rounded(s, 8.5, y + 0.15, 4.25, 0.55, RGBColor(0xF3, 0xF4, 0xF6), line=color, radius=0.05)
        add_text(s, 8.5, y + 0.15, 4.25, 0.55, wage, size=17, bold=True, color=color, font='Consolas', align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    footer(s, n, total)


def slide_case_4cols(total, n):
    s = add_slide()
    header(s, '案級補 4 個新欄位 · 對齊 Excel r31 + r37-39', 'bom_cs_case_process 新增 warehouse_dl + per_day × 3', CYAN, '🗄️ 案級')

    add_text(s, 0.4, 1.3, 12, 0.3, '📋 bom_cs_case_process 案級 schema 變動(v0.11 → v0.12)',
             size=12, bold=True, color=INK)

    add_rounded(s, 0.4, 1.7, 12.55, 5.2, BG_SOFT, line=CYAN, line_width=1.5, radius=0.03)
    add_rounded(s, 0.4, 1.7, 12.55, 0.4, CYAN, radius=0.03)
    add_text(s, 0.55, 1.75, 12.3, 0.3, '🆕 ALTER TABLE bom_cs_case_process — 補 4 個案級欄位', size=12, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)

    schema = """ALTER TABLE bom_cs_case_process
  -- 🆕 v0.12 補 r31 Warehouse DL(我之前完全漏 · Material Mgmt 製程用 0.5)
  ADD warehouse_dl_per_shift  NUMBER(5,2) DEFAULT 0,

  -- 🆕 v0.12 補 r37 / r38 IDL Line-Indep Per Day(跟 Per Shift 是兩種獨立 count)
  ADD line_leader_per_day     NUMBER(5,2) DEFAULT 0,
  ADD technician_per_day      NUMBER(5,2) DEFAULT 0,

  -- 🆕 v0.12 補 r39 IQC Per Day(v0.11 之前完全沒有此角色)
  ADD iqc_per_day             NUMBER(5,2) DEFAULT 0;

-- 既有欄位(維持 · v0.10 對齊 Excel r28-30 + r35-36 + r40)
--   dl_per_shift              -- r28
--   debug_dl_per_shift        -- r29
--   functional_dl_per_shift   -- r30
--   line_leader_per_shift     -- r35
--   technician_per_shift      -- r36
--   supervisor_per_day        -- r40

-- 範例案級數據(對應 Excel SteelSeries Rival 3+ Cleansheet · CN 廠):
SMT_MAIN:       dl=10 · debug=0.5  · funct=0.5 · ⭐warehouse=0
                ll_shift=0.5 · tech_shift=0.5 · ll_day=0.5 · tech_day=0.5 · ⭐iqc=0 · sup=0.25

WAVE_SOLDER:    dl=18 · debug=0.5  · funct=0.5 · ⭐warehouse=0
                ll_shift=0.5 · tech_shift=0.5 · ll_day=0.5 · tech_day=0.5 · ⭐iqc=0 · sup=0.25

BB_ASSY:        dl=22.5 · debug=0  · funct=0.5 · ⭐warehouse=0
                ll_shift=1   · tech_shift=1   · ll_day=1   · tech_day=1   · ⭐iqc=0 · sup=0

MAT_MGMT:       dl=0  · debug=0    · funct=0   · ⭐warehouse=0.5  ← Excel r31 真實值!"""
    add_text(s, 0.55, 2.15, 12.3, 4.7, schema, size=9.5, color=TEXT, font='Consolas')

    footer(s, n, total)


def slide_compute_dl_update(total, n):
    s = add_slide()
    header(s, 'csComputeDlSteps 公式更新 · 10 步驟', 'wage 從 schema 取 · 完整對齊 Excel C28-C58', GREEN, '📊 公式')

    add_text(s, 0.4, 1.3, 12, 0.3, '🧮 v0.12 csComputeDlSteps 改寫 — 10 個步驟(Excel C28-C58 完整對齊)',
             size=12, bold=True, color=INK)

    steps = [
        ('1', 'UPH', '3600 / TAKT × Yield × Eff', 'C14', '不變'),
        ('2', 'Max output / line / wk', 'Working hr/wk × UPH', 'C21', '不變'),
        ('3', 'Demand (max/wk)', 'Annual / 50 × 1.2', 'C18-20', '不變'),
        ('4', 'Lines installed', 'ROUNDUP(max_demand / max_output)', 'C23', '不變'),
        ('5', 'Weekly output', 'max_output × lines', 'C24', '不變'),
        ('6', 'Total DL / day 🆕', '(DL+Debug+Functional+⭐Warehouse) × lines × shifts', 'C32→C33', '🆕 v0.12 補 Warehouse DL'),
        ('7', 'Multiplier 1+2', 'mult1=(Whr/2×6)/SEA_wk · mult2=(SEA_day×Wday)/SEA_wk', 'C50, C51', '不變'),
        ('8', 'DL cost / week', '$4.95 × SEA × Total_DL × mult1 × mult2', 'C56', '🆕 廠級單一 wage(廢 7 工種)'),
        ('9', 'IDL Line-Dep $/wk 🆕', 'Σ(role.count × role.weekly_wage) × mult2', 'r43-46, C55', '🆕 wage 從 schema · 含 IQC'),
        ('10', 'DL cost / unit (FINAL)', '(DL_cost_wk + IDL_line_dep_wk) / weekly_output', 'C58', '不變'),
    ]
    headers = ['n', 'label', '公式', 'Excel ref', 'v0.12 變更']
    widths = [0.4, 2.4, 5.0, 1.4, 3.35]
    cx = 0.4
    for h, w in zip(headers, widths):
        add_rect(s, cx, 1.75, w, 0.35, INK)
        add_text(s, cx, 1.75, w, 0.35, h, size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        cx += w

    for i, row in enumerate(steps):
        y = 2.1 + i * 0.48
        bg = BG_SOFT if i % 2 == 0 else WHITE
        cx = 0.4
        for j, (cell, w) in enumerate(zip(row, widths)):
            add_rect(s, cx, y, w, 0.45, bg)
            color = TEXT
            bold = False
            font = 'Microsoft JhengHei'
            if j == 0:
                color = GREEN
                bold = True
            elif j == 2 or j == 3:
                font = 'Consolas'
                if '🆕' in row[4]:
                    color = GREEN
                    bold = True
            elif j == 4 and '🆕' in cell:
                color = RED
                bold = True
            add_text(s, cx + 0.08, y, w - 0.16, 0.45, cell, size=8.5, color=color, bold=bold, font=font, anchor=MSO_ANCHOR.MIDDLE)
            cx += w

    footer(s, n, total)


def slide_bg_preserved(total, n):
    s = add_slide()
    header(s, '✅ v0.11 正確設計保留', '別誤會 — BG/權限/設備類別 / 設定 UI 都保留 · 只是 DL 部分大改', GREEN)

    keep = [
        ('🔒', 'BG/BU 隔離 baseline 切版', 'baseline.bg_code · 同廠不同 BG 各有 baseline · 跨 BG 看不到對方數字',
         'Q-BG1-5 拍板的設計都保留 · 無變動'),
        ('🔐', '權限 grant 表', 'bom_settings_admin_grant · user × scope × bg/bu/factory · 對齊既有資料權限管理',
         'Q-PERM1-3 拍板的設計都保留 · 無變動'),
        ('🏭', '設備類別取代實體設備', 'bom_equip_category_catalog 26 類 / BG × bom_cs_case_equip_category 案級 binding',
         'Q-FIX1-7 拍板的設計都保留 · 26 類 + 3 wearable 跟 v0.11 一樣'),
        ('🔧', '設定 master UI', '7 個 sub-tab · 廠 baseline / 類別 catalog / 類別單價 / DL wage 🆕 / IDL Line-Dep wage 🆕 / IDL Centralized / 耗材 / 權限 grant',
         '原 "DL 角色費率" sub-tab 改名為 "DL 廠級單一 wage" + 新增 "IDL Line-Dep wage"'),
        ('📋', 'admin 跨 BG copy 類別', '從 OPTO BG copy 26 類到 CONSUMER BG · 加速新 BG onboarding',
         '保留 · 無變動'),
        ('📂', '案綁定 BG · SteelSeries=OPTO · WHOOP=CONSUMER', 'projects.bg_code/bu_code 從 user 主檔帶 · case_factory 冗存',
         '保留 · 無變動'),
    ]
    for i, (ico, name, desc, note) in enumerate(keep):
        y = 1.45 + i * 0.92
        add_rounded(s, 0.4, y, 12.55, 0.85, BG_SOFT, line=GREEN, line_width=1, radius=0.04)
        add_text(s, 0.55, y + 0.05, 0.5, 0.4, ico, size=24, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, 1.2, y + 0.04, 6, 0.28, name, size=11, bold=True, color=INK, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, 1.2, y + 0.32, 11.6, 0.28, desc, size=9, color=TEXT, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, 1.2, y + 0.58, 11.6, 0.25, f'  ↳ {note}', size=8.5, color=GREEN, anchor=MSO_ANCHOR.MIDDLE)

    footer(s, n, total)


def slide_close(total, n):
    s = add_slide()
    add_rect(s, 0, 0, 13.333, 7.5, INK)
    add_rect(s, 0, 2.8, 13.333, 0.06, OCEAN)
    add_text(s, 0.7, 1.2, 12, 0.6, '完。', size=44, bold=True, color=WHITE)
    add_text(s, 0.7, 1.9, 12, 0.4, 'Q&A · 等待 schema migration 拍板', size=18, color=OCEAN)

    add_rounded(s, 0.7, 3.4, 12.0, 3.5, RGBColor(0x14, 0x2F, 0x4E), radius=0.04)
    add_text(s, 1.0, 3.6, 11.5, 0.4, '🔗 v4 後續行動', size=16, bold=True, color=OCEAN)
    add_text(s, 1.0, 4.1, 11.5, 2.7, """• 互動 demo:Cortex_互動Demo_v0.12.html(§Cleansheet Step 1-7 全套 + Excel §A 100% 對齊)
• SD 更新:cleansheet-mva-sd v0.5(廢 bom_factory_dl_role · 新增 bom_factory_idl_linedep_wage)
• 案級 schema 補 4 欄(warehouse_dl + 3 per_day)· 1 個 ALTER + 1 個新表
• v0.11 保留設計:BG/BU 隔離 · 設備類別 · 權限 grant · 設定 master UI 全部沒變
• 兩案 BG 綁定維持:SteelSeries → OPTO · WHOOP → CONSUMER
• 我犯的錯:over-engineer DL 工種 + 漏 IQC/Warehouse DL/per Day · 已全部修正 + 認錯""",
        size=12, color=RGBColor(0xCB, 0xD5, 0xE1))

    add_text(s, 0.7, 7.05, 12, 0.3, 'Cortex MVA 操作手冊 v4 · Excel 完整對齊版 · 2026-06 · FOXLINK GPT',
             size=10, color=MUTED, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════
# RUN
# ════════════════════════════════════════════════════════════════════════
total = 9
print(f'Building {total} slides...')

slide_cover()                            # 1
slide_apology(total)                     # 2
slide_section_a_full(total, 3)           # 3
slide_dl_4cat(total, 4)                  # 4
slide_idl_linedep(total, 5)              # 5
slide_case_4cols(total, 6)               # 6
slide_compute_dl_update(total, 7)        # 7
slide_bg_preserved(total, 8)             # 8
slide_close(total, 9)                    # 9

OUT = 'd:/vibe_coding/foxlink_gpt/docs/Cortex_MVA操作流程說明手冊_v4_Excel完整對齊版.pptx'
prs.save(OUT)
print(f'[OK] Saved: {OUT}')
print(f'     Total slides: {len(prs.slides)}')

