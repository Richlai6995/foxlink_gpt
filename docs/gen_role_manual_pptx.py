# -*- coding: utf-8 -*-
"""
Cortex v0.13 角色工作區 · 逐角色操作測試說明書
對齊 Cortex_互動Demo_v0.13.html 實際可點操作
~22 slides · 16:9 · 給 user 照著測每個角色
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

INK=RGBColor(0x0B,0x1F,0x3A); TEXT=RGBColor(0x37,0x47,0x55); MUTED=RGBColor(0x6B,0x72,0x80)
OCEAN=RGBColor(0x02,0xC3,0x9A); NAVY=RGBColor(0x1C,0x72,0x93); LINE=RGBColor(0xE5,0xE7,0xEB)
BG=RGBColor(0xFA,0xFB,0xFC); WHITE=RGBColor(0xFF,0xFF,0xFF); GOLD=RGBColor(0xCA,0x8A,0x04)
RED=RGBColor(0xDC,0x26,0x26); GREEN=RGBColor(0x16,0xA3,0x4A); PURPLE=RGBColor(0x7C,0x3A,0xED)
CYAN=RGBColor(0x08,0x91,0xB2); SKY=RGBColor(0x0E,0xA5,0xE9); AMBER=RGBColor(0xF5,0xA5,0x24)

def slide(): return prs.slides.add_slide(prs.slide_layouts[6])
def T(s,x,y,w,h,t,sz=12,b=False,c=None,al=PP_ALIGN.LEFT,an=MSO_ANCHOR.TOP,f='Microsoft JhengHei'):
    tx=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h));tf=tx.text_frame
    tf.margin_left=tf.margin_right=tf.margin_top=tf.margin_bottom=0;tf.word_wrap=True;tf.vertical_anchor=an
    p=tf.paragraphs[0];p.alignment=al;r=p.add_run();r.text=t;r.font.name=f;r.font.size=Pt(sz);r.font.bold=b
    if c:r.font.color.rgb=c
    return tx
def box(s,x,y,w,h,fill,ln=None,lw=0.75,rad=0.04):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(w),Inches(h))
    sh.adjustments[0]=rad;sh.fill.solid();sh.fill.fore_color.rgb=fill
    if ln is None:sh.line.fill.background()
    else:sh.line.color.rgb=ln;sh.line.width=Pt(lw)
    sh.shadow.inherit=False;return sh
def rect(s,x,y,w,h,fill):
    sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(x),Inches(y),Inches(w),Inches(h))
    sh.fill.solid();sh.fill.fore_color.rgb=fill;sh.line.fill.background();sh.shadow.inherit=False;return sh

def head(s,title,sub,color=NAVY,route=None):
    rect(s,0,0,13.333,0.16,color)
    T(s,0.45,0.26,9.5,0.5,title,22,True,INK)
    if sub:T(s,0.45,0.82,9.5,0.34,sub,11,False,MUTED)
    if route:
        box(s,10.1,0.30,2.8,0.5,color,rad=0.3)
        T(s,10.1,0.31,2.8,0.48,route,11,True,WHITE,PP_ALIGN.CENTER,MSO_ANCHOR.MIDDLE,'Consolas')
def foot(s,n,tot):
    rect(s,0.45,7.08,12.45,0.02,LINE)
    T(s,0.45,7.13,9,0.22,'Cortex v0.13 角色工作區 · 逐角色測試說明書',8.5,False,MUTED)
    T(s,12.0,7.13,1,0.22,f'{n} / {tot}',8.5,False,MUTED,PP_ALIGN.RIGHT)

# 三欄卡(① landing ② menu ③ 可測試)
def role_slide(n,tot,role_name,route,ws,user,perm,color,landing,menu,testable,handoff):
    s=slide()
    head(s,role_name,f'工作區:{ws} · 登入身分:{user}',color,route)
    # perm strip
    box(s,0.45,1.28,12.45,0.4,BG,LINE,0.5,0.06)
    T(s,0.6,1.31,12.2,0.34,'🔑 '+perm,10.5,False,TEXT,MSO_ANCHOR.MIDDLE if False else MSO_ANCHOR.TOP)
    # 3 columns
    cols=[('① 登入後看到(Landing)','📍',landing,SKY),('② 左 rail menu(可點)','📋',menu,color),('③ 可測試的操作','✅',testable,GREEN)]
    cw=4.05; x0=0.45; y0=1.85; ch=4.3
    for i,(t,ico,items,cc) in enumerate(cols):
        x=x0+i*(cw+0.08)
        box(s,x,y0,cw,ch,WHITE,cc,1,0.03)
        box(s,x,y0,cw,0.42,cc,rad=0.04)
        T(s,x+0.12,y0+0.06,cw-0.24,0.3,f'{ico} {t}',11,True,WHITE,an=MSO_ANCHOR.MIDDLE)
        yy=y0+0.55
        for it in items:
            T(s,x+0.15,yy,cw-0.3,0.6,'• '+it,9.5,False,TEXT)
            yy+=0.13+0.052*max(1,len(it)//16)
    # handoff strip
    box(s,0.45,6.3,12.45,0.62,RGBColor(0xFE,0xF3,0xC7),GOLD,0.75,0.04)
    T(s,0.6,6.37,12.2,0.5,'🔗 交棒(handoff):'+handoff,10,False,RGBColor(0x92,0x40,0x0E))
    foot(s,n,tot)

TOT=22

# ---- 1 cover ----
s=slide()
rect(s,0,0,13.333,7.5,INK);rect(s,0,2.75,13.333,0.06,OCEAN)
T(s,0.7,1.25,12,0.4,'FOXLINK 正崴 · Cortex 報價/成本平台',12,False,OCEAN)
T(s,0.7,1.65,12,0.95,'角色工作區 · 操作測試說明書',40,True,WHITE)
T(s,0.7,3.0,12,0.5,'10 角色 × 各自獨立操作工作區 · 對齊 Cortex_互動Demo_v0.13.html',18,False,WHITE)
box(s,0.7,3.95,12,2.1,RGBColor(0x14,0x2F,0x4E),rad=0.04)
T(s,1.0,4.12,11.5,0.4,'怎麼用這份說明書測試',14,True,OCEAN)
for i,t in enumerate([
  '1. 開 Cortex_互動Demo_v0.13.html → 右上「切換角色登入身分」下拉選角色',
  '2. 每換一個角色 → 左 rail menu / landing / 資料投影 全部跟著換(這就是「角色工作區制」)',
  '3. 照本說明書每頁的「③ 可測試的操作」逐項點點看',
  '4. 記下哪裡怪 / 想調整 → 回來跟我討論(本 demo 是導覽範式驗證 · 不是全功能)',
]):
    T(s,1.0,4.6+i*0.34,11.5,0.32,t,11.5,False,RGBColor(0xCB,0xD5,0xE1))
T(s,0.7,6.45,12,0.4,'v0.13 · 2026-06 · 8 operational + 2 oversight · 取代 v0.12 單一 form',11,False,MUTED)

# ---- 2 導覽總覽 ----
s=slide()
head(s,'導覽總覽:10 角色 → 10 工作區','登入 → 依角色 landing 到專屬 route · 不再進單一 form')
rows=[
 ('BPM 業務專案經理','/quote/cockpit','開案 + gate + 報價 orchestration',NAVY),
 ('RD 研發','/rd/desk','上傳 BOM + AI 解析 + 料號確認',SKY),
 ('採購 Buyer(含 PKG)','/buyer/cockpit','採購策略 + RFQ 雙價 + PKG 配置',CYAN),
 ('EPM 廠製造工程經理','/epm/desk','廠 baseline + 案級 Cleansheet + compute',GREEN),
 ('DPM 開發總監','/dpm/review','cost review + BOM lock + propagate',RED),
 ('業務 Sales(兼 Finance)','/sales/quote-desk','報價 + Margin 稽核 + 議價',AMBER),
 ('採購主管','/proc/approval','採購策略審核 + 跨案比價',PURPLE),
 ('admin / IT','/admin/cost-master','設定 master + 權限 grant',MUTED),
 ('總經理 / BG 高層','/portfolio','該 BG 所有案 portfolio(唯讀)',NAVY),
 ('董事長 / 集團最高','/portfolio?scope=GLOBAL','全 BG 所有案 portfolio(唯讀)',INK),
]
y=1.45
for nm,rt,desc,c in rows:
    box(s,0.45,y,12.45,0.5,BG,c,0.75,0.06)
    box(s,0.55,y+0.08,0.34,0.34,c,rad=0.3)
    T(s,3.6,y+0.09,3.6,0.34,rt,10,True,c,an=MSO_ANCHOR.MIDDLE,f='Consolas')
    T(s,0.95,y+0.09,2.6,0.34,nm,11,True,INK,an=MSO_ANCHOR.MIDDLE)
    T(s,7.3,y+0.09,5.4,0.34,desc,10,False,TEXT,an=MSO_ANCHOR.MIDDLE)
    y+=0.555
foot(s,2,TOT)

# ---- 3 handoff ----
s=slide()
head(s,'端到端流程 · 角色怎麼交棒','QUOTE 8 stage · BPM 是 gate 唯一確認人 · 完成即交棒下一角色')
chain=[('業務','發 RFQ',GREEN),('BPM','開案 8-stage',GREEN),('RD','Stage4 BOM',GREEN),('採購∥EPM','Stage5 Collect',AMBER),('DPM','Stage6 lock',NAVY),('業務','Stage7 報價',MUTED),('業務','Stage8 結案',MUTED)]
cw=1.72;x=0.5
for nm,ds,c in chain:
    box(s,x,1.7,cw,0.95,WHITE,c,1.2,0.05)
    T(s,x,1.8,cw,0.3,nm,11,True,c,PP_ALIGN.CENTER)
    T(s,x,2.15,cw,0.4,ds,9.5,False,TEXT,PP_ALIGN.CENTER)
    if x+cw<12:T(s,x+cw-0.02,1.95,0.2,0.4,'▶',12,False,MUTED,PP_ALIGN.CENTER,MSO_ANCHOR.MIDDLE)
    x+=cw+0.04
for i,(t,c) in enumerate([
  ('🔄 reprice 回圈:業務議價破底線 → 開新 stage iteration 退 Stage5 → 採購重議+EPM 重 compute+DPM 重 lock → 回 Stage7',GOLD),
  ('📊 oversight 旁路:DPM lock+propagate 後 margin 即時聚合進 GM(/portfolio)+ 董事長(GLOBAL)· 唯讀不交棒',CYAN),
  ('⚙️ admin 旁路:不在 stage 鏈 · 月度交付 baseline/類別/wage/權限 給所有作業角色',PURPLE),
]):
    box(s,0.45,3.15+i*0.72,12.45,0.62,RGBColor(0xFA,0xFB,0xFC),c,0.75,0.04)
    T(s,0.6,3.24+i*0.72,12.2,0.5,t,10.5,False,TEXT)
T(s,0.45,5.55,12,0.34,'💡 測試提示:每個角色 menu 都有「🔗 端到端流程」可看這張 handoff 圖 + 三案目前位置',11,True,NAVY)
foot(s,3,TOT)

# ---- 4-13 各角色(operational 8 + oversight 2) ----
roles=[
 # BPM
 (4,'BPM 業務專案經理','/quote/cockpit','我的報價案工作台','Lisa 林(OPTO/MOUSE)','可看 true cost(VIEW_TRUE_COST)· owner+gate 唯一確認人',NAVY,
  ['案卡牆:我 owner 的在辦案(SteelSeries/Razer/WHOOP)','KPI:在辦/待簽 Gate/SLA/議價中','「待我處理」收件匣(待簽 stage · reprice 退回)'],
  ['🏠 Cockpit 首頁','➕ 開新報價案','📊 案進度儀表','👥 角色邀請配置','🚦 Gate 控制台(badge 2)','💰 成本檢視','📤 報價 Console','🔗 端到端流程'],
  ['點「➕ 開新報價案」→ 5 步 wizard 逐步點「下一步」(基本→model→qty→pkg→廠別+變體)','點案卡 → 案進度 swimlane(8 stage)','「🚦 Gate 控制台」→ 點放行/退回','「👥 角色邀請」看 field_grants 角色 template'],
  '接業務 RFQ → 開案 → 每 gate 放行給下一角色 → Stage8 結案'),
 # RD
 (5,'RD 研發','/rd/desk','我的 BOM 工作台','Troy 黃(OPTO/MOUSE)','不看雙價(RD 不渲染 cost 欄)· 只做 BOM 結構',SKY,
  ['指派給我的 BOM 案卡(Razer 進行中 / SteelSeries 已交棒唯讀)','KPI:指派/待提交/parse 紅旗/Reprice','4-step 進度條'],
  ['🏠 BOM Desk 首頁','🧱 BOM 建置 4-step','🧩 Module Board','🔗 端到端流程'],
  ['點案卡 → 「🧱 BOM 建置 4-step」看上傳→AI 解析→ERP 料號→Variant 4 步','「🧩 Module Board」看 WHOOP 8 子組件(6 EE + 2 ME)','看 EE→shared / ME→per_variant 自動預設'],
  '接 BPM 開案 → rd_submitted 交棒採購+EPM(仍 DRAFT)'),
 # BUYER
 (6,'採購 Buyer(含 PKG)','/buyer/cockpit','採購工作台','David 張(OPTO/MOUSE)','可看 true cost / markup(採購核心欄)',CYAN,
  ['三桶卡片牆:待採/待 Review/被退回','KPI:待報價子料/markup 異常/ERP 未拉','跨案待採(SteelSeries/WHOOP)'],
  ['🏠 採購工作台','🎯 採購策略工作台','📋 RFQ 報價錄入','🔍 ERP 歷史價','📮 PKG 配置','🔗 端到端流程'],
  ['點 SteelSeries 卡 → 「🎯 採購策略」子料矩陣(ERP 參考 / true / quote / markup / is_chosen)','「📋 RFQ 報價錄入」看雙價 tier(Low/High · markup 自動)','點 WHOOP 卡 → 「📮 PKG 配置」5 SKU × 8 module 矩陣'],
  '接 RD BOM → 採購主管核可策略 → 交 DPM(每子料選定 tier)'),
 # EPM
 (7,'EPM 廠製造工程經理','/epm/desk','EPM 製造工程工作台','Andy 陳(OPTO/MOUSE)','廠別=參數(只看 grant 的廠)· 看製造成本',GREEN,
  ['廠別切換器 CN/VN/TW','baseline 健康卡(月度燈號)','廠級 vs 案級 雙動線卡'],
  ['🏠 工作台首頁','📅 廠 Baseline 維護','🧮 案級 Cleansheet','⚙️ Compute & 廠 Lock','🔗 端到端流程'],
  ['頂部點 CN/VN/TW 切廠別','「📅 廠 Baseline 維護」看 7 sub-tab(DL 廠級單一 wage)','「🧮 案級 Cleansheet」→ 切 FULL_MVA / SIMPLIFIED → 看兩 model 不同 component','「⚙️ Compute」看 ε<0.01 + 廠 Lock'],
  '接採購材料價 → compute+廠 Lock(per case_factory)→ 交 DPM'),
 # DPM
 (8,'DPM 開發總監','/dpm/review','Cost Review 控制台','Mike 王(OPTO/MOUSE)','VIEW_TRUE_COST · cost 唯讀 · final lock 限 DPM',RED,
  ['待 Review 收件匣(SteelSeries 待 lock)','待二簽 baseline','KPI:待 Review/二簽/本月 lock/margin 紅燈'],
  ['🏠 Review 收件匣(badge 1)','🔢 Cost Matrix 控制台','🔒 BOM Lock Gate','✍️ Baseline 漲幅二簽','🔗 端到端流程'],
  ['收件匣點 SteelSeries → 「🔢 Cost Matrix」多維 pivot + margin(VN-Black ⭐ 最便宜)','「🔒 BOM Lock Gate」→ propagate dry-run 預覽 24 列','「✍️ 二簽」看 DL wage 漲幅 approve/reject'],
  '接採購+EPM(compute 完)→ final lock+propagate(寫唯一 fact)→ 交業務'),
 # SALES
 (9,'業務 Sales(兼 Finance)','/sales/quote-desk','報價工作台','John 李(OPTO/MOUSE)','VIEW_TRUE_COST(兼 Finance 才有)· 只消費 fact',AMBER,
  ['三 tab:待報價/議價循環/Margin 稽核','KPI:待報價/議價中/SLA/加權 margin','只看 DPM lock 後的 fact'],
  ['🏠 My Quotes 首頁','📝 待報價 To-Quote','🔄 議價循環','📈 Margin 稽核台','📤 客戶報價匯出','🔗 端到端流程'],
  ['首頁切三 tab(To-Quote / 議價 / Margin 稽核)','「📈 Margin 稽核」看 (quote-true)/quote amount 加權 + Top markup','「🔄 議價循環」看 A 只動 quote / B reprice 退 Stage5 兩路徑','「📤 匯出」看同引擎讀 mask 出 FULL/SIMPLIFIED 兩版面'],
  '接 DPM lock → 報價/議價 → 路徑B reprice 退採購/EPM/DPM → 成交歸檔'),
 # PROC_MGR
 (10,'採購主管','/proc/approval','採購審核台','Karen 周(OPTO/MOUSE)','VIEW_TRUE_COST · 策略核可 gate',PURPLE,
  ['跨案待核佇列(SLA+金額排序)','比價警示(同料跨案價差/離群)','KPI:待審/待核 tier/離群/核可金額'],
  ['🏠 Approval Desk','✅ 採購策略審核','💲 Price Tier 核可','📊 跨案 Supplier 比價','🔗 端到端流程'],
  ['收件匣點 SteelSeries 策略 → 看 preview 試算','「📊 跨案 Supplier 比價」看同 Foxlink P/N 跨案價差 + 離群警示','核可後 procurement_gate_status=approved'],
  '接採購提交策略 → 核可 → 放行進 DPM(或退回採購補)'),
 # ADMIN
 (11,'admin / IT','/admin/cost-master','設定 Master 控制台','Cortex Admin(HQ/IT)','跨 BG · scope=ALL · super-admin',MUTED,
  ['平台健康儀表(baseline 就緒 tiles)','9 子模組卡','跨 BG 全廠'],
  ['🏠 平台健康','📅 廠 Baseline','🏭 設備類別 Catalog','👷 工資 Master','🎛️ Costing Model × Mask','🔐 權限 Grant'],
  ['「🎛️ Costing Model × Mask」看 component 啟用對照(FULL/SIMPLIFIED + fallback_into)','「🏭 設備類別 Catalog」看跨 BG copy(OPTO→CONSUMER)','「🔐 權限 Grant」看 D10 單一源=資料權限管理'],
  '不在 stage 鏈 · 月度交付可信 master + 正確權限給所有作業角色'),
 # GM_BG
 (12,'總經理 / BG 高層','/portfolio','BG Portfolio 戰情中心','David 總(OPTO BG)','唯讀 · 該 BG 所有 BU · VTC 預設有(GM 進名單)',NAVY,
  ['BG 戰情 KPI tiles(在辦/報價額/margin/紅燈)','全案清單(唯讀 drill)','只該 BG(OPTO)的案'],
  ['🏠 BG 戰情總覽','📋 全案清單','🔴 卡關紅燈中心','📈 Margin 風險雷達','🫗 Stage 流量漏斗'],
  ['看 BG portfolio KPI 磚','點案卡 drill(唯讀)','「📈 Margin 風險雷達」看該 BG margin(VTC gate)','確認:看不到 CONSUMER BG 的案'],
  'oversight 唯讀 dead-end · 治理只走線下/訊息催 DPM/業務 · 不交棒不推進'),
 # CHAIRMAN
 (13,'董事長 / 集團最高','/portfolio?scope=GLOBAL','集團 Portfolio 戰情儀表板','Chairman(集團)','唯讀 · 全 BG 不過濾 · true cost 恆開 · 機密穿透 🛡️',INK,
  ['集團 KPI tiles(全 BG 彙總)','各 BG swimlane(OPTO + CONSUMER)','全 3 案都看得到(GLOBAL)'],
  ['🏠 集團總覽','⚖️ BG 對比','🌐 全集團專案瀏覽器','📈 Margin 與雙價分析','📑 匯出 Board Pack'],
  ['看集團 KPI + 各 BG swimlane','確認:OPTO + CONSUMER 兩 BG 案全看得到(GM 只看單 BG)','看 true cost 全穿透(右上綠盾概念)'],
  'oversight 唯讀 · 複用 /portfolio + GLOBAL toggle(D9)· 機密穿透最上層'),
]
for r in roles:
    role_slide(r[0],TOT,r[1],r[2],r[3],r[4],r[5],r[6],r[7],r[8],r[9],r[10])

# ---- 14 兩 model 差異(測試重點)----
s=slide()
head(s,'測試重點:同一系統 · 兩種計價模型','切 SteelSeries(FULL_MVA)vs WHOOP(SIMPLIFIED)· 看同角色不同投影')
box(s,0.45,1.4,6.2,5.2,RGBColor(0xF0,0xF9,0xFF),SKY,1.5,0.03)
box(s,0.45,1.4,6.2,0.5,SKY,rad=0.03)
T(s,0.6,1.48,5.9,0.36,'🖱️ SteelSeries Rival 3+(FULL_MVA · OPTO)',13,True,WHITE,an=MSO_ANCHOR.MIDDLE)
T(s,0.62,2.05,5.95,4.4,
 '在哪測:採購/EPM/DPM/業務 切到此案\n\n'
 '• 採購:商包/工包 PKG 版\n'
 '• EPM Cleansheet:選 FULL_MVA → 細製程 9×§A + IDL matrix 17×9 + 設備類別折舊\n'
 '• 啟用 component:DL_CPU/IDL_CPU/EQUIP/FACILITY/FREIGHT/VAT/LOSS/MVA_TOTAL\n'
 '• DPM Cost Matrix:3 廠×2 variant×2 qty×2 pkg = 24 cells\n'
 '• 業務:廠×variant 報價版面\n'
 '• margin ~6.7%',11,False,TEXT)
box(s,6.85,1.4,6.0,5.2,RGBColor(0xFA,0xF5,0xFF),PURPLE,1.5,0.03)
box(s,6.85,1.4,6.0,0.5,PURPLE,rad=0.03)
T(s,7.0,1.48,5.7,0.36,'⌚ WHOOP Gen4 MP(SIMPLIFIED · CONSUMER)',13,True,WHITE,an=MSO_ANCHOR.MIDDLE)
T(s,7.02,2.05,5.75,4.4,
 '在哪測:採購/EPM 切到此案 · 董事長看得到\n\n'
 '• RD:Module Board 8 子組件\n'
 '• 採購:PKG 5 SKU × 8 module 矩陣\n'
 '• EPM Cleansheet:選 SIMPLIFIED → 放大製程 macro + SMT 點數 + 材料耗損率\n'
 '• 啟用 component:PROC_MACRO/SMT_POINTS/MAT_LOSS_RATE/OVERHEAD_4PCT/TRANSPORTATION\n'
 '• disable 的 IDL/設備 fallback → OVERHEAD 4% 吸收\n'
 '• margin ~6.0%',11,False,TEXT)
foot(s,14,TOT)

# ---- 15 權限/機密測試 ----
s=slide()
head(s,'測試重點:權限 + 機密遮罩','切不同角色 → 看同一 cost matrix 顯示不同')
T(s,0.45,1.4,12,0.34,'切 RD(無 VIEW_TRUE_COST)vs DPM/業務(有)→ 看 true cost / margin 欄變化',12,True,NAVY)
data=[
 ('RD 研發','❌ 無 VTC','根本不渲染雙價欄 · 只做 BOM 結構',RED),
 ('採購 Buyer','✅ 有 VTC','true cost / markup 可見可改(核心欄)',GREEN),
 ('EPM','✅ 有 VTC','看製造成本 · 不碰材料 quote',GREEN),
 ('DPM','✅ 有 VTC','全看 · cost 唯讀 · 只能 lock',GREEN),
 ('業務(兼 Finance)','✅ 有 VTC','Margin 稽核解鎖 · markup=🔒🔒 每 access 寫 audit',GREEN),
 ('GM_BG','✅ 有 VTC','該 BG margin 雷達 · 看不到他 BG',AMBER),
 ('董事長','🛡️ 全穿透','全 BG true cost 恆開',PURPLE),
]
y=1.9
for nm,vtc,desc,c in data:
    box(s,0.45,y,12.45,0.58,BG,c,0.75,0.05)
    T(s,0.6,y+0.13,2.6,0.34,nm,11,True,INK)
    T(s,3.2,y+0.13,2.0,0.34,vtc,11,True,c)
    T(s,5.4,y+0.13,7.3,0.34,desc,10.5,False,TEXT)
    y+=0.66
T(s,0.45,6.5,12,0.34,'💡 在 DPM「🔢 Cost Matrix」或業務「📈 Margin 稽核」看完整;切 RD 對比(RD 看不到雙價)',10.5,True,NAVY)
foot(s,15,TOT)

# ---- 16 測試 checklist ----
s=slide()
head(s,'逐角色測試 Checklist','照這張勾 · 測完跟我討論要調哪裡')
cols=[
 ('作業軌(8)',[
  '☐ BPM:開案 wizard 5 步走完','☐ BPM:Gate 控制台放行/退回','☐ RD:4-step + Module Board 8 件',
  '☐ 採購:策略矩陣 + RFQ tier + PKG 矩陣','☐ EPM:廠別切換 + FULL/SIMPLIFIED 切換','☐ EPM:Compute + 廠 Lock',
  '☐ DPM:Cost Matrix + Lock Gate dry-run','☐ 業務:三 tab + 議價 A/B + 匯出','☐ 採購主管:策略審核 + 跨案比價','☐ admin:component mask 表 + 權限',
 ],NAVY),
 ('oversight(2)+ 跨切',[
  '☐ GM_BG:BG portfolio + 只看單 BG','☐ 董事長:集團 portfolio + 全 BG 穿透',
  '☐ 切 SteelSeries vs WHOOP 兩 model','☐ 切 RD vs DPM 看機密遮罩差異','☐ 每角色「🔗 端到端流程」handoff 圖',
  '','📝 測試時記下:','  · 哪個角色 menu 少了什麼','  · 哪個操作 flow 怪 / 順序不對','  · 哪些畫面要做深(目前是說明頁)',
 ],GREEN),
]
for i,(t,items,c) in enumerate(cols):
    x=0.45+i*6.25
    box(s,x,1.4,6.05,5.4,WHITE,c,1.2,0.03)
    box(s,x,1.4,6.05,0.45,c,rad=0.03)
    T(s,x+0.15,1.47,5.7,0.32,t,12,True,WHITE,an=MSO_ANCHOR.MIDDLE)
    yy=2.0
    for it in items:
        bold = it.startswith('📝')
        T(s,x+0.2,yy,5.6,0.4,it,10.5,bold,INK if bold else TEXT)
        yy+=0.46
foot(s,16,TOT)

# ---- 17 closing ----
s=slide()
rect(s,0,0,13.333,7.5,INK);rect(s,0,2.8,13.333,0.06,OCEAN)
T(s,0.7,1.3,12,0.6,'測完回來討論',38,True,WHITE)
T(s,0.7,2.0,12,0.4,'每角色操作確認後 → 我們逐角色調整',18,False,OCEAN)
box(s,0.7,3.4,12,3.3,RGBColor(0x14,0x2F,0x4E),rad=0.04)
T(s,1.0,3.6,11.5,0.4,'🔗 檔案',15,True,OCEAN)
T(s,1.0,4.1,11.5,2.5,
 '• 互動 demo:Cortex_互動Demo_v0.13.html(右上切角色測試)\n'
 '• 架構 SD:cortex-role-workspace-sd.md(10 角色完整規格)\n'
 '• 設計稿:cortex-role-operations-design.md(D1-D11 決議)\n'
 '• 底層:cortex-unified-architecture-sd.md(superset + component mask)\n\n'
 '討論方向:\n'
 '  · 哪些角色 menu / flow 要增減\n'
 '  · 哪些說明頁要做成可點深畫面\n'
 '  · handoff / gate 流程順序對不對\n'
 '  · 兩 model(FULL/SIMPLIFIED)呈現夠不夠清楚',12,False,RGBColor(0xCB,0xD5,0xE1))
T(s,0.7,7.0,12,0.3,'Cortex v0.13 角色工作區 · 操作測試說明書 · 2026-06 · FOXLINK GPT',10,False,MUTED,PP_ALIGN.CENTER)

OUT='d:/vibe_coding/foxlink_gpt/docs/Cortex_角色工作區_操作測試說明書_v1.pptx'
prs.save(OUT)
print(f'[OK] {OUT}')
print(f'     slides: {len(prs.slides)}')
