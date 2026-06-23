# -*- coding: utf-8 -*-
"""
check_overflow.py — 文字溢出色塊的「無渲染」啟發式偵測。

slide_qa.py 依賴此模組:analyze_shape(shape) -> None | {
    'overflow':      bool,    # 估算文字總高 > 容器可用高
    'text':          str,     # 該 shape 內文字(截斷預覽)
    'needed_h_pt':   float,   # 估算需要的高度 (pt)
    'avail_h_pt':    float,   # 容器扣掉內距後的可用高度 (pt)
}

原理(不靠 LibreOffice 渲染,純幾何估算):
  - 容器可用寬/高 = shape 尺寸 - text_frame 四邊內距。
  - 每段:依字級估「每行可容字數」→ 估行數(含硬換行 <a:br>)→ 累加行高 + 段後距。
  - CJK 視為全形(advance≈1.0em)、英數半形(≈0.5em)。
  - needed > avail 即判溢出。

限制:這是啟發式近似,不等於真實排版引擎。係數刻意接近實測而非過度保守,
避免把正常版面誤報成溢出害 agent 迴圈空轉。真值校驗仍須對轉出的圖人眼/vision 看一次。
"""
from pptx.util import Emu, Pt
from pptx.oxml.ns import qn
from pptx.enum.text import MSO_ANCHOR

EMU_PER_PT = 12700.0

# 字級 → 行高倍率(含行距)。python-pptx 預設單倍行距,實測 LibreOffice 約 1.2。
LINE_HEIGHT_FACTOR = 1.22
# 字元水平 advance 對字級的比例。
CJK_ADVANCE = 1.00   # 全形:一個字 ≈ 一個 em
LAT_ADVANCE = 0.52   # 英數/標點半形 ≈ 0.52 em(Calibri/Noto 拉丁約略值)
# python-pptx textbox 預設內距(未顯式設 0 時):左右 0.1in、上下 0.05in。
DEFAULT_ML = DEFAULT_MR = int(0.1 * 914400)
DEFAULT_MT = DEFAULT_MB = int(0.05 * 914400)
# 沒指定字級時的保守 fallback。
FALLBACK_SIZE_PT = 18.0


def _is_cjk(ch):
    o = ord(ch)
    return (
        0x4E00 <= o <= 0x9FFF or   # CJK 統一表意
        0x3400 <= o <= 0x4DBF or   # 擴充 A
        0x3000 <= o <= 0x303F or   # CJK 標點
        0xFF00 <= o <= 0xFFEF or   # 全形 ASCII / 半形片假名
        0x3040 <= o <= 0x30FF      # 平/片假名
    )


def _emu_to_pt(emu):
    return (emu or 0) / EMU_PER_PT


def _margin(tf, attr, default):
    v = getattr(tf, attr, None)
    return default if v is None else v


def _para_runs_size(paragraph):
    """回傳 (純文字, 該段最大字級pt)。"""
    txt = paragraph.text or ""
    sizes = []
    for r in paragraph.runs:
        if r.font.size is not None:
            sizes.append(r.font.size.pt)
    size = max(sizes) if sizes else FALLBACK_SIZE_PT
    return txt, size


def _count_hard_breaks(paragraph):
    """段落內 <a:br> 數(blue_style 用 a:br 做同段換行)。"""
    return len(paragraph._p.findall(qn('a:br')))


def _text_advance_em(txt):
    """整段文字水平總長,以 em 為單位(乘字級即得 pt)。"""
    adv = 0.0
    for ch in txt:
        if ch == '\n':
            continue
        adv += CJK_ADVANCE if _is_cjk(ch) else LAT_ADVANCE
    return adv


def _lines_for_segment(seg_txt, size_pt, avail_w_pt, word_wrap):
    """單一視覺行片段(被 <a:br> 切開後)需要幾行。"""
    if not seg_txt:
        return 1
    seg_w_pt = _text_advance_em(seg_txt) * size_pt
    if not word_wrap or avail_w_pt <= 0:
        return 1
    import math
    return max(1, math.ceil(seg_w_pt / avail_w_pt - 1e-6))


def analyze_shape(shape):
    if not getattr(shape, "has_text_frame", False):
        return None
    tf = shape.text_frame
    full_text = (tf.text or "").strip()
    if not full_text:
        return None
    # 形狀內「置中標籤」(chevron / tile / circle 的 center_text_in)是刻意短字,
    # 非矩形/置中形狀的文字可用區難以幾何估算,跳過以免誤報溢出。
    try:
        if tf.vertical_anchor == MSO_ANCHOR.MIDDLE:
            return None
    except Exception:
        pass

    w_pt = _emu_to_pt(shape.width)
    h_pt = _emu_to_pt(shape.height)
    ml = _emu_to_pt(_margin(tf, "margin_left", DEFAULT_ML))
    mr = _emu_to_pt(_margin(tf, "margin_right", DEFAULT_MR))
    mt = _emu_to_pt(_margin(tf, "margin_top", DEFAULT_MT))
    mb = _emu_to_pt(_margin(tf, "margin_bottom", DEFAULT_MB))
    avail_w_pt = w_pt - ml - mr
    avail_h_pt = h_pt - mt - mb

    word_wrap = tf.word_wrap is not False  # None/True 都當會自動換行

    needed_h_pt = 0.0
    for p in tf.paragraphs:
        txt, size_pt = _para_runs_size(p)
        # 段落被硬換行切成多個視覺片段,各自再算自動換行行數。
        segments = txt.split('\n') if txt else [""]
        # _para_runs_size 的 paragraph.text 已把 <a:br> 轉成 \n,故 split 即可;
        # 但保險加上獨立計數,取較大者。
        seg_lines = sum(_lines_for_segment(s, size_pt, avail_w_pt, word_wrap) for s in segments)
        seg_lines = max(seg_lines, _count_hard_breaks(p) + 1)
        space_after = p.space_after.pt if p.space_after is not None else 0.0
        needed_h_pt += seg_lines * size_pt * LINE_HEIGHT_FACTOR + space_after

    return {
        'overflow': needed_h_pt > avail_h_pt + 0.5,   # 0.5pt 容差
        'text': full_text[:40] + ('…' if len(full_text) > 40 else ''),
        'needed_h_pt': round(needed_h_pt, 1),
        'avail_h_pt': round(avail_h_pt, 1),
    }


if __name__ == '__main__':
    import sys
    from pptx import Presentation
    path = sys.argv[1] if len(sys.argv) > 1 else 'example_out.pptx'
    prs = Presentation(path)
    for si, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            r = analyze_shape(shape)
            if r and r['overflow']:
                print(f"[slide {si}] OVERFLOW {r['text']!r} need {r['needed_h_pt']}pt > avail {r['avail_h_pt']}pt")
