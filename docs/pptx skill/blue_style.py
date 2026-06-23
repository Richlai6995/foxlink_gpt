# -*- coding: utf-8 -*-
"""
blue_style.py — 企業級專業報告 · 藍色基調元件庫 (Executive Blue Deck helpers)

用法:
    from blue_style import *
    prs = new_deck()                       # 4:3, 10 x 7.5 in
    s = add_slide(prs)
    header(s, "副標 · 第一階段", "主標題", "重點一 · 重點二 · 重點三", badge="規劃中")
    concept_band(s, "一句話說明系統的價值與目的。")
    feature_card(s, 0.3, 2.8, 4.68, 1.72, "1", "卡片標題", "兩三行白話說明,\\n聚焦效益而非技術。")
    ...
    prs.save("out.pptx")

設計原則(務必遵守):
  1. 以藍色為主色,整份簡報只用「一個」金色重點(最關鍵的數字/節點),其餘全藍。
  2. 所有色塊都加漸層 (grad) 與柔和陰影 (add_shadow)。
  3. 文字以高階主管視角的商業語言為主,技術術語收進小字描述。
  4. 4:3 版面,標題列 navy→blue 漸層,內容卡白底藍框,留白充足。
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---------------- 配色 (blue-dominant palette) ----------------
NAVY   = "0B2D5C"   # 最深海軍藍 — 標題列起點
NAVY2  = "143C72"
BLUE   = "1E5BA8"   # 主藍 — 卡片標題列、圖示圈、流程
BLUE2  = "2E78C7"   # 亮藍 — 漸層終點
BLUE_L = "5B9BD5"   # 淺藍 — 副標、邊框
SKY    = "EAF2FB"   # 背景淺藍
SKY2   = "DCEAF9"   # 淺藍卡片漸層
CARD   = "FFFFFF"   # 卡片底色
BORDER = "D3E1F2"   # 卡片邊框
INK    = "16263D"   # 主要文字(近黑藍)
MUTE   = "5C6B80"   # 次要文字(灰藍)
GOLD   = "C99A3B"   # 唯一暖色重點(整頁只用一次)
GOLD2  = "E0B85A"
WHITE  = "FFFFFF"

import os
# 中文字型:預設 Noto Sans TC(dev 端已裝 + Linux prod 內建,渲染一致)。
# dev 想用微軟正黑體可設 BLUE_STYLE_EA_FONT="Microsoft JhengHei"。
EAFONT = os.environ.get("BLUE_STYLE_EA_FONT", "Noto Sans TC")
LAT    = os.environ.get("BLUE_STYLE_LAT_FONT", "Calibri")   # 英數字型

# ---------------- 基礎建構 ----------------
def new_deck():
    """建立 10 x 7.5 in (4:3) 空白簡報。"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    return prs

def add_slide(prs, bg=SKY):
    """加入一張空白版面投影片,並鋪滿淺藍背景。"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    r.fill.solid(); r.fill.fore_color.rgb = RGBColor.from_string(bg)
    r.line.fill.background(); r.shadow.inherit = False
    return s

# ---------------- 文字 / 字型 ----------------
def set_run_font(run, size, color, bold=False, ea=EAFONT, lat=LAT):
    """設定 run 字型,同時指定中文(ea)與英數(latin)字型。"""
    run.font.size = Pt(size); run.font.bold = bold
    run.font.color.rgb = RGBColor.from_string(color); run.font.name = lat
    rPr = run._r.get_or_add_rPr()
    for tag in ('a:latin', 'a:ea', 'a:cs'):
        if rPr.find(qn(tag)) is None:
            rPr.append(rPr.makeelement(qn(tag), {}))
    rPr.find(qn('a:latin')).set('typeface', lat)
    rPr.find(qn('a:ea')).set('typeface', ea)
    rPr.find(qn('a:cs')).set('typeface', lat)

def add_text(sl, x, y, w, h, lines, anchor=MSO_ANCHOR.TOP, align=PP_ALIGN.LEFT, wrap=True):
    """加文字框。lines = [(text, size_pt, color_hex, bold, space_after_pt), ...]
    text 內含 \\n 會自動斷行(同一段落內)。"""
    tb = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = wrap; tf.vertical_anchor = anchor
    for m in ('margin_left', 'margin_right', 'margin_top', 'margin_bottom'):
        setattr(tf, m, 0)
    for i, (txt, size, color, bold, sa) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_before = Pt(0); p.space_after = Pt(sa)
        parts = txt.split('\n')
        for j, seg in enumerate(parts):
            if j > 0:
                br = p._p.makeelement(qn('a:br'), {}); p._p.append(br)
            r = p.add_run(); r.text = seg; set_run_font(r, size, color, bold)
    return tb

def center_text_in(shape, lines):
    """在既有 shape 內置中文字。lines = [(text, size, color, bold), ...]"""
    tf = shape.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for m in ('margin_left', 'margin_right', 'margin_top', 'margin_bottom'):
        setattr(tf, m, Pt(2))
    for i, (txt, size, color, bold) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER; p.space_before = Pt(0); p.space_after = Pt(0)
        r = p.add_run(); r.text = txt; set_run_font(r, size, color, bold)

# ---------------- 視覺效果:漸層 + 陰影 ----------------
def grad(shape, c1, c2, angle=90):
    """線性漸層。angle: 0=左→右, 90=上→下, 45=左上→右下。"""
    shape.fill.gradient()
    try: shape.fill.gradient_angle = angle
    except Exception: pass
    st = shape.fill.gradient_stops
    st[0].position = 0.0; st[0].color.rgb = RGBColor.from_string(c1)
    st[1].position = 1.0; st[1].color.rgb = RGBColor.from_string(c2)

def add_shadow(shape, blur=70000, dist=32000, direction=5400000, alpha=70, color=NAVY):
    """加外部柔和陰影 (EMU: 914400/in; direction 60000/度, 5400000=正下方)。"""
    sp = shape._element.spPr
    for el in sp.findall(qn('a:effectLst')):
        sp.remove(el)
    eff = sp.makeelement(qn('a:effectLst'), {})
    shdw = eff.makeelement(qn('a:outerShdw'),
        {'blurRad': str(blur), 'dist': str(dist), 'dir': str(direction), 'rotWithShape': '0'})
    clr = shdw.makeelement(qn('a:srgbClr'), {'val': color})
    clr.append(clr.makeelement(qn('a:alpha'), {'val': str(alpha * 1000)}))
    shdw.append(clr); eff.append(shdw); sp.append(eff)

def box(sl, x, y, w, h, fill=None, g=None, line=None, lw=0.75, rounded=True,
        shadow=True, shape=None, sh_alpha=70, sh_blur=70000, sh_dist=32000, adj=0.045):
    """通用色塊。g=(c1,c2[,angle]) 用漸層;fill=hex 用實色;line=hex 加邊框。"""
    st = shape if shape else (MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE)
    s = sl.shapes.add_shape(st, Inches(x), Inches(y), Inches(w), Inches(h))
    if st == MSO_SHAPE.ROUNDED_RECTANGLE:
        try: s.adjustments[0] = adj
        except Exception: pass
    if g: grad(s, g[0], g[1], g[2] if len(g) > 2 else 90)
    elif fill: s.fill.solid(); s.fill.fore_color.rgb = RGBColor.from_string(fill)
    else: s.fill.background()
    if line: s.line.color.rgb = RGBColor.from_string(line); s.line.width = Pt(lw)
    else: s.line.fill.background()
    if shadow: add_shadow(s, alpha=sh_alpha, blur=sh_blur, dist=sh_dist)
    else: s.shadow.inherit = False
    return s

def circle(sl, x, y, d, txt, tsize, g=(BLUE, BLUE2), tcolor=WHITE):
    """漸層圓形圖示(放編號或短字)。"""
    c = sl.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    grad(c, g[0], g[1], 45); c.line.fill.background()
    add_shadow(c, alpha=55, blur=45000, dist=22000)
    center_text_in(c, [(txt, tsize, tcolor, True)])
    return c

# ---------------- 高階版面元件 ----------------
def header(sl, eyebrow, title, subtitle, badge=None):
    """頂部標題列 (navy→blue 漸層)。badge 會在右上放一個金色標籤(如「規劃中」)。"""
    box(sl, 0, 0, 10, 1.28, g=(NAVY, BLUE, 90), rounded=False, shadow=True, sh_alpha=55, sh_dist=26000)
    add_text(sl, 0.55, 0.17, 8.0, 0.26, [(eyebrow, 10.5, BLUE_L, True, 0)])
    add_text(sl, 0.55, 0.40, 7.8, 0.5, [(title, 26, WHITE, True, 0)])
    add_text(sl, 0.55, 0.93, 9.0, 0.3, [(subtitle, 12.5, "CBD9EE", False, 0)])
    if badge:
        bb = box(sl, 8.15, 0.42, 1.5, 0.42, g=(GOLD, GOLD2, 45), shadow=True,
                 sh_alpha=50, sh_blur=40000, sh_dist=20000, adj=0.3)
        center_text_in(bb, [(badge, 11, WHITE, True)])

def concept_band(sl, text, y=1.40, h=0.5):
    """一行式重點帶(淺藍底 + 左側藍色強調條)。"""
    box(sl, 0.3, y, 9.4, h, g=(SKY, SKY2, 90), line=BLUE_L, lw=1, shadow=True, sh_alpha=40, sh_dist=20000)
    box(sl, 0.3, y, 0.12, h, fill=BLUE, rounded=False, shadow=False)
    add_text(sl, 0.62, y, 8.95, h, [(text, 10.5, INK, False, 0)], anchor=MSO_ANCHOR.MIDDLE)
    return y + h

def feature_card(sl, x, y, w, h, num, title, desc, accent=False):
    """白底功能卡:編號圓 + 粗體標題 + 說明。accent=True 用金色編號圓(整頁限一次)。"""
    box(sl, x, y, w, h, fill=CARD, line=BORDER, lw=1, shadow=True)
    g = (GOLD, GOLD2) if accent else (BLUE, BLUE2)
    circle(sl, x + 0.25, y + 0.22, 0.42, num, 15, g=g)
    add_text(sl, x + 0.82, y + 0.26, w - 1.0, 0.3, [(title, 13, INK, True, 0)])
    add_text(sl, x + 0.25, y + 0.78, w - 0.5, h - 0.9, [(desc, 10.5, MUTE, False, 0)])

def metric_tile(sl, x, y, val, title, desc, accent=False, tile_w=1.18):
    """數字 callout:左側漸層數字塊 + 右側標題/說明。accent=True 用金色(整頁限一次)。"""
    g = (GOLD, GOLD2, 45) if accent else (BLUE, BLUE2, 45)
    t = box(sl, x, y, tile_w, 0.42, g=g, shadow=True, sh_alpha=50, sh_blur=40000, sh_dist=20000)
    center_text_in(t, [(val, 12.5, WHITE, True)])
    add_text(sl, x + tile_w + 0.11, y - 0.02, 3.0, 0.26, [(title, 11.5, INK, True, 0)])
    add_text(sl, x + tile_w + 0.11, y + 0.205, 3.0, 0.3, [(desc, 9, MUTE, False, 0)])

def right_arrow(sl, x, y, w=0.55, h=0.3):
    """流程用右箭頭(藍色漸層 + 陰影)。"""
    a = sl.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(h))
    grad(a, BLUE_L, BLUE, 0); a.line.fill.background(); add_shadow(a, alpha=40, blur=30000, dist=15000)
    try: a.adjustments[0] = 0.5; a.adjustments[1] = 0.55
    except Exception: pass
    return a

def db_cylinder(sl, x, y, w, h, lines, g=(BLUE, NAVY)):
    """資料庫圓柱 (MSO CAN)。lines = [(text, size, color, bold), ...]"""
    c = sl.shapes.add_shape(MSO_SHAPE.CAN, Inches(x), Inches(y), Inches(w), Inches(h))
    grad(c, g[0], g[1], 90); c.line.fill.background()
    add_shadow(c, alpha=50, blur=45000, dist=22000)
    center_text_in(c, lines)
    return c

def chevron_row(sl, stages, y=6.18, h=0.66, cw=3.0, gap=0.12):
    """流程/路線圖 V 形列,色彩由淺藍漸深(表現推進)。
    stages = [(主字, 副字), ...]"""
    n = len(stages)
    x0 = (10 - (n * cw + (n - 1) * gap)) / 2
    for i, (t, d) in enumerate(stages):
        cx = x0 + i * (cw + gap)
        f = i / (n - 1) if n > 1 else 0
        def L(a, b): return int(a + (b - a) * f)
        c1 = "%02X%02X%02X" % (L(0x6F, 0x1E), L(0xB0, 0x4F), L(0xE0, 0x96))
        c2 = "%02X%02X%02X" % (L(0x2E, 0x0B), L(0x78, 0x2D), L(0xC7, 0x5C))
        cv = sl.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(cx), Inches(y), Inches(cw), Inches(h))
        grad(cv, c1, c2, 90); cv.line.fill.background(); add_shadow(cv, alpha=45, blur=38000, dist=18000)
        try: cv.adjustments[0] = 0.18
        except Exception: pass
        lines = [(t, 12, WHITE, True)] + ([(d, 8.5, SKY, False)] if d else [])
        center_text_in(cv, lines)

def footer(sl, left, right="2026-06"):
    """頁尾說明(左)與日期(右)。"""
    add_text(sl, 0.3, 7.22, 6.5, 0.2, [(left, 8, MUTE, False, 0)])
    add_text(sl, 6.5, 7.22, 3.2, 0.2, [(right, 8, MUTE, False, 0)], align=PP_ALIGN.RIGHT)


# ==================== 版面組合器 (composers) =========================================
# 目的:把「算座標」這件 agent 手算很爛的事交給確定性函式 —— 吃結構化內容,
# 自動垂直切分內容區、每塊撐滿、欄/項目均分,密度與平衡由建構保證,排不出稀疏。
# agent 只需:把內容整理成 blocks 結構 → 呼叫 compose_slide(),不再手擺 x/y/w/h。
#
# 用法:
#   compose_slide(s, "副標", "主標", "一行重點",
#     blocks=[
#       {"type":"columns", "cols":[
#           {"title":"系統設計目標","items":[{"title":"報價快","desc":"3-4週→1週"}, ...]},
#           {"title":"預期效果","metrics":[{"value":"↓70%","label":"報價週期","desc":"...","accent":True}, ...]},
#       ]},
#       {"type":"band", "title":"系統怎麼運作 — 7 步驟",
#        "flow":[{"title":"建置","desc":"廠 master"}, ...], "caption":"一行說明…"},
#       {"type":"cards", "cards":[{"title":"可行性已驗證","desc":"…"}, {"title":"開發排程","desc":"…"}]},
#     ],
#     footer_left="出處 · 版本")
#
# block 類型:
#   columns — 左右並排面板(白底藍標題條),每欄 body 為 items / metrics / flow
#   band    — 全寬面板(藍標題條),內放 flow(chevron 自動撐滿寬)+ 選配 caption
#   cards   — 自動 grid 功能卡(填滿區塊寬高)
#   每個 block 可給 "weight" 微調佔高;不給則依類型預設。

_C_TOP, _C_BOT = 1.46, 7.04     # 內容區上下界(header 下緣 ~ footer 上緣)
_C_MX  = 0.3                     # 左右邊距
_C_GUT = 0.18                   # 欄間距
_C_BGAP = 0.16                  # block 垂直間距
_BAR_H = 0.42                   # 面板藍標題條高
_PAD   = 0.16                   # 面板內距

def _block_weight(b):
    t = b.get("type")
    if t == "columns":  return 3.0
    if t == "band":     return 1.35
    if t == "cards":    return 1.5
    return 1.0

def compose_slide(sl, eyebrow, title, subtitle, blocks, footer_left="", footer_right="2026-06", badge=None):
    """一次組好整頁:header + 依 blocks 垂直填滿內容區 + footer。"""
    header(sl, eyebrow, title, subtitle, badge=badge)
    blocks = blocks or []
    weights = [b.get("weight", _block_weight(b)) for b in blocks]
    tw = sum(weights) or 1
    usable = (_C_BOT - _C_TOP) - _C_BGAP * (len(blocks) - 1)
    W = 10 - 2 * _C_MX
    y = _C_TOP
    for b, wgt in zip(blocks, weights):
        h = usable * (wgt / tw)
        _render_block(sl, b, _C_MX, y, W, h)
        y += h + _C_BGAP
    footer(sl, footer_left, footer_right)

def _render_block(sl, b, x, y, w, h):
    t = b.get("type")
    if t == "columns":
        cols = b.get("cols", [])
        n = max(1, len(cols))
        cw = (w - _C_GUT * (n - 1)) / n
        for i, c in enumerate(cols):
            _render_panel(sl, c, x + i * (cw + _C_GUT), y, cw, h)
    elif t == "band":
        _render_band(sl, b, x, y, w, h)
    elif t == "cards":
        _render_cards(sl, b, x, y, w, h)

def _panel(sl, x, y, w, h, title):
    """白底面板 + 藍標題條,回傳內容區 (ix, iy, iw, ih)。"""
    box(sl, x, y, w, h, fill=CARD, line=BORDER, lw=1, shadow=True)
    if title:
        bar = box(sl, x, y, w, _BAR_H, g=(BLUE, BLUE2, 90), rounded=True, shadow=False, adj=0.1)
        center_text_in(bar, [(title, 12.5, WHITE, True)])
        return (x + _PAD, y + _BAR_H + _PAD, w - 2 * _PAD, h - _BAR_H - 2 * _PAD)
    return (x + _PAD, y + _PAD, w - 2 * _PAD, h - 2 * _PAD)

def _render_panel(sl, c, x, y, w, h):
    ix, iy, iw, ih = _panel(sl, x, y, w, h, c.get("title", ""))
    if c.get("metrics"):  _fill_metrics(sl, ix, iy, iw, ih, c["metrics"])
    elif c.get("flow"):   _chevron_fill(sl, ix, iy, iw, ih, c["flow"])
    else:                 _fill_items(sl, ix, iy, iw, ih, c.get("items", []))

def _fill_items(sl, ix, iy, iw, ih, items):
    """編號圓 + 標題 + 小字,均分撐滿 ih(每項在自己的 row 內垂直置中)。"""
    n = max(1, len(items)); pitch = ih / n
    voff = max(0.0, (pitch - 0.5) / 2); cd = 0.34
    for i, it in enumerate(items):
        ry = iy + i * pitch + voff
        g = (GOLD, GOLD2) if it.get("accent") else (BLUE, BLUE2)
        circle(sl, ix, ry, cd, str(i + 1), 14, g=g)
        tx = ix + cd + 0.16; tw = iw - cd - 0.16
        add_text(sl, tx, ry - 0.02, tw, 0.26, [(it.get("title", ""), 11.5, INK, True, 0)])
        if it.get("desc"):
            add_text(sl, tx, ry + 0.23, tw, 0.3, [(it["desc"], 9, MUTE, False, 0)])

def _fill_metrics(sl, ix, iy, iw, ih, metrics):
    """數字塊 + 標題 + 小字,均分撐滿 ih。accent=True 用金色(整份限一次)。"""
    n = max(1, len(metrics)); pitch = ih / n
    voff = max(0.0, (pitch - 0.5) / 2); tile_w = 1.12
    for i, m in enumerate(metrics):
        ry = iy + i * pitch + voff
        g = (GOLD, GOLD2, 45) if m.get("accent") else (BLUE, BLUE2, 45)
        tile = box(sl, ix, ry, tile_w, 0.46, g=g, shadow=True, sh_alpha=50, sh_blur=40000, sh_dist=20000)
        center_text_in(tile, [(str(m.get("value", "")), 13, WHITE, True)])
        tx = ix + tile_w + 0.16; rw = iw - tile_w - 0.16
        add_text(sl, tx, ry - 0.02, rw, 0.26, [(m.get("label", ""), 11.5, INK, True, 0)])
        if m.get("desc"):
            add_text(sl, tx, ry + 0.24, rw, 0.3, [(m["desc"], 9, MUTE, False, 0)])

def _chevron_fill(sl, ix, iy, iw, ih, steps):
    """chevron V 形列,自動算 cw 撐滿 iw、垂直置中於 ih。steps=[{title,desc} 或 字串]。"""
    n = max(1, len(steps)); gap = 0.06
    cw = (iw - gap * (n - 1)) / n
    ch = min(0.72, ih); cy = iy + (ih - ch) / 2
    for i, st in enumerate(steps):
        cx = ix + i * (cw + gap)
        f = i / (n - 1) if n > 1 else 0
        def L(a, b): return int(a + (b - a) * f)
        c1 = "%02X%02X%02X" % (L(0x6F, 0x1E), L(0xB0, 0x4F), L(0xE0, 0x96))
        c2 = "%02X%02X%02X" % (L(0x2E, 0x0B), L(0x78, 0x2D), L(0xC7, 0x5C))
        cv = sl.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(cx), Inches(cy), Inches(cw), Inches(ch))
        grad(cv, c1, c2, 90); cv.line.fill.background(); add_shadow(cv, alpha=45, blur=38000, dist=18000)
        try: cv.adjustments[0] = 0.16
        except Exception: pass
        t = st.get("title") if isinstance(st, dict) else st
        d = st.get("desc") if isinstance(st, dict) else ""
        lines = [(str(t or ""), 11, WHITE, True)] + ([(d, 8, SKY, False)] if d else [])
        center_text_in(cv, lines)

def _render_band(sl, b, x, y, w, h):
    ix, iy, iw, ih = _panel(sl, x, y, w, h, b.get("title", ""))
    cap = b.get("caption")
    flow_h = ih - (0.26 if cap else 0.0)
    if b.get("flow"):
        _chevron_fill(sl, ix, iy, iw, flow_h, b["flow"])
    if cap:
        add_text(sl, ix, iy + flow_h, iw, 0.24, [(cap, 9, MUTE, False, 0)], align=PP_ALIGN.CENTER)

def _render_cards(sl, b, x, y, w, h):
    cards = b.get("cards", [])
    n = max(1, len(cards))
    cols = b.get("cols", 2 if n > 1 else 1)
    rows = (n + cols - 1) // cols
    gx, gy = 0.18, 0.16
    cw = (w - gx * (cols - 1)) / cols
    chh = (h - gy * (rows - 1)) / rows
    for i, c in enumerate(cards):
        r, cc = divmod(i, cols)
        feature_card(sl, x + cc * (cw + gx), y + r * (chh + gy), cw, chh,
                     c.get("num", chr(65 + i)), c.get("title", ""), c.get("desc", ""),
                     accent=c.get("accent", False))
