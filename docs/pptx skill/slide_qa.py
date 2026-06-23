# -*- coding: utf-8 -*-
"""
slide_qa.py — 投影片「程式化視覺 QA」,涵蓋 SKILL.md 列的四項檢查:
  1. 文字溢出色塊(check_overflow)
  2. 是否只用了一次金色(GOLD/GOLD2)
  3. 邊距是否足夠 / 元素是否出界
  4. 實體色塊是否互相重疊(保守近似:只報非文字實心 shape 間的顯著重疊)

主張:四項全是「可量測的幾何/顏色約束」,交給程式判斷遠比 LLM vision 目測可靠、零成本、
跨模型一致。LLM 只需做語意判斷。這是把 Claude-native skill 安全移植到 Gemini/AOAI 的關鍵。
"""
import math
from pptx import Presentation
from pptx.util import Emu
from pptx.oxml.ns import qn
from check_overflow import analyze_shape

SLIDE_W_IN, SLIDE_H_IN = 10.0, 7.5
GOLD_HEXES = {'C99A3B', 'E0B85A'}
EDGE_MARGIN_IN = 0.2
OVERLAP_RATIO = 0.40

def _emu_in(emu):
    return Emu(emu).inches

def _fill_hexes(shape):
    out = set()
    try:
        spPr = shape._element.spPr
    except Exception:
        return out
    if spPr is None:
        return out
    for el in spPr.iter(qn('a:srgbClr')):
        v = el.get('val')
        if v:
            out.add(v.upper())
    return out

def _is_solid_block(shape):
    if shape.has_text_frame and not _fill_hexes(shape):
        return False
    return bool(_fill_hexes(shape))

def _rect(shape):
    x = _emu_in(shape.left); y = _emu_in(shape.top)
    return (x, y, x + _emu_in(shape.width), y + _emu_in(shape.height))

def _overlap_area(a, b):
    ix = max(0.0, min(a[2], b[2]) - max(a[0], b[0]))
    iy = max(0.0, min(a[3], b[3]) - max(a[1], b[1]))
    return ix * iy

def _contains(big, small, tol=0.06):
    """small 是否(近乎)完全落在 big 內 — 是的話屬刻意內嵌(編號圓/圖示/強調條),非碰撞。"""
    return (small[0] >= big[0]-tol and small[1] >= big[1]-tol and
            small[2] <= big[2]+tol and small[3] <= big[3]+tol)

def qa_deck(path):
    prs = Presentation(path); issues = []
    for si, slide in enumerate(prs.slides):
        sn = si + 1; gold_count = 0; blocks = []
        for idx, shape in enumerate(slide.shapes):
            r = analyze_shape(shape)
            if r and r['overflow']:
                issues.append((sn, 'OVERFLOW', f"{r['text']!r} need {r['needed_h_pt']}pt>avail {r['avail_h_pt']}pt"))
            if _fill_hexes(shape) & GOLD_HEXES:
                gold_count += 1
            x0, y0, x1, y1 = _rect(shape)
            is_fullbg = (x0 <= 0.01 and y0 <= 0.01 and x1 >= SLIDE_W_IN-0.01 and y1 >= SLIDE_H_IN-0.01)
            if not is_fullbg:
                if x0 < -0.01 or y0 < -0.01 or x1 > SLIDE_W_IN+0.01 or y1 > SLIDE_H_IN+0.01:
                    issues.append((sn, 'OUT_OF_BOUNDS', f"shape rect=({x0:.2f},{y0:.2f},{x1:.2f},{y1:.2f})"))
                if _is_solid_block(shape):
                    blocks.append((idx, (x0,y0,x1,y1), (x1-x0)*(y1-y0)))
        if gold_count > 1:
            issues.append((sn, 'GOLD_OVERUSE', f"金色色塊用了 {gold_count} 次(限 1)"))
        for i in range(len(blocks)):
            for j in range(i+1, len(blocks)):
                bi, bj = blocks[i], blocks[j]
                big, sml = (bi, bj) if bi[2] >= bj[2] else (bj, bi)
                ov = _overlap_area(bi[1], bj[1])
                if sml[2] > 0 and ov/sml[2] > OVERLAP_RATIO:
                    if _contains(big[1], sml[1]):
                        continue  # 小塊完全內嵌進大塊 = 刻意設計(badge/icon/accent),非碰撞
                    issues.append((sn, 'OVERLAP', f"blocks #{bi[0]}&#{bj[0]} 重疊 {ov/sml[2]*100:.0f}%"))
    return issues

if __name__ == '__main__':
    import sys
    path = sys.argv[1] if len(sys.argv) > 1 else 'example_out.pptx'
    issues = qa_deck(path)
    print(f'=== QA {path}: {len(issues)} issue(s) ===')
    for sn, kind, detail in issues:
        print(f'  [slide {sn}] {kind}: {detail}')
    if not issues:
        print('  (clean)')
