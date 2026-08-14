# -*- coding: utf-8 -*-
"""
Cortex BOM × MVA — 董事長單頁簡報
格式:全螢幕 4:3 · 寬 25.4 cm × 高 19.05 cm (10 x 7.5 inches)
重點:設計目標(業務快接接單) + 期待效果 + 系統功能 7-Phase 端到端
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

prs = Presentation()
prs.slide_width = Inches(10)      # 25.4 cm
prs.slide_height = Inches(7.5)    # 19.05 cm

# Palette
INK    = RGBColor(0x0B, 0x1F, 0x3A)
TEXT   = RGBColor(0x37, 0x47, 0x55)
MUTED  = RGBColor(0x6B, 0x72, 0x80)
OCEAN  = RGBColor(0x02, 0xC3, 0x9A)
NAVY   = RGBColor(0x1C, 0x72, 0x93)
LINE   = RGBColor(0xE5, 0xE7, 0xEB)
BG     = RGBColor(0xFA, 0xFB, 0xFC)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GOLD   = RGBColor(0xCA, 0x8A, 0x04)
RED    = RGBColor(0xDC, 0x26, 0x26)
GREEN  = RGBColor(0x16, 0xA3, 0x4A)

PHASE_COLORS = [
    RGBColor(0x7C, 0x3A, 0xED),
    RGBColor(0x0E, 0xA5, 0xE9),
    RGBColor(0x16, 0xA3, 0x4A),
    RGBColor(0x08, 0x91, 0xB2),
    RGBColor(0xDC, 0x26, 0x26),
    RGBColor(0xB4, 0x53, 0x09),
    RGBColor(0x93, 0x33, 0xEA),
]


def add_text(slide, x, y, w, h, text, size=12, bold=False, color=None,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font_name='Microsoft JhengHei'):
    tx = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tx.text_frame
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font_name
    r.font.size = Pt(size)
    r.font.bold = bold
    if color:
        r.font.color.rgb = color
    return tx


def add_rect(slide, x, y, w, h, fill, line=None, line_width=0.5):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(x), Inches(y), Inches(w), Inches(h))
    rect.fill.solid()
    rect.fill.fore_color.rgb = fill
    if line is None:
        rect.line.fill.background()
    else:
        rect.line.color.rgb = line
        rect.line.width = Pt(line_width)
    rect.shadow.inherit = False
    return rect


def add_rounded(slide, x, y, w, h, fill, line=None, line_width=0.5, radius=0.05):
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(x), Inches(y), Inches(w), Inches(h))
    rect.adjustments[0] = radius
    rect.fill.solid()
    rect.fill.fore_color.rgb = fill
    if line is None:
        rect.line.fill.background()
    else:
        rect.line.color.rgb = line
        rect.line.width = Pt(line_width)
    rect.shadow.inherit = False
    return rect


def add_chevron(slide, x, y, w, h, fill, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.CHEVRON,
                                Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
    sh.shadow.inherit = False
    return sh


# ═══ Single slide ═══════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])

# 1. 整頁底色
add_rect(s, 0, 0, 10, 7.5, BG)

# 2. 頂部 banner(深色 + 口號)
add_rect(s, 0, 0, 10, 1.35, INK)
# 左側 logo bar
add_rect(s, 0, 0, 0.15, 1.35, OCEAN)

# Eyebrow
add_text(s, 0.45, 0.18, 9.3, 0.28, 'FOXLINK 正崴 · Cortex AI 整合平台', size=10, color=OCEAN)
# Main title
add_text(s, 0.45, 0.42, 9.3, 0.5, 'Cortex BOM × MVA  ▏  業務快接接單平台',
         size=22, bold=True, color=WHITE)
# Tagline
add_text(s, 0.45, 0.95, 9.3, 0.32, '報價週期 3-4 週 → 1 週 ‧ 跨案重用知識資產 ‧ Margin 即時透明 ‧ EPM 走人不重做',
         size=11, color=RGBColor(0xCB, 0xD5, 0xE1))

# ════════════════════════════════════════════════════════════
# 3. 中段:左「設計目標」+ 右「期待效果」
# ════════════════════════════════════════════════════════════
mid_y = 1.55
mid_h = 2.9

# --- 左:設計目標 ---
add_rounded(s, 0.3, mid_y, 4.7, mid_h, WHITE, line=NAVY, line_width=1.2, radius=0.03)
add_rounded(s, 0.3, mid_y, 4.7, 0.4, NAVY, radius=0.03)
add_text(s, 0.5, mid_y + 0.06, 4.5, 0.3, '🎯  系統設計目標 (Why)',
         size=12, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)

goals = [
    ('⚡', '業務快接接單', '報價週期從 3-4 週壓到 1 週 · 同步出 6 cells 多廠多色多批量報價'),
    ('🔄', '跨案 / 跨廠重用', '廠級 master(設備/耗材/IDL/wage)結構化 · EPM 走人不重做 1500 行 Excel'),
    ('📊', 'Margin 即時透明', 'True Cost / Quote Price 雙價分層 · 議價當下看到毛利浮動'),
    ('🔒', '雙價 + 4 級權限', '業務只看 Quote · 採購/CFO 看 True · 機密欄 VIEW_TRUE_COST 控管'),
]
for i, (ico, t, d) in enumerate(goals):
    y = mid_y + 0.55 + i * 0.58
    add_text(s, 0.5, y, 0.4, 0.4, ico, size=18, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, 0.95, y, 4.0, 0.25, t, size=11, bold=True, color=INK)
    add_text(s, 0.95, y + 0.25, 4.0, 0.32, d, size=8.5, color=TEXT)

# --- 右:期待效果 ---
add_rounded(s, 5.05, mid_y, 4.65, mid_h, WHITE, line=GREEN, line_width=1.2, radius=0.03)
add_rounded(s, 5.05, mid_y, 4.65, 0.4, GREEN, radius=0.03)
add_text(s, 5.25, mid_y + 0.06, 4.45, 0.3, '📈  期待達成效果 (Impact)',
         size=12, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)

impacts = [
    ('↓ 70%', '報價週期', '報價交件:3-4 週 → 1 週', RED),
    ('↑ 3 廠', '同步並行 RFQ', 'CN / VN / TW 同時跑 Cleansheet,不再 EPM 個人 Excel 各自寫', NAVY),
    ('24 cells', '一鍵多視角', '3 廠 × 2 色 × 2 批量 × 2 包裝 = 24 種報價 pivot 即看', OCEAN),
    ('0 漏單', 'Lock 自動 propagate', 'BOM 鎖版 → MVA × SGA × Profit 全自動寫 fact table', GOLD),
]
for i, (kpi, t, d, c) in enumerate(impacts):
    y = mid_y + 0.55 + i * 0.58
    add_rounded(s, 5.25, y, 0.95, 0.5, c, radius=0.1)
    add_text(s, 5.25, y, 0.95, 0.5, kpi, size=11, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, 6.3, y, 3.3, 0.25, t, size=11, bold=True, color=INK)
    add_text(s, 6.3, y + 0.25, 3.3, 0.32, d, size=8.5, color=TEXT)

# ════════════════════════════════════════════════════════════
# 4. 中段下:7-Phase 端到端 pipeline
# ════════════════════════════════════════════════════════════
pipe_y = 4.7
add_rounded(s, 0.3, pipe_y, 9.4, 1.7, WHITE, line=GOLD, line_width=1.2, radius=0.02)
add_rounded(s, 0.3, pipe_y, 9.4, 0.4, GOLD, radius=0.02)
add_text(s, 0.5, pipe_y + 0.06, 9.0, 0.3, '🛠️  系統功能 — 端到端 7 Phase × 33 步驟 (What)',
         size=12, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)

phases = [
    ('A', '建置', '一次性\n廠 master'),
    ('B', '月維', '月度\nbaseline'),
    ('C', '開案', '選模板\n設 qty/pkg'),
    ('D', '配置', '案級\n9×17 matrix'),
    ('E', '算', '30 秒\ncompute'),
    ('F', '鎖', 'DPM lock\npropagate'),
    ('G', '報', 'Margin\n+ 業務 gate'),
]

chev_y = pipe_y + 0.55
chev_w = 1.2
gap = 0.05
start_x = 0.55
for i, (code, name, desc) in enumerate(phases):
    x = start_x + i * (chev_w + gap)
    color = PHASE_COLORS[i]
    add_chevron(s, x, chev_y, chev_w, 0.65, color)
    add_text(s, x, chev_y, chev_w - 0.18, 0.3, code,
             size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x, chev_y + 0.35, chev_w - 0.18, 0.28, name,
             size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # description below
    add_text(s, x - 0.05, chev_y + 0.75, chev_w + 0.05, 0.6, desc,
             size=8, color=TEXT, align=PP_ALIGN.CENTER)

# Bottom summary line under pipeline
add_text(s, 0.5, pipe_y + 1.4, 9.0, 0.25,
         '⇨ 廠級 SCD Type 2 baseline ‧ 案級 24 cells fact table ‧ Excel→DB 結構化(100+ 條公式對應 · ε<0.01)',
         size=9, color=MUTED, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════
# 5. 底部:兩案範本驗證 + 上線時程
# ════════════════════════════════════════════════════════════
bot_y = 6.55

# 兩案 badge
add_rounded(s, 0.3, bot_y, 4.7, 0.75, RGBColor(0xF0, 0xF9, 0xFF), line=RGBColor(0x03, 0x69, 0xA1), line_width=1, radius=0.05)
add_text(s, 0.5, bot_y + 0.05, 4.5, 0.28, '✓ 已用兩案範本驗證可行性', size=10, bold=True, color=RGBColor(0x03, 0x69, 0xA1))
add_text(s, 0.5, bot_y + 0.35, 4.5, 0.35,
         '🖱️ SteelSeries Rival 3+ Wired (MOUSE_STD · 9 製程) ‧ ⌚ WHOOP Gen4 MP (WEARABLE · 8 sub-BOM)',
         size=8.5, color=TEXT)

# 時程 badge
add_rounded(s, 5.05, bot_y, 4.65, 0.75, RGBColor(0xFE, 0xF3, 0xC7), line=GOLD, line_width=1, radius=0.05)
add_text(s, 5.25, bot_y + 0.05, 4.45, 0.28, '📅 開發排程', size=10, bold=True, color=GOLD)
add_text(s, 5.25, bot_y + 0.35, 4.45, 0.35,
         '6.5 sprint ≈ 5-6 週 / 1 人 · Phase 1-4 核心 (基礎+引擎+UI+propagate) · Phase 5-6 跨廠擴充',
         size=8.5, color=TEXT)

# ════════════════════════════════════════════════════════════
# 6. 最底細線 footer
# ════════════════════════════════════════════════════════════
add_rect(s, 0.3, 7.36, 9.4, 0.02, LINE)
add_text(s, 0.3, 7.4, 6, 0.18, 'Cortex BOM × MVA 模組 · v0.9 demo · 對應 cleansheet-mva-sd v0.3 + bom-collection-sd v0.4',
         size=7.5, color=MUTED)
add_text(s, 7.5, 7.4, 2.2, 0.18, '2026-06 · 內部呈核版',
         size=7.5, color=MUTED, align=PP_ALIGN.RIGHT)

OUT = 'd:/vibe_coding/foxlink_gpt/docs/Cortex_BOM_MVA_董事長簡報_v1.pptx'
prs.save(OUT)
print(f'[OK] Saved: {OUT}')
print(f'     Slide size: {prs.slide_width/914400:.2f}in x {prs.slide_height/914400:.2f}in (4:3 fullscreen)')
print(f'     [~] {prs.slide_width/360000:.2f}cm x {prs.slide_height/360000:.2f}cm')
