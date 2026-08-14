# -*- coding: utf-8 -*-
"""
Cortex MVA 操作流程說明手冊 v1
跨 SteelSeries (MOUSE_STD) + WHOOP (WEARABLE) 兩案整合
~32 slides · Phase A-G × 33 steps · 內部討論用
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ─── Slide size 16:9 ───────────────────────────────────
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ─── Palette (Foxlink Cortex) ──────────────────────────
INK = RGBColor(0x0B, 0x1F, 0x3A)
TEXT = RGBColor(0x37, 0x47, 0x55)
MUTED = RGBColor(0x6B, 0x72, 0x80)
OCEAN = RGBColor(0x02, 0xC3, 0x9A)
NAVY = RGBColor(0x1C, 0x72, 0x93)
LINE = RGBColor(0xE5, 0xE7, 0xEB)
BG_SOFT = RGBColor(0xFA, 0xFB, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# Phase colors (對齊 HTML)
PHASE_COLORS = {
    'A': RGBColor(0x7C, 0x3A, 0xED),   # purple
    'B': RGBColor(0x0E, 0xA5, 0xE9),   # sky
    'C': RGBColor(0x16, 0xA3, 0x4A),   # green
    'D': RGBColor(0x08, 0x91, 0xB2),   # cyan
    'E': RGBColor(0xDC, 0x26, 0x26),   # red
    'F': RGBColor(0xB4, 0x53, 0x09),   # amber
    'G': RGBColor(0x93, 0x33, 0xEA),   # violet
}
STEEL_COLOR = RGBColor(0x03, 0x69, 0xA1)
WHOOP_COLOR = RGBColor(0x6B, 0x21, 0xA8)
STEEL_BG = RGBColor(0xF0, 0xF9, 0xFF)
WHOOP_BG = RGBColor(0xFA, 0xF5, 0xFF)


def add_blank_slide():
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_text(slide, x, y, w, h, text, size=14, bold=False, color=None,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font_name='Microsoft JhengHei'):
    tx = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tx.text_frame
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font_name
    r.font.size = Pt(size)
    r.font.bold = bold
    if color:
        r.font.color.rgb = color
    return tx


def add_rect(slide, x, y, w, h, fill, line=None, line_width=0.75):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(x), Inches(y), Inches(w), Inches(h))
    rect.fill.solid()
    rect.fill.fore_color.rgb = fill
    if line is None:
        rect.line.fill.background()
    else:
        rect.line.color.rgb = line
        rect.line.width = Pt(line_width)
    rect.shadow.inherit = False
    return rect


def add_rounded(slide, x, y, w, h, fill, line=None, line_width=0.75, radius=0.05):
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(x), Inches(y), Inches(w), Inches(h))
    rect.adjustments[0] = radius
    rect.fill.solid()
    rect.fill.fore_color.rgb = fill
    if line is None:
        rect.line.fill.background()
    else:
        rect.line.color.rgb = line
        rect.line.width = Pt(line_width)
    rect.shadow.inherit = False
    return rect


def add_pill(slide, x, y, text, fill, fg, size=9):
    """Small pill label."""
    w = max(0.5, len(text) * 0.085 + 0.15)
    pill = add_rounded(slide, x, y, w, 0.22, fill, radius=0.5)
    add_text(slide, x + 0.04, y + 0.025, w - 0.08, 0.18, text,
             size=size, bold=True, color=fg, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return w


def slide_header(slide, title, subtitle=None, phase_code=None):
    """Standard slide header bar."""
    # Top accent bar
    color = PHASE_COLORS.get(phase_code, NAVY) if phase_code else NAVY
    add_rect(slide, 0, 0, 13.333, 0.16, color)
    # Title
    add_text(slide, 0.4, 0.25, 11, 0.5, title, size=22, bold=True, color=INK)
    if subtitle:
        add_text(slide, 0.4, 0.75, 11, 0.32, subtitle, size=11, color=MUTED)
    # Phase code badge (top-right)
    if phase_code:
        add_rounded(slide, 12.1, 0.32, 0.85, 0.42, color, radius=0.3)
        add_text(slide, 12.1, 0.33, 0.85, 0.4, f'Phase {phase_code}',
                 size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def slide_footer(slide, n, total, side_text='Cortex MVA 操作流程 v1'):
    add_text(slide, 0.4, 7.1, 6, 0.3, side_text, size=9, color=MUTED)
    add_text(slide, 12.0, 7.1, 1, 0.3, f'{n} / {total}', size=9, color=MUTED, align=PP_ALIGN.RIGHT)
    add_rect(slide, 0.4, 7.05, 12.55, 0.02, LINE)


# ════════════════════════════════════════════════════════════════════════
# PHASE DATA(同 HTML MVA_PHASES)
# ════════════════════════════════════════════════════════════════════════

PHASES = [
    {
        'code': 'A', 'title': '公司 / 廠級一次性建置',
        'timing': '初始 · 一次性 · ~2 週/廠',
        'owner': 'IT + Admin EPM(per 廠)',
        'summary': '建立公司級主檔(製程目錄 / 模板)+ 各廠首份 baseline(SupplierBaseInput / IDL 17 角色 / 設備庫 / 耗材庫)。完成後其他案才能開。',
        'prereq': [
            ('📄 文件', '廠 EPM Cleansheet xlsx 原檔', '本 demo: Rival 3+ Wired Mouse Cleansheet_China-20241011-19S.xlsx'),
            ('📄 文件', '設備折舊年限政策', '各分類年限 SMT / SMT_TEST / BB_ASSY / BB_TEST / IT'),
            ('📊 資料', 'IDL 17 種薪資表', 'Ops Mgr / Section Mgr / Engineer / Tech 等 17 角色'),
            ('📊 資料', '廠別 master', 'CN/VN/TW 幣別 + ledger book_id (對齊 ERP)'),
            ('📐 政策', '9 製程目錄 enum', 'SMT/Wave/Router/Laser/BB Assy/BB Test/Mat Mgmt/Q-SMT/Q-BB'),
            ('📐 政策', '產品模板初始(2 種)', 'MOUSE_STD + WHOOP_WEARABLE'),
        ],
        'steps': [
            ('A1', '廠基本資料建檔', 'Admin', 'INSERT bom_factory', 'CN/VN/TW 各一列', 30),
            ('A2', '9 製程目錄 enum 初始化', 'Admin', 'INSERT bom_process_catalog', '依 SD §3.2 · 10 列(含 COMMON)· 三語言名稱', 60),
            ('A3', '2 產品模板(含 step list)', 'Admin', 'INSERT bom_process_template + ..._step',
                'MOUSE: SMT→Wave→Router→Laser→BB Assy→Test→Mat→Q · WHOOP: Multi-Module SMT→Module Integ→Final Assy→FATP 7', 90),
            ('A4', '廠 EPM 上傳首份 Cleansheet xlsx', 'CN EPM Andy', 'multer upload → stage',
                'SupplierBaseInput / Equipment / Consumables / Cleansheet / Price Summary 5 sheet 必備', 10),
            ('A5', '系統解析 → 寫 baseline + IDL + 折舊年限', '系統', 'bom_cs_baseline_import.py',
                'INSERT bom_factory_baseline (status=draft) + bom_factory_idl_role × 17 + bom_factory_dep_years × N', 5),
            ('A6', '系統解析 Equipment List → upsert 廠設備庫', '系統', 'INSERT/UPDATE bom_factory_equipment',
                '58 件設備 · bucket/type/desc/manuf/model/pn/useful_life/acq_cost', 5),
            ('A7', '系統解析 Consumables → 廠耗材庫', '系統', 'INSERT bom_factory_consumable',
                '17 種耗材 · unit_cost / annual_usage / UOM', 5),
            ('A8', 'Admin approve baseline', 'Admin', 'UPDATE baseline status=active + effective_from',
                '設 effective_from = today · 該廠首份 active baseline 鎖定', 15),
        ],
        'delta_steel': 'CN baseline: DL $4.95/hr · 58 equipment · 17 consumable · 17 IDL roles · 9 製程',
        'delta_whoop': '同廠 baseline 共用(不重建)· template 用 WHOOP_WEARABLE · FATP 7 station 設備庫獨立加 7 件',
        'schemas': ['bom_factory', 'bom_factory_baseline', 'bom_factory_idl_role', 'bom_factory_dep_years',
                    'bom_factory_equipment', 'bom_factory_consumable',
                    'bom_process_catalog', 'bom_process_template', 'bom_process_template_step'],
    },
    {
        'code': 'B', 'title': '月度 Baseline 維護',
        'timing': '每月底 · ~1 天/廠',
        'owner': '各廠 EPM + admin 簽核',
        'summary': '廠級 wage / 設備購置 / 耗材單價 月度異動 · SCD Type 2 切版 · 不影響已 lock 案,僅 draft 案 EPM 自選 reprice。',
        'prereq': [
            ('📄 文件', '新月 Cleansheet xlsx', 'EPM 本月更新版 · 主異動:SupplierBaseInput + Equipment + Consumables'),
            ('📊 資料', '現行 active baseline', 'SELECT WHERE factory_code=? AND status=active'),
            ('📐 政策', 'Diff 比對策略', 'wage 漲>5% / 設備 acq>10% / 折舊年限改 → EPM+廠財雙簽'),
        ],
        'steps': [
            ('B1', 'EPM 上傳新版 xlsx', '廠 EPM', 'POST /api/cs/baselines/import', '建 draft baseline · 不影響 active', 5),
            ('B2', '系統解析 + Diff 比對', '系統', 'bom_cs_baseline_diff 寫入', 'wage/IDL/設備/耗材/折舊 5 維度 · ADDED/CHANGED/RETIRED', 3),
            ('B3', 'Diff Preview UI', 'EPM', 'GET /api/cs/baselines/:id/diff',
                '卡片式: DL $4.95→$5.20(+5.05%) · DEK acq $84k→$92k · 新增 FUJI M6S · 汰換 TR7500', 20),
            ('B4', 'EPM 逐項勾選 approve / reject', 'EPM', 'PATCH /api/cs/baselines/:id/items',
                '單獨勾選 · 沒勾不套用 · 漲幅>5% 系統高亮需二簽', 15),
            ('B5', 'admin / 廠財務 sign-off(漲幅大者)', '廠財務+admin', 'audit_log entry',
                '漲幅>5% block 直到二簽', 30),
            ('B6', 'SCD Type 2 切版', '系統', 'UPDATE old effective_to + status=superseded · UPDATE new status=active',
                'today/today+1 切割', 1),
            ('B7', '影響案 notify', '系統', 'user_notifications', '推 chip 給所有 draft case_factory', 1),
        ],
        'delta_steel': '本案 draft → EPM 收 chip · 可選 reprice 或保留(baseline_locked_manually=1)',
        'delta_whoop': '同 SteelSeries · 但 reprice 觸發 8 sub-BOM × 24 cells 重算 · ~3x 慢',
        'schemas': ['bom_factory_baseline', 'bom_cs_baseline_diff', 'bom_cs_audit_log', 'user_notifications'],
    },
    {
        'code': 'C', 'title': '新案開立 (Case Setup)',
        'timing': '案啟動 · 1 天',
        'owner': 'BPM + 三廠 EPM',
        'summary': '從 project 出發建 case_factory · 選 process_template · 配 qty_scenarios + pkg_versions + variants(色款)。此 phase 結束才能進案級配置。',
        'prereq': [
            ('📊 資料', 'project 主表已建', 'customer_id / BPN / annual_demand / 期望出貨地'),
            ('📊 資料', 'variants 清單', 'SteelSeries: black/white · WHOOP: 3 色'),
            ('📊 資料', 'qty_scenarios 拍板', 'Low/High(本案 100K/418K)+ region 適用範圍'),
            ('📊 資料', 'pkg_versions 拍板', 'SteelSeries: 商包+工包 · WHOOP: 5 PKG SKU × 8 module'),
            ('📐 政策', '廠別配置範圍', '3 廠都跑 / 只 CN+VN / 只 TW · BPM 決定 RFQ scope'),
        ],
        'steps': [
            ('C1', 'BPM 開 project + customer/BPN', 'BPM Lisa', 'INSERT projects', '§customer 填完 → ACTIVE', 20),
            ('C2', 'RD 上傳 BOM xlsx → 抽 variants', 'RD Troy', '§BOM 流程',
                'AI 解析 EE/ME BOM · §variants 標 shared/per_variant scope', 60),
            ('C3', 'BPM 建 case_factory × 3 廠', 'BPM', 'INSERT bom_cs_case_factory × 3',
                'baseline_id 鎖當前 active · status=draft', 10),
            ('C4', '選 process_template', 'BPM + EPM', 'UPDATE bom_cs_case_factory.process_template_id',
                'MOUSE → MOUSE_STD · 可穿戴 → WHOOP_WEARABLE · 自訂 → CUSTOM', 5),
            ('C5', '設 qty_scenarios', 'BPM', 'INSERT bom_cs_case_qty_scenario',
                '本 demo: LOW=100K(US) · HIGH=418K(global) · is_baseline=HIGH', 5),
            ('C6', '設 pkg_versions', 'BPM + PKG 採購 Ken', 'INSERT bom_cs_case_pkg',
                'RETAIL=商包(retail+online) · BULK=工包(spare+b2b)', 10),
        ],
        'delta_steel': 'MOUSE_STD · 1-BOM · 3 廠 × 2 variant × 2 qty × 2 pkg = 24 cells / case',
        'delta_whoop': 'WHOOP_WEARABLE · 8 sub-BOM (bom_module) · 5 PKG SKU × 8 module 矩陣 · FATP 7 station',
        'schemas': ['bom_cs_case_factory', 'bom_cs_case_qty_scenario', 'bom_cs_case_pkg', 'projects'],
    },
    {
        'code': 'D', 'title': '案級配置 (Case-Level Config)',
        'timing': '案級填單 · 3-5 天 / 廠',
        'owner': '各廠 EPM(自家廠)',
        'summary': '從模板帶入製程清單 → 各製程調 9 個變數 · IDL × 製程 multiplier matrix · 設備與耗材 binding(line_qty/qty_per_line)· PKG items 上傳。',
        'prereq': [
            ('📊 資料', '廠 master 已建(Phase A 完)', '設備/耗材庫必須存在'),
            ('📊 資料', 'case_factory + qty/pkg (Phase C)', 'process_template_id 已選'),
            ('📄 文件', 'RD Process Flow xlsx', 'takt 來源 · 文字 ref (takt_source 欄)'),
            ('📄 文件', 'EE BOM Cleansheet 工時 章節', 'industrial engineer 給的標準工時'),
            ('📊 資料', '廠設備 master 子集', '案內 line_qty / qty_per_line'),
            ('📄 文件', 'PKG BOM xlsx', '商包 16 item / 工包 3 item'),
        ],
        'steps': [
            ('D1', '從 template 帶入製程清單', 'EPM', 'INSERT bom_cs_case_process FROM ..._template_step',
                '一鍵 copy · 9 製程 · 帶 step_order', 5),
            ('D2', 'EPM 調各製程 9 變數', 'EPM', 'UPDATE bom_cs_case_process',
                'takt/yield/eff/工時/工作日/shifts/DL/debug DL/functional DL/lines', 60),
            ('D3', 'IDL × 製程 multiplier matrix', 'EPM', 'INSERT bom_cs_case_idl_alloc',
                '17 role × 9 process · 首次抄 Excel J64/L64/N64...', 90),
            ('D4', '案級設備 binding', 'EPM', 'INSERT bom_cs_case_equipment',
                '從廠 master 選 32 件 · 設 line_qty/qty_per_line · 可 override life/acq', 30),
            ('D5', '案級耗材 binding', 'EPM', 'INSERT bom_cs_case_consumable',
                '從廠 master 選 · 設 annual_usage_qty · 可 override unit_cost', 20),
            ('D6', 'PKG items 上傳(per pkg_version)', 'PKG 採購 Ken', 'INSERT bom_cs_case_pkg_item',
                '商包 16 / 工包 3 · qty/unit/spec/vendor/true_cost/quote_price/markup', 45),
        ],
        'delta_steel': '9 製程 × 17 IDL · 32 設備 binding · 17 耗材 · 商包 16 + 工包 3 item',
        'delta_whoop': '每 sub-BOM(8 module)各自 SMT 配置 · FATP 7 station 串接 · PKG SKU × Module 矩陣 5×8=40 點',
        'schemas': ['bom_cs_case_process', 'bom_cs_case_idl_alloc', 'bom_cs_case_equipment',
                    'bom_cs_case_consumable', 'bom_cs_case_pkg_item'],
    },
    {
        'code': 'E', 'title': '計算 (Compute)',
        'timing': '配置完按鈕 · 30 秒 / case_factory',
        'owner': 'EPM 觸發 · 系統執行',
        'summary': '5 區公式跑(DL / IDL / Equip+Facility / Others / Total)· 寫 cs_run header + cs_run_cell × (process × component × qty_scenario) · Regression test vs Excel < 0.01ε。',
        'prereq': [
            ('📊 資料', 'Phase D 全完成', 'case_process/idl_alloc/equipment/consumable/pkg_item 齊'),
            ('📊 資料', 'BOM main board cost', 'motherboard_cost_ref ≈ $8.683 · 算 VAT 用'),
            ('⚙️ 設定', 'compute engine 版本', 'db_v1 · 對齊 §App-A 100+ Excel 公式'),
            ('⚙️ 設定', 'Regression test 容差', '每 cell vs Excel · ε<0.01 PASS · >0.01 → run.status=failed'),
        ],
        'steps': [
            ('E1', 'EPM 按 「🔄 Compute」', 'EPM', 'POST /api/cs/case-factories/:id/compute',
                '新 cs_run(status=computing) · case_factory.dirty=0', 1),
            ('E2', 'compute_dl_cost(per process × qty)', '系統', 'compute_dl_cost.py',
                'UPH → max_output → lines → weekly_output → DL/day → multipliers → /unit · C58', 5),
            ('E3', 'compute_idl_cost(per process × qty)', '系統', 'compute_idl_cost.py',
                'centralized + line-dep IDL · rows 64-83', 5),
            ('E4', 'compute_equipment_cost(MRO + Depr)', '系統', 'compute_equipment_cost.py',
                'SUM(P:P) × (annual / weekly) · H90/H91', 3),
            ('E5', 'compute_facility + ind_materials', '系統', 'compute_facility_cost.py',
                'sqft × $unit × IDL multiplier × (annual / weekly) · J100', 3),
            ('E6', 'compute_others(Freight+VAT+Loss)', '系統', 'compute_others_cost.py',
                'inbound_freight/annual · motherboard×vat/18 · loss×material · J107/J110/J114', 2),
            ('E7', '寫 cs_run_cell × N rows', '系統', 'INSERT bom_cs_run_cell',
                '9 process × 9 component × qty_scenarios', 1),
            ('E8', 'Roll-up MVA + SGA + Profit', '系統', 'compute_mva_total + compute_sga_profit',
                'MVA = SUM(cell) · SGA = motherboard × 0.02 · Profit = (MVA+motherboard) × 0.14', 1),
            ('E9', 'Regression test vs Excel(自動)', '系統 CI', 'pytest test_compute_regression.py',
                '每 cell vs Excel · ε>0.01 → status=failed + warnings_json', 5),
            ('E10', 'cs_run.status = ready', '系統', 'UPDATE bom_cs_run',
                'computed_at + 觸發 notification', 1),
        ],
        'delta_steel': '本案 = 3 廠 × 2 qty = 6 cs_run · 各 ~50 cell rows',
        'delta_whoop': '6 cs_run · 但每 run 含 8 sub-BOM × 各自 SMT process · FATP 7 station · ~300+ cell rows · compute ~3 min',
        'schemas': ['bom_cs_run', 'bom_cs_run_cell'],
    },
    {
        'code': 'F', 'title': 'BOM Lock Propagate',
        'timing': 'DPM gate · 5 分鐘',
        'owner': 'DPM(觸發)+ 系統(自動)',
        'summary': 'BOM 鎖版時自動跑 material × MVA × SGA × Profit 整合 · 寫 cs_run_result fact table(SteelSeries:24 列 · WHOOP:依 PKG×module 展開)· True/Quote 雙價 + Margin via Oracle GENERATED VIRTUAL。',
        'prereq': [
            ('📊 資料', 'cs_run 已 ready (Phase E)', '各廠各 qty · status≠failed'),
            ('📊 資料', 'BOM 採購策略已選定', '每 item 已 select_mfg + select_tier'),
            ('📊 資料', 'PKG BOM items 雙價齊', '商包 16 / 工包 3 各 true/quote'),
            ('⚙️ 設定', 'variant scope 對齊', 'shared / per_variant · 影響 cell material 算法'),
            ('📐 政策', 'multi-PKG lock 順序', '全 lock / 只 primary · DPM UI 選'),
        ],
        'steps': [
            ('F1', 'DPM 按 「🔒 Lock BOM」', 'DPM Mike', 'POST /api/bom/:id/lock',
                '前置 check(run ready / 雙價齊 / variant 對) · pass propagate', 3),
            ('F2', '列每 case_factory 維度組合', '系統', 'in-memory 展開',
                'variants × qty × pkg × factory · SteelSeries=24 · WHOOP 視 SKU 矩陣展開', 1),
            ('F3', 'compute_material_per_combo', '系統', 'SUMPRODUCT(item × tier match × variant)',
                'tier: qty_min≤target≤qty_max · 出 true/quote', 5),
            ('F4', '取 PKG cost(per pkg_version)', '系統', 'SELECT bom_cs_case_pkg',
                'Phase D 已算 · 商包 $0.706 / 工包 $0.182', 1),
            ('F5', '取 MVA(per qty_scenario)', '系統', 'SUM(bom_cs_run_cell WHERE run_id AND qty)', 'Phase E 已算', 1),
            ('F6', '算 SGA + Profit', '系統', '公式 hard-code',
                'sga = motherboard × 0.02 · profit = (MVA+motherboard) × 0.14', 1),
            ('F7', 'UPSERT bom_cs_run_result × N rows', '系統', 'MERGE INTO bom_cs_run_result',
                'UNIQUE (run, factory, variant, qty, pkg) · 24 列 · Margin/Total GENERATED', 2),
            ('F8', 'audit_log 寫 PROPAGATE event', '系統', 'INSERT bom_cs_audit_log',
                'event=BOM_PROPAGATE · payload', 1),
            ('F9', 'notify 影響干係人', '系統', 'user_notifications', 'EPM/BPM/採購主管 各 chip', 1),
        ],
        'delta_steel': 'propagate 寫 24 列 · pivot view 直接看 6 cells (3 廠 × 2 variant)',
        'delta_whoop': 'propagate 寫多 PKG × 多 module · 5 SKU × 24 cells = 120 列 · 看 module 攤分(由 SKU_MODULE_MATRIX 加權)',
        'schemas': ['bom_cs_run_result', 'bom_cs_audit_log', 'user_notifications'],
    },
    {
        'code': 'G', 'title': '報價 + Margin 分析',
        'timing': '業務 Gate · 4 小時',
        'owner': '業務 + BPM + 採購主管 + finance',
        'summary': 'Factory matrix pivot view · 24 cells 切 qty/pkg/variant 動態 refresh · True Cost 走 VIEW_TRUE_COST 權限(finance/採購主管/dpm 可見)· 業務只看 Quote · 議價時走 reprice / unlock 開新 cs_run cycle。',
        'prereq': [
            ('📊 資料', 'BOM 已 lock (Phase F)', 'cs_run_result 已有 24 列'),
            ('🔐 權限', 'VIEW_TRUE_COST 已分配', 'finance/dpm/採購主管/EPM · 一般業務 ❌'),
            ('📊 資料', '競品 / 客戶歷史報價', 'BPM 從 CRM 拉 · 對齊 markup % 上下限'),
            ('📐 政策', 'reprice 觸發條件', '議價 → 新 cs_run · 舊 archived · UI 只列 ready'),
        ],
        'steps': [
            ('G1', '§factory_matrix UI 開啟', 'BPM/採購主管/finance', 'GET /api/cs/cases/:id/matrix',
                '24 cells × {factory,variant,qty,pkg,true,quote,mva,margin%}', 5),
            ('G2', 'Qty/PKG/Variant toggle 切換', '用戶', '前端 state pivot',
                'v07.qty / v07.pkg → re-fetch · 6 cell 主表 refresh', 10),
            ('G3', 'Margin Heatmap 視角', 'finance/採購主管', 'GET /api/cs/margin/:run',
                '24 cells heatmap · margin% 顏色 · Top Markup Items per PKG · 需 VIEW_TRUE_COST', 15),
            ('G4', 'BPM 匯出 RFQ Cost Excel', 'BPM Lisa', 'POST /api/quotes/:id/export',
                '只匯 Quote(True/Margin strip)· region/qty/pkg 分頁', 20),
            ('G5', '業務/BPM 簽核 submit', '業務+BPM', 'INSERT rfq_submissions',
                'submit quote · 系統建 RFQ record · audit log · status=submitted', 5),
            ('G6', '(議價)客戶回 → 採購 reprice', '採購+EPM', 'POST /api/cs/case-factories/:id/reprice',
                '改 tier or mfg · 新 compute · 舊 archived · 新 run 接續', '循環'),
        ],
        'delta_steel': 'BPM 用 §factory_matrix toggle 24 cells · finance 看 §margin_analysis heatmap',
        'delta_whoop': '同 SteelSeries · 但 §whoop_summary 用 5 SKU 拆解視角 · 看「同 module 跨 SKU 攤分」效應',
        'schemas': ['bom_cs_run_result', 'rfq_submissions', 'ai_data_policy_rules (VIEW_TRUE_COST)'],
    },
]


# ════════════════════════════════════════════════════════════════════════
# COVER SLIDE
# ════════════════════════════════════════════════════════════════════════

def slide_cover():
    s = add_blank_slide()
    add_rect(s, 0, 0, 13.333, 7.5, INK)
    # gradient accent
    add_rect(s, 0, 2.6, 13.333, 0.06, OCEAN)
    # subtitle eyebrow
    add_text(s, 0.7, 1.2, 12, 0.4, 'FOXLINK GPT · Cortex BOM 模組', size=12, color=OCEAN, font_name='Microsoft JhengHei')
    add_text(s, 0.7, 1.6, 12, 0.85, 'Cortex MVA 操作流程說明手冊', size=42, bold=True, color=WHITE)
    add_text(s, 0.7, 2.85, 12, 0.5, 'Cleansheet × Manufacturing Value Add — 細部操作指南', size=20, color=WHITE)

    # 兩案 banner
    add_rounded(s, 0.7, 4.0, 5.8, 1.5, RGBColor(0x1E, 0x40, 0xAF), radius=0.1)
    add_text(s, 0.95, 4.2, 5.3, 0.4, '🖱️  SteelSeries Rival 3+ Wired', size=16, bold=True, color=WHITE)
    add_text(s, 0.95, 4.65, 5.3, 0.8, 'MOUSE_STD 模板 · 1-BOM · 9 製程\n3 廠 × 2 variant × 2 qty × 2 pkg = 24 result cells', size=11, color=RGBColor(0xCB, 0xD5, 0xE1))

    add_rounded(s, 6.85, 4.0, 5.8, 1.5, RGBColor(0x6B, 0x21, 0xA8), radius=0.1)
    add_text(s, 7.1, 4.2, 5.3, 0.4, '⌚  WHOOP Gen4 MP', size=16, bold=True, color=WHITE)
    add_text(s, 7.1, 4.65, 5.3, 0.8, 'WHOOP_WEARABLE · 8 sub-BOM · FATP 7 stations\n5 PKG SKU × 8 module 矩陣 = 40 點配置', size=11, color=RGBColor(0xE9, 0xD5, 0xFF))

    add_text(s, 0.7, 6.0, 12, 0.4, 'Phase A → G · 33 操作步驟 · 兩案差異對照 · 內部討論用', size=14, color=OCEAN)
    add_text(s, 0.7, 6.5, 12, 0.4, 'v1 · 2026-06 · 對應 cleansheet-mva-sd v0.3 + bom-collection-sd v0.4', size=11, color=MUTED)


def slide_toc():
    s = add_blank_slide()
    slide_header(s, '本手冊涵蓋什麼?', '7 個 Phase · 33 步驟 · 兩案視角整合')

    # 7 phase cards
    for i, ph in enumerate(PHASES):
        row = i // 4
        col = i % 4
        x = 0.5 + col * 3.15
        y = 1.5 + row * 2.5
        color = PHASE_COLORS[ph['code']]
        add_rounded(s, x, y, 3.0, 2.2, WHITE, line=color, line_width=2, radius=0.04)
        add_rounded(s, x, y, 3.0, 0.45, color, radius=0.04)
        add_text(s, x + 0.15, y + 0.06, 2.7, 0.33, f"Phase {ph['code']}", size=12, bold=True, color=WHITE)
        add_text(s, x + 0.15, y + 0.55, 2.7, 0.4, ph['title'].split(' (')[0], size=12, bold=True, color=INK)
        add_text(s, x + 0.15, y + 0.95, 2.7, 0.3, f"👥 {ph['owner']}", size=8.5, color=MUTED)
        add_text(s, x + 0.15, y + 1.25, 2.7, 0.3, f"⏱️ {ph['timing']}", size=8.5, color=MUTED)
        add_text(s, x + 0.15, y + 1.6, 2.7, 0.55, f"📋 {len(ph['prereq'])} 素材 · 🎯 {len(ph['steps'])} 步驟",
                 size=9, bold=True, color=color)

    slide_footer(s, 2, total_pages)


def slide_what_is_mva():
    s = add_blank_slide()
    slide_header(s, 'MVA 是什麼?為什麼要做?', 'spec cleansheet-mva-sd §0 · §1')

    # 左:What
    add_rounded(s, 0.5, 1.4, 6.0, 5.3, BG_SOFT, line=NAVY, line_width=1, radius=0.02)
    add_text(s, 0.7, 1.55, 5.7, 0.45, '💡 MVA = Manufacturing Value Add', size=18, bold=True, color=NAVY)
    add_text(s, 0.7, 2.1, 5.7, 0.35, '製造轉換成本(廠端把料變成品的所有費用)', size=11, color=TEXT)

    # Formula box
    add_rounded(s, 0.7, 2.65, 5.6, 1.6, WHITE, line=NAVY, radius=0.04)
    add_text(s, 0.85, 2.78, 5.4, 0.35, 'TC = Material + MVA + SG&A + Profit', size=13, bold=True, color=INK)
    add_text(s, 0.85, 3.2, 5.4, 0.95, """MVA = DL + IDL + Equipment(MRO+Depr) + Indirect Mat
        + Facility + Inbound Freight + VAT + Loss
        ━━━━━━━━━━━━━━━━━━━━━━━━━━
        9 個 cost component × N 個製程""", size=10, color=TEXT)

    add_text(s, 0.7, 4.45, 5.7, 0.4, '🎯 分工:', size=12, bold=True, color=INK)
    add_text(s, 0.7, 4.85, 5.7, 1.8, """• Material → BOM 端算(per variant · 跨廠共用採購)
• MVA → Cleansheet 端算(per factory · 跨色共用)
• SG&A + Profit → 廠 baseline(集團政策可統一)
• Total → §factory_matrix UI 整合呈現""", size=10, color=TEXT)

    # 右:Why(痛點)
    add_rounded(s, 6.85, 1.4, 6.0, 5.3, RGBColor(0xFE, 0xF2, 0xF2), line=RGBColor(0xDC, 0x26, 0x26), line_width=1, radius=0.02)
    add_text(s, 7.05, 1.55, 5.7, 0.45, '🔥 EPM Excel 現況的 5 個痛點', size=18, bold=True, color=RGBColor(0xB9, 0x1C, 0x1C))

    pains = [
        ('1', '檔案分散', '每案每廠 EPM 各自一份 xlsx · 公司無 master'),
        ('2', '人走人重做', '計算引擎 1500 行 Excel 公式 · 新 EPM 看不懂'),
        ('3', '月度更新漏', 'wage/設備異動 → 各 EPM 手動同步 · 常漏'),
        ('4', '跨案無比對', '同款 DEK 跨案 utilization 無法查詢'),
        ('5', '改料人工 propagate', 'BOM 變動 → Cleansheet 不自動更新'),
    ]
    for i, (n, t, d) in enumerate(pains):
        y = 2.2 + i * 0.85
        add_rounded(s, 7.05, y, 0.45, 0.45, RGBColor(0xDC, 0x26, 0x26), radius=0.5)
        add_text(s, 7.05, y, 0.45, 0.45, n, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, 7.6, y, 5.2, 0.35, t, size=11, bold=True, color=INK)
        add_text(s, 7.6, y + 0.32, 5.2, 0.5, d, size=9, color=MUTED)

    slide_footer(s, 3, total_pages)


def slide_two_cases():
    s = add_blank_slide()
    slide_header(s, '兩個案子的差異 (Why both?)', '本手冊整合 SteelSeries (MOUSE_STD) 與 WHOOP (WEARABLE) 兩種模板視角')

    # SteelSeries
    add_rounded(s, 0.5, 1.4, 6.1, 5.3, STEEL_BG, line=STEEL_COLOR, line_width=2, radius=0.02)
    add_rounded(s, 0.5, 1.4, 6.1, 0.55, STEEL_COLOR, radius=0.02)
    add_text(s, 0.7, 1.48, 5.8, 0.45, '🖱️  SteelSeries Rival 3+ Wired Mouse', size=15, bold=True, color=WHITE)

    rows_steel = [
        ('產品類別', 'Wired Mouse · 消費電子'),
        ('process_template', 'MOUSE_STD(9 製程)'),
        ('BOM 結構', '1-BOM(單層 EE/ME merge)'),
        ('Sub-BOM', '無(單一主板組裝)'),
        ('變體 (variants)', 'black / white(2 色)'),
        ('PKG 版本', '商包(16 item) · 工包(3 item)'),
        ('Qty Scenarios', 'Low=100K(US) · High=418K(global)'),
        ('result cells', '3 廠 × 2 variant × 2 qty × 2 pkg = 24'),
        ('CN baseline MVA', '$1.86 / unit'),
        ('採購來源 xlsx', 'Rival 3+ Wired Mouse Cleansheet_China-20241011-19S'),
    ]
    for i, (k, v) in enumerate(rows_steel):
        y = 2.1 + i * 0.46
        add_text(s, 0.75, y, 1.85, 0.35, k, size=10, bold=True, color=STEEL_COLOR)
        add_text(s, 2.65, y, 3.85, 0.4, v, size=10, color=TEXT)

    # WHOOP
    add_rounded(s, 6.85, 1.4, 6.0, 5.3, WHOOP_BG, line=WHOOP_COLOR, line_width=2, radius=0.02)
    add_rounded(s, 6.85, 1.4, 6.0, 0.55, WHOOP_COLOR, radius=0.02)
    add_text(s, 7.05, 1.48, 5.8, 0.45, '⌚  WHOOP Gen4 MP(可穿戴)', size=15, bold=True, color=WHITE)

    rows_whoop = [
        ('產品類別', 'Wearable Device · 健康追蹤'),
        ('process_template', 'WHOOP_WEARABLE(模組化)'),
        ('BOM 結構', 'Multi-Sub-BOM(8 模組層)'),
        ('Sub-BOM', '8 個 bom_module(各自 SMT)'),
        ('變體 (variants)', '3 色款'),
        ('PKG 版本', '5 PKG SKU × 8 module 包含矩陣'),
        ('Qty Scenarios', 'TBD(視 MP 拍板)'),
        ('result cells', '~120(5 SKU × 24 multi-dim)'),
        ('Process 特色', 'Module Integ + Final Assy + FATP 7 station'),
        ('採購來源 xlsx', 'WHOOP_Gen4 MP quotation book_ for AI'),
    ]
    for i, (k, v) in enumerate(rows_whoop):
        y = 2.1 + i * 0.46
        add_text(s, 7.1, y, 1.85, 0.35, k, size=10, bold=True, color=WHOOP_COLOR)
        add_text(s, 9.0, y, 3.85, 0.4, v, size=10, color=TEXT)

    slide_footer(s, 4, total_pages)


def slide_architecture():
    s = add_blank_slide()
    slide_header(s, 'MVA 5-Layer 架構', 'spec cleansheet-mva-sd §2 · 全結構化取代 Excel')

    layers = [
        ('Layer 1', '製程目錄 + 模板', '公司級 · 跨案共用',
            'bom_process_catalog · bom_process_template · bom_process_template_step',
            RGBColor(0xF3, 0xE8, 0xFF), RGBColor(0x7C, 0x3A, 0xED)),
        ('Layer 2', '廠級資產 Baseline', 'per factory · SCD Type 2 月度更新',
            'bom_factory_baseline · bom_factory_idl_role · bom_factory_equipment · bom_factory_consumable',
            RGBColor(0xDB, 0xEA, 0xFE), RGBColor(0x1E, 0x40, 0xAF)),
        ('Layer 3', '案級配置', '案 × 廠 × variant',
            'bom_cs_case_factory · bom_cs_case_process · bom_cs_case_equipment · bom_cs_case_pkg',
            RGBColor(0xD1, 0xFA, 0xE5), RGBColor(0x16, 0xA3, 0x4A)),
        ('Layer 4', '計算結果 snapshot', '一次 compute · 可多 run',
            'bom_cs_run · bom_cs_run_cell · bom_cs_run_result(v0.3 加 fact table)',
            RGBColor(0xFE, 0xE2, 0xE2), RGBColor(0xDC, 0x26, 0x26)),
        ('Layer 5', 'Audit + Propagate', '記錄變更 · BOM lock 自動展開',
            'bom_cs_audit_log · bom_cs_baseline_diff',
            RGBColor(0xFE, 0xF3, 0xC7), RGBColor(0xB4, 0x53, 0x09)),
    ]
    for i, (n, t, d, schemas, bg, fg) in enumerate(layers):
        y = 1.5 + i * 1.05
        add_rounded(s, 0.7, y, 12.0, 0.95, bg, line=fg, line_width=1, radius=0.03)
        # Layer badge
        add_rounded(s, 0.85, y + 0.15, 1.6, 0.65, fg, radius=0.04)
        add_text(s, 0.85, y + 0.15, 1.6, 0.65, n, size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, 2.65, y + 0.13, 6, 0.35, t, size=13, bold=True, color=INK)
        add_text(s, 2.65, y + 0.46, 6, 0.3, d, size=10, color=MUTED)
        # schemas
        add_text(s, 8.85, y + 0.13, 3.7, 0.7, schemas, size=8, color=TEXT, font_name='Consolas')

    # Bottom note
    add_text(s, 0.7, 6.7, 12, 0.35, '👉 本手冊 Phase A 建 Layer 1+2 · Phase C-D 建 Layer 3 · Phase E 算 Layer 4 · Phase F 寫 Layer 4+5',
             size=10, color=NAVY, bold=True)

    slide_footer(s, 5, total_pages)


# ════════════════════════════════════════════════════════════════════════
# PHASE SLIDES(每 Phase 3-4 slides)
# ════════════════════════════════════════════════════════════════════════

def slide_phase_intro(ph, page_n):
    s = add_blank_slide()
    color = PHASE_COLORS[ph['code']]

    # Full-bleed banner
    add_rect(s, 0, 0, 13.333, 7.5, color)
    # Layered effect
    add_rounded(s, 0.4, 0.4, 12.55, 6.7, WHITE, radius=0.02)

    # Phase indicator
    add_rounded(s, 0.8, 0.95, 2.4, 0.7, color, radius=0.04)
    add_text(s, 0.8, 0.95, 2.4, 0.7, f"PHASE {ph['code']}", size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Title
    add_text(s, 0.8, 1.85, 12, 0.7, ph['title'], size=32, bold=True, color=INK)

    # Meta strip
    add_rounded(s, 0.8, 2.8, 5.8, 0.5, RGBColor(0xF3, 0xF4, 0xF6), radius=0.05)
    add_text(s, 1.0, 2.85, 5.6, 0.4, f"👥 主負責: {ph['owner']}", size=11, bold=True, color=TEXT, anchor=MSO_ANCHOR.MIDDLE)

    add_rounded(s, 6.8, 2.8, 5.8, 0.5, RGBColor(0xF3, 0xF4, 0xF6), radius=0.05)
    add_text(s, 7.0, 2.85, 5.6, 0.4, f"⏱️ 時程: {ph['timing']}", size=11, bold=True, color=TEXT, anchor=MSO_ANCHOR.MIDDLE)

    # Summary block
    add_rounded(s, 0.8, 3.55, 11.8, 1.3, RGBColor(0xFA, 0xFB, 0xFC), line=color, line_width=2, radius=0.03)
    add_text(s, 1.05, 3.7, 11.3, 1.1, ph['summary'], size=14, color=INK)

    # KPI 3 boxes
    kpis = [
        ('📋', str(len(ph['prereq'])), '輸入素材'),
        ('🎯', str(len(ph['steps'])), '操作步驟'),
        ('🗄️', str(len(ph['schemas'])), '影響表'),
    ]
    for i, (ico, n, t) in enumerate(kpis):
        x = 0.8 + i * 4.0
        add_rounded(s, x, 5.05, 3.8, 1.5, color, radius=0.04)
        add_text(s, x, 5.15, 3.8, 0.55, ico, size=28, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(s, x, 5.7, 3.8, 0.55, n, size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(s, x, 6.25, 3.8, 0.3, t, size=12, color=RGBColor(0xFE, 0xFE, 0xFE), align=PP_ALIGN.CENTER)

    slide_footer(s, page_n, total_pages, side_text=f"Phase {ph['code']} 開始 · Cortex MVA v1")


def slide_phase_prereq(ph, page_n):
    s = add_blank_slide()
    color = PHASE_COLORS[ph['code']]
    slide_header(s, f"① 需要準備的素材", f"Phase {ph['code']} · {ph['title']}", ph['code'])

    # Subtitle
    add_text(s, 0.4, 1.25, 12, 0.35, f"進此 Phase 前必須備齊 · 共 {len(ph['prereq'])} 項",
             size=11, color=MUTED)

    # Prereq cards (2 column)
    for i, (kind, name, detail) in enumerate(ph['prereq']):
        row = i // 2
        col = i % 2
        x = 0.5 + col * 6.3
        y = 1.75 + row * 1.05
        # kind colors
        kind_color = {
            '📄 文件': (RGBColor(0xFE, 0xF3, 0xC7), RGBColor(0x92, 0x40, 0x0E)),
            '📊 資料': (RGBColor(0xDB, 0xEA, 0xFE), RGBColor(0x1E, 0x40, 0xAF)),
            '📐 政策': (RGBColor(0xFC, 0xE7, 0xF3), RGBColor(0x9D, 0x17, 0x4D)),
            '⚙️ 設定': (RGBColor(0xE0, 0xE7, 0xFF), RGBColor(0x37, 0x30, 0xA3)),
            '🔐 權限': (RGBColor(0xFE, 0xE2, 0xE2), RGBColor(0x99, 0x1B, 0x1B)),
        }.get(kind, (RGBColor(0xF3, 0xF4, 0xF6), RGBColor(0x37, 0x47, 0x55)))

        add_rounded(s, x, y, 6.1, 0.95, kind_color[0], line=kind_color[1], line_width=0.75, radius=0.04)
        # kind pill
        add_rounded(s, x + 0.15, y + 0.15, 1.05, 0.3, kind_color[1], radius=0.5)
        add_text(s, x + 0.15, y + 0.15, 1.05, 0.3, kind, size=9, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # name
        add_text(s, x + 1.35, y + 0.13, 4.6, 0.35, name, size=12, bold=True, color=INK)
        # detail
        add_text(s, x + 0.15, y + 0.5, 5.8, 0.42, detail, size=9, color=TEXT)

    slide_footer(s, page_n, total_pages)


def slide_phase_steps(ph, page_n):
    s = add_blank_slide()
    color = PHASE_COLORS[ph['code']]
    slide_header(s, f"② 操作步驟 (Step-by-step)", f"Phase {ph['code']} · {ph['title']}", ph['code'])

    # Total time
    total_mins = sum(s_[5] for s_ in ph['steps'] if isinstance(s_[5], (int, float)))
    cyclic = any(not isinstance(s_[5], (int, float)) for s_ in ph['steps'])
    sub = f"{len(ph['steps'])} 步 · 預估總計 {total_mins} 分鐘"
    if cyclic:
        sub += ' · 含循環步驟'
    add_text(s, 0.4, 1.25, 12, 0.35, sub, size=11, color=MUTED)

    # Steps list
    n_steps = len(ph['steps'])
    row_h = min(0.62, 5.4 / n_steps)
    for i, (num, name, who, sys_, detail, mins) in enumerate(ph['steps']):
        y = 1.75 + i * row_h
        bg = BG_SOFT if i % 2 == 0 else WHITE
        add_rect(s, 0.4, y, 12.55, row_h - 0.04, bg)
        # step number badge
        add_rounded(s, 0.5, y + 0.05, 0.7, row_h - 0.14, color, radius=0.1)
        add_text(s, 0.5, y + 0.05, 0.7, row_h - 0.14, num, size=11, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # name + who
        add_text(s, 1.3, y + 0.05, 5.3, 0.3, name, size=11, bold=True, color=INK)
        mins_str = f"{mins} min" if isinstance(mins, (int, float)) else mins
        add_text(s, 1.3, y + 0.32, 5.3, 0.25, f"👤 {who} · ⏱ {mins_str}", size=8, color=MUTED)
        # detail
        add_text(s, 6.7, y + 0.05, 4.5, row_h - 0.14, detail, size=8.5, color=TEXT)
        # sys (mono)
        add_text(s, 11.3, y + 0.1, 1.55, row_h - 0.2, sys_, size=7.5, color=color, font_name='Consolas')

    slide_footer(s, page_n, total_pages)


def slide_phase_delta(ph, page_n):
    s = add_blank_slide()
    color = PHASE_COLORS[ph['code']]
    slide_header(s, f"③④ 兩案差異 + 影響的 DB 表", f"Phase {ph['code']} · {ph['title']}", ph['code'])

    # 兩案 side by side
    add_rounded(s, 0.5, 1.4, 6.1, 3.2, STEEL_BG, line=STEEL_COLOR, line_width=2, radius=0.02)
    add_rounded(s, 0.5, 1.4, 6.1, 0.55, STEEL_COLOR, radius=0.02)
    add_text(s, 0.7, 1.48, 5.8, 0.45, '🖱️ SteelSeries (MOUSE_STD)', size=14, bold=True, color=WHITE)
    add_text(s, 0.75, 2.1, 5.85, 2.4, ph['delta_steel'], size=12, color=TEXT)

    add_rounded(s, 6.85, 1.4, 6.0, 3.2, WHOOP_BG, line=WHOOP_COLOR, line_width=2, radius=0.02)
    add_rounded(s, 6.85, 1.4, 6.0, 0.55, WHOOP_COLOR, radius=0.02)
    add_text(s, 7.05, 1.48, 5.8, 0.45, '⌚ WHOOP (WEARABLE)', size=14, bold=True, color=WHITE)
    add_text(s, 7.1, 2.1, 5.75, 2.4, ph['delta_whoop'], size=12, color=TEXT)

    # Schemas
    add_text(s, 0.4, 4.85, 12, 0.45, '🗄️ 本 Phase 結束後 · 以下 DB 表的狀態被改', size=14, bold=True, color=INK)

    # schema chips
    cy = 5.45
    cx = 0.5
    for sch in ph['schemas']:
        w = len(sch) * 0.085 + 0.4
        if cx + w > 12.85:
            cx = 0.5
            cy += 0.55
        add_rounded(s, cx, cy, w, 0.4, color, radius=0.2)
        add_text(s, cx, cy, w, 0.4, sch, size=10, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name='Consolas')
        cx += w + 0.15

    slide_footer(s, page_n, total_pages)


# ════════════════════════════════════════════════════════════════════════
# CLOSING SLIDES
# ════════════════════════════════════════════════════════════════════════

def slide_timeline(page_n):
    s = add_blank_slide()
    slide_header(s, '端到端 Timeline', '從 IT 建置 → 案級配置 → 計算 → 報價 · 時程一覽')

    # Horizontal flow
    y = 2.5
    bar_y = y + 0.45
    add_rect(s, 0.5, bar_y, 12.3, 0.04, LINE)

    items = [
        ('A', '一次性建置', '2 週/廠', '初始'),
        ('B', '月度維護', '1 天/廠', '每月底'),
        ('C', '新案開立', '1 天', '案啟動'),
        ('D', '案級配置', '3-5 天/廠', '主要工作'),
        ('E', '計算', '30 秒', '按鈕'),
        ('F', 'BOM Lock', '5 分鐘', 'DPM gate'),
        ('G', '報價 + Margin', '4 小時', '業務 gate'),
    ]
    seg_w = 12.3 / len(items)
    for i, (code, name, dur, when) in enumerate(items):
        cx = 0.5 + i * seg_w + seg_w / 2 - 0.5
        color = PHASE_COLORS[code]
        # Dot
        add_rounded(s, cx, bar_y - 0.18, 0.4, 0.4, color, radius=0.5)
        add_text(s, cx, bar_y - 0.18, 0.4, 0.4, code, size=13, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # Name + dur (above or below alternating)
        if i % 2 == 0:
            add_text(s, cx - 0.45, bar_y - 0.95, 1.3, 0.3, name, size=10, bold=True, color=INK, align=PP_ALIGN.CENTER)
            add_text(s, cx - 0.45, bar_y - 0.65, 1.3, 0.25, dur, size=8, color=color, align=PP_ALIGN.CENTER, bold=True)
            add_text(s, cx - 0.45, bar_y - 0.4, 1.3, 0.22, when, size=7, color=MUTED, align=PP_ALIGN.CENTER)
        else:
            add_text(s, cx - 0.45, bar_y + 0.3, 1.3, 0.3, name, size=10, bold=True, color=INK, align=PP_ALIGN.CENTER)
            add_text(s, cx - 0.45, bar_y + 0.6, 1.3, 0.25, dur, size=8, color=color, align=PP_ALIGN.CENTER, bold=True)
            add_text(s, cx - 0.45, bar_y + 0.85, 1.3, 0.22, when, size=7, color=MUTED, align=PP_ALIGN.CENTER)

    # Bottom note
    add_rounded(s, 0.5, 4.5, 12.3, 2.4, BG_SOFT, line=NAVY, line_width=1, radius=0.02)
    add_text(s, 0.7, 4.65, 12, 0.35, '💡 整體建議排程(以本 SteelSeries 案為例)', size=13, bold=True, color=NAVY)
    add_text(s, 0.7, 5.1, 12.0, 1.7, """• Week -2 ~ -1:Phase A 公司+CN 廠首建(IT/Admin 人力)
• Week 0:案啟動 → Phase C(BPM 開案 1 天)
• Week 1-2:Phase D 三廠 EPM 並行(各 3-5 天)
• Week 2:Phase E EPM 按 Compute(各 30 秒)+ regression test
• Week 3 結束前:Phase F DPM lock + propagate(5 分鐘)
• Week 3:Phase G 業務 gate + 客戶議價(4 小時 ~ 數天迭代)""", size=10, color=TEXT)

    slide_footer(s, page_n, total_pages)


def slide_checklist(page_n):
    s = add_blank_slide()
    slide_header(s, '🎯 操作 Checklist 速查表', '對齊 Cortex_互動Demo_v0.9.html §mva_workflow section')

    # Header row
    headers = ['Phase', '主負責人', '時程', '輸入', '步驟', '影響表']
    widths = [1.0, 3.5, 2.0, 1.3, 1.3, 3.5]
    cx = 0.4
    for h, w in zip(headers, widths):
        add_rect(s, cx, 1.5, w, 0.5, INK)
        add_text(s, cx, 1.5, w, 0.5, h, size=11, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        cx += w

    # Rows
    for i, ph in enumerate(PHASES):
        y = 2.0 + i * 0.65
        bg = BG_SOFT if i % 2 == 0 else WHITE
        cx = 0.4
        cells = [ph['code'], ph['owner'], ph['timing'], str(len(ph['prereq'])), str(len(ph['steps'])), str(len(ph['schemas']))]
        for j, (cell, w) in enumerate(zip(cells, widths)):
            add_rect(s, cx, y, w, 0.6, bg)
            if j == 0:
                color = PHASE_COLORS[ph['code']]
                add_rounded(s, cx + 0.15, y + 0.1, w - 0.3, 0.4, color, radius=0.1)
                add_text(s, cx + 0.15, y + 0.1, w - 0.3, 0.4, cell, size=12, bold=True, color=WHITE,
                         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            else:
                add_text(s, cx + 0.1, y, w - 0.2, 0.6, cell, size=9.5, color=TEXT, anchor=MSO_ANCHOR.MIDDLE)
            cx += w

    # Total
    total_steps = sum(len(p['steps']) for p in PHASES)
    total_prereq = sum(len(p['prereq']) for p in PHASES)
    total_schemas = sum(len(p['schemas']) for p in PHASES)
    y = 2.0 + len(PHASES) * 0.65
    add_rounded(s, 0.4, y, sum(widths), 0.6, OCEAN, radius=0.04)
    add_text(s, 0.6, y, sum(widths) - 0.4, 0.6,
             f"  TOTAL · 7 Phase · {total_prereq} 輸入素材 · {total_steps} 操作步驟 · 影響 {total_schemas} 個 DB 表 entries",
             size=12, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)

    slide_footer(s, page_n, total_pages)


def slide_excel_mapping(page_n):
    s = add_blank_slide()
    slide_header(s, 'Backup: Excel 公式 → DB Service 對照(摘錄)', 'spec cleansheet-mva-sd Appendix A · 100+ 條完整對應')

    mappings = [
        ('Cleansheet!C14', '=3600/C13*C15*C16', 'compute_uph(case_factory, process)', 'UPH/hr'),
        ('Cleansheet!C21', '=C10*C14', 'compute_max_output_per_line(...)', 'max output / week / line'),
        ('Cleansheet!C23', '=ROUNDUP(C20/C21, 0)', 'compute_lines_installed(...)', 'lines installed'),
        ('Cleansheet!C58', '=(C56+C55)/C24', 'compute_dl_cost(...)', 'DL cost / unit'),
        ('Cleansheet!H90', "=SUM('Equipment List'!P3:P7)*(C20/C24)", 'compute_equipment_mro(...)', 'Equip MRO'),
        ('Cleansheet!H91', "=SUM('Equipment List'!P8:P34)*(C20/C24)", 'compute_equipment_depr(...)', 'Equip Depr'),
        ('Cleansheet!J100', '=(I100*$C$100/$C$18)*I97*(C20/C24)', 'compute_facility_cost(...)', 'Facility'),
        ('Cleansheet!J107', '=C107/I107', 'compute_inbound_freight(...)', 'Freight'),
        ('Cleansheet!J110', '=I110*C110', 'compute_loss_factor(...)', 'Loss'),
        ('Cleansheet!J114', '=C114*I114/C18', 'compute_vat(...)', 'VAT'),
        ('Cleansheet!C118', '=J112+J109+X103+X94+X83+Y70+L58', 'compute_mva_total(...)', 'MVA total'),
        ('Cleansheet!J121', '=I121*(C121+C120)', 'compute_sga_profit(...)', 'SGA+Profit'),
        ('Cleansheet!C125', '=J120+C118+J121', 'compute_total_tc(...)', 'TC total'),
    ]
    headers = ['Excel 位置', '公式', 'DB Service', '對應 Output']
    widths = [1.7, 3.6, 4.0, 3.5]
    cx = 0.4
    for h, w in zip(headers, widths):
        add_rect(s, cx, 1.4, w, 0.4, INK)
        add_text(s, cx, 1.4, w, 0.4, h, size=10, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        cx += w

    for i, row in enumerate(mappings):
        y = 1.8 + i * 0.38
        bg = BG_SOFT if i % 2 == 0 else WHITE
        cx = 0.4
        for j, (cell, w) in enumerate(zip(row, widths)):
            add_rect(s, cx, y, w, 0.38, bg)
            font = 'Consolas' if j in (0, 1, 2) else 'Microsoft JhengHei'
            size = 8 if j == 1 else 9
            add_text(s, cx + 0.08, y, w - 0.16, 0.38, cell, size=size, color=TEXT,
                     font_name=font, anchor=MSO_ANCHOR.MIDDLE)
            cx += w

    slide_footer(s, page_n, total_pages)


def slide_close():
    s = add_blank_slide()
    add_rect(s, 0, 0, 13.333, 7.5, INK)
    add_rect(s, 0, 2.8, 13.333, 0.06, OCEAN)
    add_text(s, 0.7, 1.2, 12, 0.6, '完。', size=42, bold=True, color=WHITE)
    add_text(s, 0.7, 1.85, 12, 0.4, 'Q&A · Internal Discussion', size=18, color=OCEAN)

    add_rounded(s, 0.7, 3.4, 12.0, 3.5, RGBColor(0x14, 0x2F, 0x4E), radius=0.04)
    add_text(s, 1.0, 3.6, 11.5, 0.4, '🔗 相關資源', size=16, bold=True, color=OCEAN)
    add_text(s, 1.0, 4.1, 11.5, 2.7, """• HTML demo:Cortex_互動Demo_v0.9.html(新增 §MVA 操作流程 section)
• 設計文件:cleansheet-mva-sd v0.3(14 表 schema · 5 Layer · 7 Phase 完整流程)
• 對應 BOM SD:bom-collection-sd v0.4(§17 multi-PKG · §18 qty scenario · §19 dual price · §20 VIEW_TRUE_COST)
• 兩案範本:
    - Rival 3+ Wired Mouse Cleansheet_China-20241011-19S.xlsx
    - WHOOP_Gen4 MP quotation book_ for AI.xlsx
• 開發 Roadmap:6.5 sprint ≈ 26 天(Phase 1-6 詳見 SD §13)""", size=12, color=RGBColor(0xCB, 0xD5, 0xE1))

    add_text(s, 0.7, 7.0, 12, 0.3, 'Cortex MVA 操作流程說明手冊 v1 · 2026-06 · FOXLINK GPT',
             size=10, color=MUTED, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════
# RUN
# ════════════════════════════════════════════════════════════════════════

# Pre-calculate total pages
total_pages = 1  # cover
total_pages += 1  # toc
total_pages += 1  # what is MVA
total_pages += 1  # two cases
total_pages += 1  # architecture
total_pages += 4 * len(PHASES)  # each phase: intro + prereq + steps + delta
total_pages += 1  # timeline
total_pages += 1  # checklist
total_pages += 1  # excel mapping
total_pages += 1  # close

print(f"Building {total_pages} slides...")

slide_cover()  # 1
slide_toc()    # 2
slide_what_is_mva()  # 3
slide_two_cases()    # 4
slide_architecture() # 5

page = 6
for ph in PHASES:
    slide_phase_intro(ph, page); page += 1
    slide_phase_prereq(ph, page); page += 1
    slide_phase_steps(ph, page); page += 1
    slide_phase_delta(ph, page); page += 1

slide_timeline(page); page += 1
slide_checklist(page); page += 1
slide_excel_mapping(page); page += 1
slide_close(); page += 1

OUT = 'd:/vibe_coding/foxlink_gpt/docs/Cortex_MVA操作流程說明手冊_v1.pptx'
prs.save(OUT)
print(f'[OK] Saved: {OUT}')
print(f'     Total slides: {len(prs.slides)}')
