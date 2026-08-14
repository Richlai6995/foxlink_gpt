# -*- coding: utf-8 -*-
"""
CORTEX 開發進度報告 — 董事長單頁簡報（4:3 全螢幕）
左：C. 執行狀況（已上線 1,650 人 + 4 項持續執行）
右：D. 2026 下半年規劃（自規劃文件萃取的真實 roadmap）
執行：python docs/gen_chairman_progress_2026h2.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

prs = Presentation()
prs.slide_width = Inches(10)      # 4:3
prs.slide_height = Inches(7.5)

# ── Palette ──
INK   = RGBColor(0x0B, 0x1F, 0x3A)
TEXT  = RGBColor(0x37, 0x47, 0x55)
MUTED = RGBColor(0x6B, 0x72, 0x80)
OCEAN = RGBColor(0x02, 0xC3, 0x9A)
NAVY  = RGBColor(0x1C, 0x72, 0x93)
GOLD  = RGBColor(0xCA, 0x8A, 0x04)
BG    = RGBColor(0xFA, 0xFB, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LINE  = RGBColor(0xE5, 0xE7, 0xEB)
BLUE  = RGBColor(0x0E, 0xA5, 0xE9)
PURPLE= RGBColor(0x7C, 0x3A, 0xED)
GREEN = RGBColor(0x16, 0xA3, 0x4A)
TEAL  = RGBColor(0x08, 0x91, 0xB2)
AMBER = RGBColor(0xB4, 0x53, 0x09)
SLATE = RGBColor(0xCB, 0xD5, 0xE1)
FONT  = 'Microsoft JhengHei'


def add_text(s, x, y, w, h, text, size=12, bold=False, color=None,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tx = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tx.text_frame
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = FONT
    r.font.size = Pt(size)
    r.font.bold = bold
    if color:
        r.font.color.rgb = color
    return tx


def add_runs(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """runs = [(text, size, bold, color), ...] 同一段多 run"""
    tx = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tx.text_frame
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    for (text, size, bold, color) in runs:
        r = p.add_run()
        r.text = text
        r.font.name = FONT
        r.font.size = Pt(size)
        r.font.bold = bold
        if color:
            r.font.color.rgb = color
    return tx


def add_rect(s, x, y, w, h, fill, line=None, lw=0.5):
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    r.fill.solid(); r.fill.fore_color.rgb = fill
    if line is None:
        r.line.fill.background()
    else:
        r.line.color.rgb = line; r.line.width = Pt(lw)
    r.shadow.inherit = False
    return r


def add_round(s, x, y, w, h, fill, line=None, lw=0.5, radius=0.06):
    r = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    r.adjustments[0] = radius
    r.fill.solid(); r.fill.fore_color.rgb = fill
    if line is None:
        r.line.fill.background()
    else:
        r.line.color.rgb = line; r.line.width = Pt(lw)
    r.shadow.inherit = False
    return r


s = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(s, 0, 0, 10, 7.5, BG)

# ═══ Banner ═══
add_rect(s, 0, 0, 10, 1.15, INK)
add_rect(s, 0, 0, 0.16, 1.15, OCEAN)
add_text(s, 0.45, 0.15, 7.0, 0.26, 'FOXLINK 正崴 · CORTEX AI 整合平台', size=10, color=OCEAN)
add_text(s, 0.45, 0.38, 8.2, 0.46, 'CORTEX 開發進度報告 — 執行狀況與下半年規劃',
         size=21, bold=True, color=WHITE)
add_text(s, 0.45, 0.85, 8.6, 0.26,
         '對話 × 知識庫 × AI 戰情 × 教育訓練 × 辦公室自動化 — 一站式內部整合平台',
         size=10.5, color=SLATE)
add_text(s, 7.7, 0.18, 2.1, 0.24, '2026-06 · 董事長簡報', size=9.5, color=SLATE,
         align=PP_ALIGN.RIGHT)

# ═══ 版面參數 ═══
LX, LW = 0.30, 4.55          # 左欄
RX, RW = 5.00, 4.70          # 右欄
HEAD_Y, HEAD_H = 1.30, 0.40

# ════════════════════ 左欄：C. 執行狀況 ════════════════════
add_round(s, LX, HEAD_Y, LW, HEAD_H, NAVY, radius=0.10)
add_text(s, LX + 0.18, HEAD_Y, LW - 0.3, HEAD_H, 'C ▏ 執行狀況（現況）',
         size=13, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)

# Hero metric：已上線 1,650 人
hy = 1.78
add_round(s, LX, hy, LW, 0.84, WHITE, line=OCEAN, lw=1.4, radius=0.05)
add_rect(s, LX, hy, 0.10, 0.84, OCEAN)
add_runs(s, LX + 0.28, hy + 0.08, 2.35, 0.68,
         [('1,650', 38, True, OCEAN), (' 人', 16, True, OCEAN)],
         anchor=MSO_ANCHOR.MIDDLE)
add_text(s, LX + 2.55, hy + 0.15, LW - 2.7, 0.3, '已上線使用人數', size=12, bold=True, color=INK)
add_text(s, LX + 2.55, hy + 0.45, LW - 2.7, 0.34, '全公司日常導入 · 持續成長中',
         size=9, color=MUTED)

# 4 項持續執行（USER C 段原文）
exec_items = [
    ('🪙', '貴金屬情報平台持續優化', NAVY,
     '持續調整運作,讓使用者更快速、正確取得貴金屬行情與決策資訊'),
    ('🔗', '擴大內部系統連接', BLUE,
     '持續串接公司既有系統,充分發揮「整合平台」一站式綜效'),
    ('⚙️', '辦公室自動化已落地', OCEAN,
     '已完成:複雜 Excel 解析 ‧ 長時間語音轉文字 ‧ 各式文檔輸出'),
    ('📊', '深化 AI 資料分析整合', PURPLE,
     '強化 AI 戰情應用,把內部數據變成可即時決策的洞察'),
]
cy = 2.74
ch, cg = 0.90, 0.14
for i, (ico, t, c, d) in enumerate(exec_items):
    y = cy + i * (ch + cg)
    add_round(s, LX, y, LW, ch, WHITE, line=LINE, lw=0.75, radius=0.05)
    add_rect(s, LX, y, 0.10, ch, c)
    add_text(s, LX + 0.24, y + 0.12, 0.5, 0.5, ico, size=17, anchor=MSO_ANCHOR.TOP)
    add_text(s, LX + 0.78, y + 0.13, LW - 0.95, 0.3, t, size=11.5, bold=True, color=INK)
    add_text(s, LX + 0.78, y + 0.45, LW - 0.95, 0.40, d, size=8.8, color=TEXT)

# ════════════════════ 右欄：D. 下半年規劃 ════════════════════
add_round(s, RX, HEAD_Y, RW, HEAD_H, GOLD, radius=0.10)
add_text(s, RX + 0.18, HEAD_Y, RW - 0.3, HEAD_H, 'D ▏ 2026 下半年規劃',
         size=13, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)

# 5 項策略 initiative（自規劃文件萃取）
plan_items = [
    ('通用專案管理平台', '開案 30→5 分', PURPLE,
     '業務開案提速、群聊協作;進階上線贏單預測 + 損益模擬,支撐報價決策'),
    ('BOM × MVA 業務快接接單', '報價週期 −50%', AMBER,
     '中/越/台多廠成本矩陣一眼對比 + AI 智能料號推薦,報價簽核自動鎖版'),
    ('行動化支援', '手機 / PWA', BLUE,
     '手機端 AI 對話與報價查詢,支援現場與越南廠同仁,擴展跨廠區使用'),
    ('辦公室自動化再升級', '3 小時會議逐字', GREEN,
     '長音檔完整轉錄無截斷 + AI 戰情報表模板 / 超大資料集自動分檔輸出'),
    ('ERP 整合自助化', '降低 IT 依賴', TEAL,
     '自然語言探索 ERP 資料表、自助設定同步,加速新系統與資料源接入'),
]
py = 1.78
ph, pg = 0.90, 0.124
for i, (t, kpi, c, d) in enumerate(plan_items):
    y = py + i * (ph + pg)
    add_round(s, RX, y, RW, ph, WHITE, line=LINE, lw=0.75, radius=0.05)
    add_rect(s, RX, y, 0.10, ph, c)
    add_text(s, RX + 0.26, y + 0.12, RW - 1.95, 0.3, t, size=11.5, bold=True, color=INK)
    # KPI pill
    pw = 1.55
    add_round(s, RX + RW - pw - 0.14, y + 0.12, pw, 0.32, c, radius=0.30)
    add_text(s, RX + RW - pw - 0.14, y + 0.12, pw, 0.32, kpi, size=8.5, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, RX + 0.26, y + 0.48, RW - 0.45, 0.40, d, size=8.8, color=TEXT)

# ═══ Footer：策略主軸 ═══
fy = 6.96
add_round(s, LX, fy, 9.40, 0.30, RGBColor(0xF0, 0xFD, 0xFA), line=OCEAN, lw=0.75, radius=0.20)
add_runs(s, LX + 0.22, fy, 9.0, 0.30,
         [('策略主軸　', 9.5, True, NAVY),
          ('下半年聚焦「業務接單提速(專案平台 + BOM)」與「行動化 + 系統整合深化」,'
           '讓 1,650 位使用者的日常工作再省一大段時間。', 9.5, False, TEXT)],
         anchor=MSO_ANCHOR.MIDDLE)

# 最底 footer
add_rect(s, LX, 7.33, 9.40, 0.015, LINE)
add_text(s, LX, 7.35, 6.5, 0.15,
         'CORTEX AI 整合平台 · 開發進度報告', size=7.5, color=MUTED)
add_text(s, 7.0, 7.35, 2.7, 0.15, '2026-06-16 · 內部呈核版',
         size=7.5, color=MUTED, align=PP_ALIGN.RIGHT)

OUT = 'd:/vibe_coding/foxlink_gpt/docs/CORTEX_董事長進度報告_2026H2_v1.pptx'
prs.save(OUT)
print(f'[OK] Saved: {OUT}')
print(f'     Slide: {prs.slide_width/914400:.2f}in x {prs.slide_height/914400:.2f}in (4:3)')
